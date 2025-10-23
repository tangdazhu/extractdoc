# -*- coding: utf-8 -*-
"""
学术风格PPT生成器

适合学术报告、论文演示、研究汇报
配色：深绿色系，专业、严谨
"""
import logging
from typing import List, Dict, Any, Optional
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from utils.config_manager import config
from .base_ppt_generator import BasePPTGenerator

logger = logging.getLogger(__name__)


class AcademicStylePPTGenerator(BasePPTGenerator):
    """
    学术风格PPT生成器
    
    设计特点:
    - 深绿色渐变背景
    - 金色装饰条
    - 白色内容区域
    - 简洁专业
    """
    
    # 配色方案
    COLOR_PRIMARY_DARK = RGBColor(0, 102, 68)  # 深绿
    COLOR_PRIMARY_LIGHT = RGBColor(76, 175, 80)  # 浅绿
    COLOR_ACCENT = RGBColor(255, 193, 7)  # 金色强调
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(50, 50, 50)
    COLOR_TEXT_LIGHT = RGBColor(150, 150, 150)
    
    def __init__(self):
        """初始化生成器"""
        super().__init__()  # 调用基类初始化
        logger.info("初始化学术风格PPT生成器")
    
    def _clean_markdown_text(self, text: str) -> tuple:
        """
        清理Markdown格式标记
        
        Args:
            text: 原始文本（可能包含**加粗**标记）
        
        Returns:
            (清理后的文本, 是否加粗)
        """
        import re
        
        # 检测是否有加粗标记
        is_bold = '**' in text
        
        # 移除加粗标记 **文本**
        cleaned_text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        
        # 移除单个星号 *文本*
        cleaned_text = re.sub(r'\*(.+?)\*', r'\1', cleaned_text)
        
        return cleaned_text, is_bold
    
    def _add_gradient_background(self, slide, angle=90.0):
        """添加渐变背景"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = angle
        fill.gradient_stops[0].color.rgb = self.COLOR_PRIMARY_DARK
        fill.gradient_stops[1].color.rgb = self.COLOR_PRIMARY_LIGHT
    
    def _add_top_accent_bar(self, slide):
        """添加顶部金色装饰条"""
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(13.33), Inches(0.3)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.COLOR_ACCENT
        top_bar.line.fill.background()
    
    def create_cover_slide(self, title: str, subtitle: str = "", 
                          reporter: str = "", date: str = ""):
        """
        创建封面页
        
        Args:
            title: 主标题
            subtitle: 副标题
            reporter: 汇报人/作者
            date: 日期
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8),
            Inches(10), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.2),
                Inches(9.33), Inches(0.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(20)
            subtitle_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # 作者信息（左下）
        if reporter:
            author_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(6.2),
                Inches(4), Inches(0.8)
            )
            author_frame = author_box.text_frame
            author_frame.text = f"作者: {reporter}"
            author_frame.paragraphs[0].font.size = Pt(16)
            author_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # 日期信息（右下）
        if date:
            date_box = slide.shapes.add_textbox(
                Inches(8), Inches(6.5),
                Inches(3.5), Inches(0.5)
            )
            date_frame = date_box.text_frame
            date_frame.text = f"日期: {date}"
            date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            date_frame.paragraphs[0].font.size = Pt(16)
            date_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        logger.info(f"创建封面页: {title}")
    
    def create_catalog_slide(self, catalog_items: List[Dict[str, str]]):
        """
        创建目录页（动态高度，单列布局）
        
        Args:
            catalog_items: 目录项列表，每项包含 {"number": "01", "title": "标题"}
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8),
            Inches(3), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "目录"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # CONTENTS副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(4), Inches(1),
            Inches(3), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "CONTENTS"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = self.COLOR_TEXT_LIGHT
        
        # 从配置读取参数
        max_items = config.get("ppt_generation.generation_preferences.catalog_max_items")
        min_height = config.get("ppt_generation.generation_preferences.catalog_min_item_height")
        max_height = config.get("ppt_generation.generation_preferences.catalog_max_item_height")
        available_height = config.get("ppt_generation.generation_preferences.catalog_available_height")
        start_y = config.get("ppt_generation.generation_preferences.catalog_start_y")
        
        # 限制显示数量
        items_to_show = catalog_items[:max_items]
        total_items = len(items_to_show)
        
        # 动态计算每项高度
        if total_items > 0:
            calculated_height = available_height / total_items
            # 限制在最小和最大高度之间
            item_height = max(min_height, min(calculated_height, max_height))
        else:
            item_height = max_height
        
        # 动态调整字体大小
        if item_height >= 0.4:
            font_size = 18
        elif item_height >= 0.3:
            font_size = 16
        else:
            font_size = 14
        
        for i, item in enumerate(items_to_show):
            number = item.get("number", f"{i+1:02d}")
            title = item.get("title", "")
            
            y_pos = start_y + i * item_height
            
            # 目录项框（动态高度）
            item_box_height = item_height * 0.75
            item_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(y_pos),
                Inches(10), Inches(item_box_height)
            )
            item_shape.fill.solid()
            item_shape.fill.fore_color.rgb = self.COLOR_WHITE
            item_shape.fill.fore_color.brightness = -0.1
            item_shape.line.color.rgb = self.COLOR_ACCENT
            item_shape.line.width = Pt(1.5)
            
            item_text = item_shape.text_frame
            item_text.text = f"{number}  {title}"
            item_text.paragraphs[0].font.size = Pt(font_size)
            item_text.paragraphs[0].font.color.rgb = self.COLOR_TEXT_DARK
            item_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            item_text.margin_left = Inches(0.3)
        
        logger.info(f"创建目录页: {total_items}项，每项高度{item_height:.2f}英寸")
    
    def create_content_slide(self, title: str, content_lines: List[str]):
        """
        创建内容页
        
        Args:
            title: 页面标题
            content_lines: 内容行列表
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(0.3),
            Inches(13.33), Inches(0.9)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        # 标题文字
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(32)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        # 内容区域
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.8),
            Inches(11.73), Inches(5)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = self.COLOR_WHITE
        content_box.line.fill.background()
        
        # 内容文字（使用基类统一方法）
        content_text = content_box.text_frame
        self._setup_text_frame_margins(content_text)
        
        # 从配置读取字体大小
        font_sizes = {
            0: config.get("text_formatting.font_size_level_0", 20),
            1: config.get("text_formatting.font_size_level_1", 18),
            2: config.get("text_formatting.font_size_level_2", 16)
        }
        
        # 使用基类统一方法添加内容
        self._add_bullet_content(content_text, content_lines, font_sizes, self.COLOR_TEXT_DARK)
        
        logger.debug(f"创建内容页: {title} ({len(content_lines)}行)")
    
    def create_two_column_slide(self, title: str, left_content: List[str], right_content: List[str], 
                                left_title: str = "传统方式", right_title: str = "AI方式"):
        """创建左右对比页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(0.3),
            Inches(13.33), Inches(0.9)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(32)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        split_ratio = config.get("ppt_generation.layout_types.two_column.split_ratio", 0.5)
        total_width = 11.73
        gap = 0.5
        left_width = (total_width - gap) * split_ratio
        right_width = total_width - gap - left_width
        
        # 左侧
        left_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.8),
            Inches(left_width), Inches(5)
        )
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = self.COLOR_WHITE
        left_box.line.color.rgb = self.COLOR_ACCENT
        left_box.line.width = Pt(2)
        
        left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(left_width), Inches(0.4))
        left_title_box.text_frame.text = left_title
        left_title_box.text_frame.paragraphs[0].font.size = Pt(18)
        left_title_box.text_frame.paragraphs[0].font.bold = True
        left_title_box.text_frame.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        left_title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        left_text = left_box.text_frame
        left_text.word_wrap = True
        left_text.margin_left = Inches(0.3)
        left_text.margin_top = Inches(0.5)
        for i, line in enumerate(left_content):
            if i > 0:
                left_text.add_paragraph()
            left_text.paragraphs[i].text = f"● {line}"
            left_text.paragraphs[i].font.size = Pt(16)
            left_text.paragraphs[i].font.color.rgb = self.COLOR_TEXT_DARK
        
        # 右侧
        right_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + left_width + gap), Inches(1.8),
            Inches(right_width), Inches(5)
        )
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = self.COLOR_WHITE
        right_box.line.color.rgb = self.COLOR_ACCENT
        right_box.line.width = Pt(2)
        
        right_title_box = slide.shapes.add_textbox(Inches(0.8 + left_width + gap), Inches(1.5), Inches(right_width), Inches(0.4))
        right_title_box.text_frame.text = right_title
        right_title_box.text_frame.paragraphs[0].font.size = Pt(18)
        right_title_box.text_frame.paragraphs[0].font.bold = True
        right_title_box.text_frame.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        right_title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        right_text = right_box.text_frame
        right_text.word_wrap = True
        right_text.margin_left = Inches(0.3)
        right_text.margin_top = Inches(0.5)
        for i, line in enumerate(right_content):
            if i > 0:
                right_text.add_paragraph()
            right_text.paragraphs[i].text = f"● {line}"
            right_text.paragraphs[i].font.size = Pt(16)
            right_text.paragraphs[i].font.color.rgb = self.COLOR_TEXT_DARK
        
        logger.debug(f"创建左右对比页: {title}")
    
    def create_three_column_slide(self, title: str, cards: List[Dict[str, str]]):
        """创建三列卡片页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), Inches(13.33), Inches(0.9))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(32)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        card_width = 3.5
        card_gap = 0.5
        start_x = 1.0
        cards_to_show = cards[:3]
        
        for i, card in enumerate(cards_to_show):
            x_pos = start_x + i * (card_width + card_gap)
            
            card_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2.0), Inches(card_width), Inches(4.5))
            card_box.fill.solid()
            card_box.fill.fore_color.rgb = self.COLOR_WHITE
            card_box.line.color.rgb = self.COLOR_ACCENT
            card_box.line.width = Pt(2)
            
            icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 1.0), Inches(2.5), Inches(1.5), Inches(1.5))
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            icon_circle.line.fill.background()
            
            icon_text = icon_circle.text_frame
            icon_text.text = card.get("icon", f"{i+1}")
            icon_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            icon_text.paragraphs[0].font.size = Pt(48)
            icon_text.paragraphs[0].font.bold = True
            icon_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            icon_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            card_title_box = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(4.2), Inches(card_width - 0.6), Inches(0.5))
            card_title_box.text_frame.text = card.get("title", "")
            card_title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            card_title_box.text_frame.paragraphs[0].font.size = Pt(18)
            card_title_box.text_frame.paragraphs[0].font.bold = True
            card_title_box.text_frame.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
            
            card_content_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(4.8), Inches(card_width - 0.4), Inches(1.8))
            # 智能截断：优先在标点符号处截断
            content = card.get("content", "")
            if len(content) > 30:
                # 尝试在句号、逗号处截断
                for i in range(27, 15, -1):
                    if i < len(content) and content[i] in '。，、；':
                        content = content[:i+1]
                        break
                else:
                    # 没找到标点，直接截断
                    content = content[:27] + "..."
            card_content_box.text_frame.text = content
            card_content_box.text_frame.word_wrap = True
            card_content_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            card_content_box.text_frame.paragraphs[0].font.size = Pt(12)
            card_content_box.text_frame.paragraphs[0].font.color.rgb = self.COLOR_TEXT_DARK
        
        logger.debug(f"创建三列卡片页: {title}, {len(cards_to_show)}张卡片")
    
    def create_flow_diagram_slide(self, title: str, steps: List[Dict[str, str]]):
        """创建流程图页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), Inches(13.33), Inches(0.9))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(32)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        # 从配置读取参数
        max_steps = config.get("ppt_generation.layout_types.flow_diagram.max_steps")
        base_step_width = config.get("ppt_generation.layout_types.flow_diagram.base_step_width")
        base_arrow_width = config.get("ppt_generation.layout_types.flow_diagram.base_arrow_width")
        min_step_width = config.get("ppt_generation.layout_types.flow_diagram.min_step_width")
        min_arrow_width = config.get("ppt_generation.layout_types.flow_diagram.min_arrow_width")
        content_area_width = config.get("ppt_generation.layout_types.flow_diagram.content_area_width")
        step_title_font_size = config.get("ppt_generation.layout_types.flow_diagram.step_title_font_size")
        step_desc_font_size = config.get("ppt_generation.layout_types.flow_diagram.step_desc_font_size")
        step_desc_max_chars = config.get("ppt_generation.layout_types.flow_diagram.step_desc_max_chars")
        
        steps_to_show = steps[:max_steps]
        step_count = len(steps_to_show)
        
        # 动态计算步骤框和箭头宽度
        base_total_width = step_count * base_step_width + (step_count - 1) * base_arrow_width
        
        if base_total_width > content_area_width:
            scale_factor = content_area_width / base_total_width
            step_width = max(min_step_width, base_step_width * scale_factor)
            arrow_width = max(min_arrow_width, base_arrow_width * scale_factor)
            total_width = step_count * step_width + (step_count - 1) * arrow_width
            
            if scale_factor < 0.6:
                step_title_font_size = int(step_title_font_size * 0.7)
                step_desc_font_size = int(step_desc_font_size * 0.7)
                step_desc_max_chars = int(step_desc_max_chars * 0.6)
            elif scale_factor < 0.8:
                step_title_font_size = int(step_title_font_size * 0.85)
                step_desc_font_size = int(step_desc_font_size * 0.85)
                step_desc_max_chars = int(step_desc_max_chars * 0.8)
        else:
            step_width = base_step_width
            arrow_width = base_arrow_width
            total_width = base_total_width
        
        start_x = (13.33 - total_width) / 2
        
        for i, step in enumerate(steps_to_show):
            x_pos = start_x + i * (step_width + arrow_width)
            
            step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2.5), Inches(step_width), Inches(1.5))
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            step_box.line.fill.background()
            
            step_title_text = step_box.text_frame
            step_title_text.text = step.get("title", f"步骤{i+1}")
            step_title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            step_title_text.paragraphs[0].font.size = Pt(step_title_font_size)
            step_title_text.paragraphs[0].font.bold = True
            step_title_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            step_title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            step_title_text.word_wrap = True
            
            desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(4.2), Inches(step_width), Inches(2.0))
            desc_text = desc_box.text_frame
            # 智能截断：优先在标点符号处截断
            description = step.get("description", "")
            if len(description) > step_desc_max_chars:
                truncate_pos = int(step_desc_max_chars * 0.9)
                for j in range(truncate_pos, max(truncate_pos - 10, 0), -1):
                    if j < len(description) and description[j] in '。，、；':
                        description = description[:j+1]
                        break
                else:
                    description = description[:truncate_pos] + "..."
            desc_text.text = description
            desc_text.word_wrap = True
            desc_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_text.paragraphs[0].font.size = Pt(step_desc_font_size)
            desc_text.paragraphs[0].font.color.rgb = self.COLOR_TEXT_DARK
            
            if i < step_count - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_pos + step_width), Inches(2.9), Inches(arrow_width), Inches(0.7))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.COLOR_ACCENT
                arrow.line.fill.background()
        
        logger.debug(f"创建流程图页: {title}, {step_count}个步骤")
    
    def create_timeline_slide(self, title: str, timeline_items: List[Dict[str, str]]):
        """创建时间线页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), Inches(13.33), Inches(0.9))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(32)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        line_x = 3.0
        start_y = 2.0
        item_height = 1.2
        # 从配置读取最大项目数
        max_items = config.get("ppt_generation.layout_types.timeline.max_items", 6)
        items_to_show = timeline_items[:max_items]
        
        for i in range(len(items_to_show)):
            y_pos = start_y + i * item_height
            
            node_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(line_x - 0.15), Inches(y_pos), Inches(0.3), Inches(0.3))
            node_circle.fill.solid()
            node_circle.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            node_circle.line.fill.background()
            
            if i < len(items_to_show) - 1:
                connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(line_x - 0.02), Inches(y_pos + 0.3), Inches(0.04), Inches(item_height - 0.3))
                connector.fill.solid()
                connector.fill.fore_color.rgb = self.COLOR_ACCENT
                connector.line.fill.background()
            
            item = items_to_show[i]
            content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(line_x + 0.5), Inches(y_pos - 0.1), Inches(9), Inches(1.0))
            content_box.fill.solid()
            content_box.fill.fore_color.rgb = self.COLOR_WHITE
            content_box.line.color.rgb = self.COLOR_ACCENT
            content_box.line.width = Pt(1.5)
            
            content_text = content_box.text_frame
            content_text.margin_left = Inches(0.3)
            content_text.margin_top = Inches(0.1)
            content_text.text = item.get("title", "")
            content_text.paragraphs[0].font.size = Pt(16)
            content_text.paragraphs[0].font.bold = True
            content_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
            
            if item.get("content"):
                content_text.add_paragraph()
                # 限制文本长度，避免溢出
                content = item.get("content", "")
                if len(content) > 60:
                    content = content[:57] + "..."
                content_text.paragraphs[1].text = content
                content_text.paragraphs[1].font.size = Pt(14)
                content_text.paragraphs[1].font.color.rgb = self.COLOR_TEXT_DARK
            
            # 启用自动换行
            content_text.word_wrap = True
        
        logger.debug(f"创建时间线页: {title}, {len(items_to_show)}项")
    
    def create_section_slide(self, number: str, title: str):
        """
        创建章节分隔页
        
        Args:
            number: 章节编号（如 "01"）
            title: 章节标题
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        # 大号数字（金色）
        number_box = slide.shapes.add_textbox(
            Inches(3), Inches(2.5),
            Inches(7.33), Inches(1.5)
        )
        number_frame = number_box.text_frame
        number_frame.text = number
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        number_frame.paragraphs[0].font.size = Pt(100)
        number_frame.paragraphs[0].font.bold = True
        number_frame.paragraphs[0].font.color.rgb = self.COLOR_ACCENT
        
        # 章节标题
        section_title = slide.shapes.add_textbox(
            Inches(3), Inches(4.2),
            Inches(7.33), Inches(0.8)
        )
        section_frame = section_title.text_frame
        section_frame.text = title
        section_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        section_frame.paragraphs[0].font.size = Pt(40)
        section_frame.paragraphs[0].font.bold = True
        section_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        logger.debug(f"创建章节页: {number} - {title}")
    
    def create_picture_slide(self, title: str, image_path: str, caption: str = ""):
        """
        创建图片页（左右布局：左侧图片，右侧文字说明）
        
        Args:
            title: 页面标题
            image_path: 图片路径
            caption: 图片说明（文字内容）
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_top_accent_bar(slide)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(0.3),
            Inches(13.33), Inches(0.9)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(32)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        # 左侧：图片容器
        pic_container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.8),
            Inches(6.5), Inches(5.2)
        )
        pic_container.fill.solid()
        pic_container.fill.fore_color.rgb = self.COLOR_WHITE
        pic_container.line.fill.background()
        
        # 插入图片
        try:
            pic = slide.shapes.add_picture(
                image_path,
                Inches(1), Inches(2),
                width=Inches(6)
            )
            # 调整图片大小以适应容器
            if pic.height > Inches(4.8):
                ratio = Inches(4.8) / pic.height
                pic.height = Inches(4.8)
                pic.width = int(pic.width * ratio)
            
            # 居中图片
            pic.left = Inches(4.05) - int(pic.width / 2)
            pic.top = Inches(4.4) - int(pic.height / 2)
            
        except Exception as e:
            logger.error(f"插入图片失败: {e}")
        
        # 右侧：文字说明区域
        if caption:
            text_container = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7.5), Inches(1.8),
                Inches(5.5), Inches(5.2)
            )
            text_container.fill.solid()
            text_container.fill.fore_color.rgb = self.COLOR_WHITE
            text_container.line.fill.background()
            
            text_frame = text_container.text_frame
            self._setup_text_frame_margins(text_frame)
            text_frame.margin_bottom = Inches(0.3)
            
            # 使用基类统一方法添加内容
            lines = caption.split('\n') if caption else []
            font_sizes = {0: 14, 1: 12}
            self._add_bullet_content(text_frame, lines, font_sizes, self.COLOR_TEXT_DARK)
        
        logger.debug(f"创建图片页: {title}")
    
    def save(self, output_path: str):
        """保存PPT"""
        self.prs.save(output_path)
        logger.info(f"PPT已保存: {output_path}")
