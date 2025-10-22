# -*- coding: utf-8 -*-
"""
PPT模板管理器

统一管理PPT模板的加载逻辑
"""

import logging
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches

from utils.config_manager import config

logger = logging.getLogger("converter")


class TemplateManager:
    """PPT模板管理器"""
    
    @staticmethod
    def load_template(template_config: dict, base_dir: Path) -> Presentation:
        """
        加载PPT模板
        
        统一的模板加载逻辑，用于URL模式和文件模式
        
        Args:
            template_config: 模板配置字典（从config.yaml读取）
            base_dir: 项目根目录（通常是settings.BASE_DIR.parent）
            
        Returns:
            Presentation对象
        """
        template_path = template_config.get("template_path")
        
        # 1. 尝试加载配置的模板
        if template_path:
            full_template_path = base_dir / template_path
            
            if full_template_path.exists():
                logger.info(f"使用配置的模板: {full_template_path}")
                return Presentation(str(full_template_path))
            else:
                logger.warning(f"配置的模板不存在: {full_template_path}")
        
        # 2. 使用默认模板
        default_template = base_dir / "config" / "templates" / "business_template.pptx"
        if default_template.exists():
            logger.info(f"使用默认模板: {default_template}")
            return Presentation(str(default_template))
        
        # 3. 创建空白PPT
        logger.warning("默认模板不存在，创建空白PPT")
        prs = Presentation()
        
        # 设置幻灯片尺寸
        slide_size = config.get("ppt_generation.slide_size", {})
        prs.slide_width = Inches(slide_size.get("width", 10.0))
        prs.slide_height = Inches(slide_size.get("height", 7.5))
        
        return prs
    
    @staticmethod
    def get_template_path(template_config: dict, base_dir: Path) -> Optional[Path]:
        """
        获取模板文件路径（不加载）
        
        Args:
            template_config: 模板配置字典
            base_dir: 项目根目录
            
        Returns:
            模板文件路径，如果不存在返回默认模板路径
        """
        template_path = template_config.get("template_path")
        
        # 1. 尝试配置的模板
        if template_path:
            full_template_path = base_dir / template_path
            if full_template_path.exists():
                return full_template_path
        
        # 2. 返回默认模板
        default_template = base_dir / "config" / "templates" / "business_template.pptx"
        if default_template.exists():
            return default_template
        
        # 3. 没有可用模板
        return None
