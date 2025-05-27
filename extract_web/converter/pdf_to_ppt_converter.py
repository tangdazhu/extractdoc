import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import uuid # For generating unique filenames for intermediate images
import logging
import subprocess

# Assuming libreoffice_converter.py exists and has a function like convert_pdf_to_pptx_libreoffice
# from .libreoffice_converter import convert_pdf_to_pptx_libreoffice # This will be a new function

logger = logging.getLogger('converter')

def get_slide_dimensions(prs):
    return prs.slide_width, prs.slide_height

def convert_pdf_to_ppt(pdf_path, output_folder, mode='screenshot', desired_filename_base=None):
    logger.info(f"Starting PDF to PPT conversion. Mode: {mode}, PDF: {pdf_path}, Output: {output_folder}, DesiredBase: {desired_filename_base}")
    base_name_from_pdf = os.path.splitext(os.path.basename(pdf_path))[0]
    
    if desired_filename_base:
        # Ensure desired_filename_base does not already contain .pptx
        if desired_filename_base.lower().endswith('.pptx'):
            actual_base = os.path.splitext(desired_filename_base)[0]
        else:
            actual_base = desired_filename_base
        output_pptx_filename = f"{actual_base}.pptx"
    else:
        output_pptx_filename = f"{base_name_from_pdf}_converted.pptx"
    output_pptx_path = os.path.join(output_folder, output_pptx_filename)

    os.makedirs(output_folder, exist_ok=True)

    if mode == 'libreoffice':
        logger.info(f"Using LibreOffice mode to convert {pdf_path} to {output_pptx_path}")
        # Ensure output_pptx_path is an absolute path for soffice, or that output_folder is.
        # soffice --outdir expects a directory, and it will place the converted file there with the original basename + new extension.
        actual_output_dir = os.path.dirname(output_pptx_path)
        expected_output_filename_by_soffice = f"{os.path.splitext(os.path.basename(pdf_path))[0]}.pptx"
        # The final output path we want is output_pptx_path. Soffice might name it differently if pdf_path's basename is different from output_pptx_path's basename.
        # Let's ensure soffice outputs to the correct directory, then rename if necessary to match output_pptx_path.

        temp_soffice_output_path = os.path.join(actual_output_dir, expected_output_filename_by_soffice)

        try:
            cmd = [
                'soffice', # Or full path to soffice.exe on Windows if not in PATH
                '--headless',
                '--convert-to', 'pptx:impress_pdf_import',
                '--outdir', actual_output_dir, # soffice will use the input PDF's basename
                pdf_path
            ]
            logger.info(f"Executing LibreOffice command: {' '.join(cmd)}")
            process = subprocess.run(cmd, capture_output=True, text=True, timeout=120) # Added timeout

            if process.returncode == 0:
                logger.info(f"LibreOffice process completed successfully for {pdf_path}.")
                # Check if soffice created the file with the expected name (based on pdf_path's basename)
                if os.path.exists(temp_soffice_output_path):
                    # If the soffice output name is different from our desired output_pptx_path (e.g. due to desired_filename_base)
                    # we need to rename it.
                    if temp_soffice_output_path != output_pptx_path:
                        if os.path.exists(output_pptx_path):
                            logger.warning(f"Target file {output_pptx_path} already exists. Overwriting for LibreOffice conversion.")
                            os.remove(output_pptx_path)
                        os.rename(temp_soffice_output_path, output_pptx_path)
                        logger.info(f"Renamed LibreOffice output from {temp_soffice_output_path} to {output_pptx_path}")
                    
                    if os.path.exists(output_pptx_path):
                        logger.info(f"LibreOffice conversion successful: {output_pptx_path}")
                        return True, output_pptx_path, "LibreOffice conversion successful."
                    else:
                        logger.error(f"LibreOffice conversion error: Expected output file {output_pptx_path} not found after potential rename.")
                        return False, None, "LibreOffice conversion failed: Output file not found after rename."
                else:
                    logger.error(f"LibreOffice conversion error: Expected intermediate output {temp_soffice_output_path} not found. Soffice stdout: {process.stdout}, stderr: {process.stderr}")
                    return False, None, f"LibreOffice conversion failed: Soffice did not produce expected output file. Details: {process.stderr[:200]}"
            else:
                logger.error(f"LibreOffice conversion failed for {pdf_path}. Return code: {process.returncode}")
                logger.error(f"LibreOffice stdout: {process.stdout}")
                logger.error(f"LibreOffice stderr: {process.stderr}")
                return False, None, f"LibreOffice conversion failed. Error: {process.stderr[:200]}" # Return first 200 chars of stderr

        except FileNotFoundError:
            logger.error("LibreOffice (soffice) command not found. Ensure LibreOffice is installed and in PATH.")
            return False, None, "LibreOffice (soffice) command not found. Please install LibreOffice and add it to your system PATH."
        except subprocess.TimeoutExpired:
            logger.error(f"LibreOffice conversion timed out for {pdf_path}.")
            return False, None, "LibreOffice conversion timed out."
        except Exception as e_lo:
            logger.error(f"Error during LibreOffice conversion for {pdf_path}: {e_lo}", exc_info=True)
            # Attempt to clean up if soffice created its temp output but rename failed, or if output_pptx_path was created by a previous failed attempt.
            if os.path.exists(temp_soffice_output_path) and temp_soffice_output_path != output_pptx_path: 
                try: os.remove(temp_soffice_output_path) 
                except: pass
            if os.path.exists(output_pptx_path):
                 try: os.remove(output_pptx_path)
                 except: pass
            return False, None, f"LibreOffice conversion error: {str(e_lo)}"

    elif mode == 'screenshot':
        logger.info(f"Using Screenshot mode to convert {pdf_path} to {output_pptx_path}")
        try:
            doc = fitz.open(pdf_path)
            prs = Presentation()
            slide_width_emu, slide_height_emu = get_slide_dimensions(prs)

            temp_image_folder = os.path.join(output_folder, f"temp_imgs_{uuid.uuid4().hex}")
            os.makedirs(temp_image_folder, exist_ok=True)
            logger.debug(f"Temporary image folder created: {temp_image_folder}")

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=300)
                image_path = os.path.join(temp_image_folder, f"page_{page_num + 1}.png")
                pix.save(image_path)

                blank_slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(blank_slide_layout)

                img_width_px = pix.width
                img_height_px = pix.height

                # Convert EMU to inches as float for calculations
                slide_width_val = slide_width_emu / 914400.0
                slide_height_val = slide_height_emu / 914400.0

                scale_w = slide_width_val / img_width_px if img_width_px > 0 else 1
                scale_h = slide_height_val / img_height_px if img_height_px > 0 else 1
                # Ensure scale is not zero if dimensions are positive, to prevent zero-size images
                if scale_w == 0 and slide_width_val > 0 and img_width_px > 0: scale_w = 1.0 / img_width_px # Avoid zero scale if possible
                if scale_h == 0 and slide_height_val > 0 and img_height_px > 0: scale_h = 1.0 / img_height_px # Avoid zero scale if possible
                
                scale = min(scale_w, scale_h)
                if scale == 0 and (scale_w > 0 or scale_h > 0): # If min resulted in 0 but one of them was positive
                    scale = max(scale_w, scale_h) # Try to use a non-zero scale if one exists

                # If scale is still zero, it implies an issue, but proceed cautiously
                if scale == 0:
                    logger.warning(f"Calculated scale is 0 for page {page_num+1}. Image: {img_width_px}x{img_height_px}, Slide: {slide_width_val:.2f}x{slide_height_val:.2f}. This might result in a tiny/invisible image.")
                    # Fallback to a very small scale if absolutely zero, to avoid division by zero in width/height later if they were used directly with scale
                    if img_width_px > 0 and img_height_px > 0 : scale = 0.00001 


                scaled_img_width_val = img_width_px * scale
                scaled_img_height_val = img_height_px * scale

                left_val = (slide_width_val - scaled_img_width_val) / 2
                top_val = (slide_height_val - scaled_img_height_val) / 2
                
                # Ensure left and top are not negative
                left_val = max(0.0, left_val)
                top_val = max(0.0, top_val)

                # Convert final values to Inches for add_picture
                final_left_in = Inches(left_val)
                final_top_in = Inches(top_val)
                final_width_in = Inches(scaled_img_width_val)
                final_height_in = Inches(scaled_img_height_val)

                # Additional check for very small or zero dimensions before adding picture
                if final_width_in <= 0 or final_height_in <= 0:
                    logger.warning(f"Skipping add_picture for page {page_num+1} due to zero or negative calculated dimensions: W={final_width_in}, H={final_height_in}")
                    continue # Skip adding this picture

                try:
                    slide.shapes.add_picture(image_path, final_left_in, final_top_in, width=final_width_in, height=final_height_in)
                except Exception as e_add_pic:
                    logger.error(f"Error adding picture {image_path} to slide: {e_add_pic}", exc_info=True)

            prs.save(output_pptx_path)
            logger.info(f"Screenshot-based PPTX saved to {output_pptx_path}")

            try:
                for img_file in os.listdir(temp_image_folder):
                    os.remove(os.path.join(temp_image_folder, img_file))
                os.rmdir(temp_image_folder)
                logger.debug(f"Temporary image folder {temp_image_folder} cleaned up.")
            except Exception as e_cleanup:
                logger.warning(f"Could not clean up temporary image folder {temp_image_folder}: {e_cleanup}")
            
            doc.close()
            return True, output_pptx_path, "Screenshot conversion successful."

        except Exception as e:
            logger.error(f"Error during screenshot PDF to PPT conversion for {pdf_path}: {e}", exc_info=True)
            if os.path.exists(output_pptx_path):
                try: os.remove(output_pptx_path)
                except Exception as e_del: logger.error(f"Failed to delete partial PPTX {output_pptx_path} on error: {e_del}")
            return False, None, f"Screenshot mode error: {str(e)}"
    else:
        logger.warning(f"Unknown mode '{mode}' for PDF to PPT conversion of {pdf_path}.")
        return False, None, f"Unknown conversion mode: {mode}"

# Example usage (for testing purposes if run directly):
# if __name__ == '__main__':
#     test_pdf_path = "path_to_your_test.pdf"  # Replace with a real PDF path
#     test_output_folder = "./test_output"
#     os.makedirs(test_output_folder, exist_ok=True)

#     # Test screenshot mode
#     print("\nTesting Screenshot Mode...")
#     success_s, path_s, msg_s = convert_pdf_to_ppt(test_pdf_path, test_output_folder, mode='screenshot')
#     if success_s:
#         print(f"Screenshot Success: Output at {path_s}, Message: {msg_s}")
#     else:
#         print(f"Screenshot Failure: Message: {msg_s}")

#     # Test screenshot mode with desired filename
#     print("\nTesting Screenshot Mode with desired filename...")
#     success_sf, path_sf, msg_sf = convert_pdf_to_ppt(test_pdf_path, test_output_folder, mode='screenshot', desired_filename_base="my_screenshot_ppt_123")
#     if success_sf:
#         print(f"Screenshot (Custom Name) Success: Output at {path_sf}, Message: {msg_sf}")
#     else:
#         print(f"Screenshot (Custom Name) Failure: Message: {msg_sf}")


#     # Test LibreOffice mode (will be simulated)
#     print("\nTesting LibreOffice Mode (Simulated)...")
#     success_lo, path_lo, msg_lo = convert_pdf_to_ppt(test_pdf_path, test_output_folder, mode='libreoffice')
#     if success_lo:
#         print(f"LibreOffice (Simulated) Success: Output at {path_lo}, Message: {msg_lo}")
#     else:
#         print(f"LibreOffice (Simulated) Failure: Message: {msg_lo}")
    
#     # Test LibreOffice mode with desired filename (will be simulated)
#     print("\nTesting LibreOffice Mode with desired filename (Simulated)...")
#     success_lof, path_lof, msg_lof = convert_pdf_to_ppt(test_pdf_path, test_output_folder, mode='libreoffice', desired_filename_base="my_libre_ppt_xyz")
#     if success_lof:
#         print(f"LibreOffice (Simulated & Custom Name) Success: Output at {path_lof}, Message: {msg_lof}")
#     else:
#         print(f"LibreOffice (Simulated & Custom Name) Failure: Message: {msg_lof}")


#     print(f"\nCheck the '{test_output_folder}' directory for results.")



# Functions for merging PPTX files (currently used by views.py for screenshot method)
# This might need adjustment if LibreOffice produces PPTs that need a different merge strategy.
def append_ppt(source_prs, target_prs):
    """Appends all slides from source_prs to target_prs."""
    for slide in source_prs.slides:
        target_slide_layout = target_prs.slide_layouts[5] # Assuming blank layout
        # A more robust way would be to copy the source slide's layout if possible,
        # or at least its type. For now, using a default blank layout.
        new_slide = target_prs.slides.add_slide(target_slide_layout)

        # Copy shapes from source slide to new slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text_frame.text = shape.text_frame.text
                # TODO: Copy more text properties (font, size, color, etc.)
            elif hasattr(shape, 'image'): # Check if shape is an image
                # This is a simplified way; direct image data copy might be complex.
                # python-pptx doesn't directly support copying an image object from one pres to another easily.
                # It would involve saving the image to a temp file and re-adding it, or handling EmbeddedPackagePart.
                # For screenshot method, this is fine as we are adding images from files.
                # For LibreOffice outputs, this merge function might not be suitable if they contain complex objects.
                
                # If we are merging PPTs created by the screenshot method, shapes are pictures.
                # For other types of PPTs, this part would need to be much more sophisticated.
                try:
                    # This assumes the shape is a picture and we can access its properties
                    # to re-add it. This is highly dependent on how `python-pptx` handles
                    # shapes added by other tools or by itself.
                    if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                        # Getting the image bytes and re-adding is non-trivial.
                        # For screenshot based slides, the shape *is* the picture we added.
                        # If we have to merge two PPTs made by screenshot method, we are essentially copying slides with one picture.
                        
                        # Simplification: if it's a picture, we need its stream to add it.
                        # This part is tricky and might not work as expected for arbitrary PPTs.
                        # Since our screenshot method adds pictures from file paths, if we were to implement
                        # a generic PPTX merge, we would need to handle images more carefully.
                        # However, `convert_and_merge_pdfs_to_pptx` calls `convert_pdf_to_ppt` for each PDF (which creates a PPTX with images from files)
                        # and then uses this `append_ppt` function.
                        pass # For now, skip direct image copy in this generic append_ppt if not straightforward.
                         # The current convert_and_merge_pdfs_to_pptx actually re-creates slides for merging
                         # by re-adding pictures. This append_ppt is perhaps a misnomer or intended for a different flow.

                except Exception as e_img_copy:
                    logger.warning(f"Could not copy image shape: {e_img_copy}")
                    
            # TODO: Add support for other shape types (tables, charts, etc.)


def convert_and_merge_pdfs_to_pptx(pdf_paths, merged_pptx_path, request_id="", ppt_creation_mode='screenshot'):
    """Converts multiple PDFs to individual PPTX files (using specified mode) and then merges them into one PPTX file."""
    logger.info(f"Starting convert_and_merge_pdfs_to_pptx. Mode for individual conversions: {ppt_creation_mode}. RequestID: {request_id}")
    
    if not pdf_paths:
        logger.warning(f"No PDF paths provided for merging to PPTX. RequestID: {request_id}")
        return False, "No PDF files provided for merging."

    output_folder = os.path.dirname(merged_pptx_path)
    os.makedirs(output_folder, exist_ok=True)

    intermediate_pptx_files = []
    all_conversions_successful = True

    for pdf_path in pdf_paths:
        temp_pptx_base = f"{os.path.splitext(os.path.basename(pdf_path))[0]}_{uuid.uuid4().hex}_temp" # Unique base for intermediate file
        
        success, temp_pptx_output_path, msg = convert_pdf_to_ppt(pdf_path, output_folder, mode=ppt_creation_mode, desired_filename_base=temp_pptx_base)
        
        if success and temp_pptx_output_path:
            intermediate_pptx_files.append(temp_pptx_output_path)
            logger.info(f"Intermediate PPTX created: {temp_pptx_output_path} using mode '{ppt_creation_mode}'. RequestID: {request_id}")
        else:
            all_conversions_successful = False
            logger.error(f"Failed to convert {pdf_path} to PPTX using mode '{ppt_creation_mode}'. Message: {msg}. RequestID: {request_id}")
            # Clean up any intermediate files created so far on failure
            for f_path in intermediate_pptx_files:
                if os.path.exists(f_path): 
                    try: os.remove(f_path)
                    except Exception as e_clean_fail: logger.warning(f"Failed to cleanup intermediate {f_path} on early exit: {e_clean_fail}")
            return False, f"Failed to convert one or more PDFs ({os.path.basename(pdf_path)}). Error: {msg}"

    if not all_conversions_successful or not intermediate_pptx_files:
        logger.error(f"Not all PDF to PPTX conversions were successful or no intermediate files were generated. Cannot merge. RequestID: {request_id}")
        return False, "One or more PDF to PPTX conversions failed."

    try:
        final_prs = Presentation()
        slide_width_emu, slide_height_emu = get_slide_dimensions(final_prs)
        slide_width_in = Inches(slide_width_emu / 914400)
        slide_height_in = Inches(slide_height_emu / 914400)

        logger.info(f"Merging {len(intermediate_pptx_files)} intermediate PPTX files into {merged_pptx_path}. RequestID: {request_id}")

        if ppt_creation_mode == 'screenshot':
            for pptx_file_path in intermediate_pptx_files:
                if not os.path.exists(pptx_file_path):
                    logger.warning(f"Intermediate PPTX file {pptx_file_path} not found. Skipping. RequestID: {request_id}")
                    continue
                
                try:
                    prs_part = Presentation(pptx_file_path)
                    for slide_part in prs_part.slides:
                        added_picture_to_final_slide = False
                        for shape in slide_part.shapes:
                            if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                                try:
                                    blank_slide_layout = final_prs.slide_layouts[5] 
                                    final_slide = final_prs.slides.add_slide(blank_slide_layout)
                                    
                                    from io import BytesIO
                                    image_stream = BytesIO(shape.image.blob)
                                    
                                    img_width_shape_in = Inches(shape.width / 914400)
                                    img_height_shape_in = Inches(shape.height / 914400)

                                    left = (slide_width_in - img_width_shape_in) / 2
                                    top = (slide_height_in - img_height_shape_in) / 2
                                    left = max(Inches(0), left)
                                    top = max(Inches(0), top)

                                    final_slide.shapes.add_picture(image_stream, left, top, width=img_width_shape_in, height=img_height_shape_in)
                                    added_picture_to_final_slide = True
                                    break # Assuming one main picture per slide from screenshot method
                                except Exception as e_shape_copy:
                                    logger.error(f"Error copying picture shape from {pptx_file_path} to merged PPTX: {e_shape_copy}. RequestID: {request_id}", exc_info=True)
                        if not added_picture_to_final_slide:
                            logger.warning(f"No picture found or copied from a slide in {pptx_file_path}. RequestID: {request_id}")
                except Exception as e_open_part:
                    logger.error(f"Error opening or processing intermediate PPTX {pptx_file_path}: {e_open_part}. RequestID: {request_id}", exc_info=True)
        
        elif ppt_creation_mode == 'libreoffice':
            logger.warning(f"Merging of PPTX files created by LibreOffice is not robustly supported by this specific merge function ('{__name__}'). The merged file may be incomplete. A different strategy (merging PDFs first) is used in views.py for LibreOffice merge. RequestID: {request_id}")
            # Attempt a very naive merge, understanding it's limited.
            for pptx_file_path in intermediate_pptx_files:
                try:
                    prs_part = Presentation(pptx_file_path)
                    for slide_part_master in prs_part.slides:
                        target_slide_layout = final_prs.slide_layouts[5] 
                        new_target_slide = final_prs.slides.add_slide(target_slide_layout)
                        for shape_master in slide_part_master.shapes:
                            if shape_master.has_text_frame:
                                try:
                                    new_el = new_target_slide.shapes.add_textbox(shape_master.left, shape_master.top, shape_master.width, shape_master.height)
                                    new_el.text_frame.text = shape_master.text_frame.text
                                except Exception as e_copy_textbox:
                                    logger.warning(f"Naive LO merge: Failed to copy textbox: {e_copy_textbox}")
                            # This does not copy images or other complex shapes for LO parts.
                except Exception as e_merge_lo_part:
                     logger.error(f"Error during naive merge of LibreOffice-created PPTX part {pptx_file_path}: {e_merge_lo_part}. RequestID: {request_id}")

        final_prs.save(merged_pptx_path)
        logger.info(f"Successfully merged intermediate PPTXs into {merged_pptx_path} (Primary mode used: {ppt_creation_mode}). RequestID: {request_id}")
        return True, f"Successfully merged PDFs into {os.path.basename(merged_pptx_path)} using {ppt_creation_mode} for parts."

    except Exception as e:
        logger.error(f"Error merging PPTX files: {e}. RequestID: {request_id}", exc_info=True)
        return False, f"Error during merging of PPTX files: {str(e)}"
    finally:
        logger.debug(f"Cleaning up {len(intermediate_pptx_files)} intermediate PPTX files. RequestID: {request_id}")
        for temp_file in intermediate_pptx_files:
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e_clean:
                    logger.warning(f"Failed to clean up temporary PPTX {temp_file}: {e_clean}. RequestID: {request_id}") 