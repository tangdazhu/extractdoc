# -*- coding: utf-8 -*-
"""
智能PPT生成器 - 新版本
使用固定布局 + Auto-Fit,简单可靠
"""

import logging
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import MSO_AUTO_SIZE

from .fixed_layout_manager import FixedLayoutManager
from .autofit_renderer import AutoFitRenderer

logger = logging.getLogger("converter")


class SmartPPTGenerator:
    """
    智能PPT生成器(新版)
    
    核心理念:
    1. 使用固定布局区域
    2. 利用PowerPoint的auto-fit功能
    3. 简单可靠,不会溢出
    """
    
    def __init__(self, config: dict = None):
        """
        初始化智能PPT生成器
        
        Args:
            config: 配置
        """
        self.config = config or {}
        self.layout_manager = FixedLayoutManager()
        self.renderer = AutoFitRenderer()
    
    def generate_ppt(
        self,
        template_path: Path,
        document_structure: dict,
        page_analyses: List[dict],
        multimodal_data: dict,
        request_id: str
    ) -> Presentation:
        """
        生成PPT
        
        Args:
            template_path: 模板路径
            document_structure: AI分析的文档结构
            page_analyses: AI分析的各页内容
            multimodal_data: 多模态数据
            request_id: 请求ID
            
        Returns:
            生成的Presentation对象
        """
        logger.info("开始智能PPT生成(新版),RequestID=%s", request_id)
        
        # 1. 加载模板
        presentation = Presentation(str(template_path))
        
        # 删除模板中除第一页外的示例页面
        if len(presentation.slides) > 1:
            xml_slides = presentation.slides._sldIdLst
            slides_to_delete = list(range(1, len(xml_slides)))
            for idx in reversed(slides_to_delete):
                rId = xml_slides[idx].rId
                presentation.part.drop_rel(rId)
                del xml_slides[idx]
            logger.debug("已清除模板示例页面(保留第一页)")
        
        # 2. 生成标题页
        self._create_title_slide(
            presentation,
            document_structure.get("title_page", {}),
            multimodal_data
        )
        
        # 3. 生成内容页
        for page_analysis in page_analyses:
            page_num = page_analysis.get("page_number")
            
            # 跳过标题页
            if page_num == document_structure.get("title_page", {}).get("page_number", 1):
                continue
            
            # 跳过低重要度页面
            if page_analysis.get("importance") == "low" and self.config.get("skip_low_importance", False):
                logger.debug("跳过低重要度页面: 第%d页", page_num)
                continue
            
            self._create_content_slide(
                presentation,
                page_analysis,
                multimodal_data
            )
        
        logger.info("智能PPT生成完成,共%d页,RequestID=%s", len(presentation.slides), request_id)
        return presentation
    
    def _create_title_slide(
        self,
        presentation: Presentation,
        title_page_info: dict,
        multimodal_data: dict
    ):
        """创建标题页"""
        
        # 使用模板第一页
        if len(presentation.slides) == 0:
            # 查找标题布局(不hardcode索引)
            title_layout = self._find_layout(presentation, 'title')
            title_slide = presentation.slides.add_slide(title_layout)
        else:
            title_slide = presentation.slides[0]
        
        # 设置标题和副标题
        elements = title_page_info.get("elements", {})
        title_text = elements.get("title", "文档标题")
        subtitle_text = elements.get("subtitle", "AI智能生成")
        
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title_text
        
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = subtitle_text
        
        logger.info("已创建标题页: %s", title_text)
        
        # 如果有元数据表,添加到标题页
        metadata_table = elements.get("metadata_table")
        if metadata_table and metadata_table.get("should_include", False):
            table_page = metadata_table.get("page")
            page_tables = [t for t in multimodal_data.get("tables", []) if t["page"] == table_page]
            
            if page_tables:
                table_data = page_tables[0]["data"]
                
                # 使用固定区域渲染表格(向上移动,避免溢出)
                zone = self.layout_manager.to_inches({
                    'left': 3.0,
                    'top': 4.5,  # 从5.5上移到4.5
                    'width': 4.0,
                    'height': 1.5
                })
                
                self.renderer.render_table(title_slide, table_data, zone, enable_autofit=True)
                logger.info("已将元数据表添加到标题页")
    
    def _create_content_slide(
        self,
        presentation: Presentation,
        page_analysis: dict,
        multimodal_data: dict
    ):
        """
        创建内容页(使用固定布局+auto-fit)
        """
        
        page_num = page_analysis.get("page_number")
        title = page_analysis.get("title", f"第{page_num}页")
        layout_type = page_analysis.get("suggested_layout", "title_and_text")
        
        logger.debug("创建第%d页: %s, 布局=%s", page_num, title, layout_type)
        
        # 1. 创建幻灯片(不hardcode索引)
        content_layout = self._find_layout(presentation, 'content')
        slide = presentation.slides.add_slide(content_layout)
        
        # 2. 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # 3. 清空内容占位符(但不删除,因为auto-size只对placeholder有效)
        self._clear_content_placeholders(slide)
        
        # 4. 获取布局区域
        zones = self.layout_manager.get_zones(layout_type)
        
        # 5. 收集内容
        content = self._collect_page_content(page_num, page_analysis, multimodal_data)
        
        # 6. 渲染内容到固定区域
        self._render_content_to_zones(slide, content, zones, layout_type)
        
        logger.debug("已创建内容页: 第%d页 - %s", page_num, title)
    
    def _clear_content_placeholders(self, slide):
        """
        清空内容占位符(但不删除)
        
        关键: MSO_AUTO_SIZE只对placeholder有效,所以我们保留placeholder
        但清空其内容,避免显示默认文本
        """
        for shape in slide.shapes:
            # 跳过标题
            if shape == slide.shapes.title:
                continue
            
            # 清空placeholder的文本
            if hasattr(shape, 'placeholder_format') and shape.has_text_frame:
                shape.text_frame.clear()
                logger.debug("已清空placeholder: %s", shape.name)
    
    def _collect_page_content(self, page_num: int, page_analysis: dict, multimodal_data: dict) -> dict:
        """
        收集页面内容
        
        Returns:
            {'table': [...], 'images': [...], 'text': '...'}
        """
        content = {
            'table': None,
            'images': [],
            'text': ''
        }
        
        # 收集表格
        page_tables = [t for t in multimodal_data.get("tables", []) if t["page"] == page_num]
        if page_tables:
            content['table'] = page_tables[0]['data']
        
        # 收集图片(根据AI判断过滤)
        ai_approved_images = []
        for element in page_analysis.get("elements", []):
            if element.get("type") == "image" and element.get("should_keep", False):
                size_str = element.get("size", "")
                # 找到对应的图片
                for img in multimodal_data.get("images", []):
                    if img["page"] == page_num:
                        img_size = f"{img.get('width', 0)}x{img.get('height', 0)}"
                        if img_size == size_str and img.get("path") and img["path"].exists():
                            ai_approved_images.append(img["path"])
        
        content['images'] = ai_approved_images
        
        # 收集文本
        page_data = next((p for p in multimodal_data.get("pages", []) if p["page"] == page_num), None)
        if page_data:
            # 优先使用AI重组的文本
            formatted_content = page_analysis.get("formatted_content")
            if formatted_content:
                content['text'] = formatted_content
                logger.debug("第%d页使用AI重新组织的文本", page_num)
            else:
                content['text'] = page_data.get("text", "")
        
        return content
    
    def _render_content_to_zones(self, slide, content: dict, zones: list, layout_type: str):
        """
        将内容渲染到固定区域
        
        Args:
            slide: 幻灯片
            content: 内容字典
            zones: 区域列表
            layout_type: 布局类型
        """
        # 根据布局类型决定渲染策略
        if layout_type == 'title_and_table':
            # 表格 + 文本
            self._render_table_text_layout(slide, content, zones)
        
        elif layout_type == 'title_and_image':
            # 图片 + 文本
            self._render_image_text_layout(slide, content, zones)
        
        elif layout_type == 'title_and_text':
            # 纯文本
            self._render_text_only_layout(slide, content, zones)
        
        elif layout_type == 'title_text_and_image':
            # 文本 + 图片
            self._render_text_image_layout(slide, content, zones)
        
        else:
            # 默认: 文本
            self._render_text_only_layout(slide, content, zones)
    
    def _render_table_text_layout(self, slide, content: dict, zones: list):
        """渲染表格布局(只渲染表格,表格已包含所有信息)"""
        # 找到表格区域
        table_zone = None
        
        for zone in zones:
            zone_inches = self.layout_manager.to_inches(zone)
            if zone['type'] == 'table':
                table_zone = zone_inches
                break
        
        # 只渲染表格(表格已包含所有信息,不需要额外文本)
        if content.get('table') and table_zone:
            self.renderer.render_table(slide, content['table'], table_zone, enable_autofit=True)
            logger.debug("title_and_table布局:只渲染表格")
    
    def _render_image_text_layout(self, slide, content: dict, zones: list):
        """
        渲染图片+文本布局
        
        注意:图片会覆盖placeholder,所以需要删除placeholder并使用textbox
        """
        image_zone = None
        text_zone = None
        
        for zone in zones:
            zone_inches = self.layout_manager.to_inches(zone)
            if zone['type'] == 'image':
                image_zone = zone_inches
            elif zone['type'] == 'text':
                text_zone = zone_inches
        
        # 删除内容placeholder(因为图片会覆盖它)
        for shape in list(slide.shapes):
            if (hasattr(shape, 'placeholder_format') and 
                shape.has_text_frame and 
                shape != slide.shapes.title):
                sp = shape.element
                sp.getparent().remove(sp)
                logger.debug("删除placeholder以避免被图片遮挡")
        
        # 渲染图片
        if content.get('images') and image_zone:
            if len(content['images']) > 1:
                # 多张图片并排
                self.renderer.render_images_side_by_side(slide, content['images'], image_zone)
            elif len(content['images']) == 1:
                # 单张图片
                self.renderer.render_image(slide, content['images'][0], image_zone)
        
        # 渲染文本(使用textbox,因为placeholder已删除)
        if content.get('text') and text_zone:
            # 创建textbox
            textbox = slide.shapes.add_textbox(
                text_zone['left'],
                text_zone['top'],
                text_zone['width'],
                text_zone['height']
            )
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            
            # 清除边距
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            # 添加文本
            lines = content['text'].split('\n')
            for idx, line in enumerate(lines):
                if not line.strip():
                    continue
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = PptPt(9)  # 固定字体9pt
            
            logger.debug("已渲染文本到textbox: %d行", len(lines))
    
    def _render_text_only_layout(self, slide, content: dict, zones: list):
        """渲染纯文本布局"""
        text_zone = None
        
        for zone in zones:
            zone_inches = self.layout_manager.to_inches(zone)
            if zone['type'] == 'text':
                text_zone = zone_inches
                break
        
        # 渲染文本
        if content.get('text') and text_zone:
            self.renderer.render_text(slide, content['text'], text_zone, enable_autofit=True)
    
    def _render_text_image_layout(self, slide, content: dict, zones: list):
        """渲染文本+图片布局"""
        text_zone = None
        image_zone = None
        
        for zone in zones:
            zone_inches = self.layout_manager.to_inches(zone)
            if zone['type'] == 'text' and not text_zone:
                text_zone = zone_inches
            elif zone['type'] == 'image' and not image_zone:
                image_zone = zone_inches
        
        # 渲染文本
        if content.get('text') and text_zone:
            self.renderer.render_text(slide, content['text'], text_zone, enable_autofit=True)
        
        # 渲染图片
        if content.get('images') and image_zone:
            if len(content['images']) > 1:
                self.renderer.render_images_side_by_side(slide, content['images'], image_zone)
            elif len(content['images']) == 1:
                self.renderer.render_image(slide, content['images'][0], image_zone)
    
    def _find_layout(self, presentation: Presentation, layout_type: str):
        """
        根据类型查找合适的布局(不hardcode索引)
        
        Args:
            presentation: 演示文稿对象
            layout_type: 布局类型 ('title' 或 'content')
        
        Returns:
            找到的布局对象
        """
        layouts = presentation.slide_layouts
        
        if layout_type == 'title':
            # 查找标题布局
            # 常见名称: "Title Slide", "标题幻灯片", "Title Only"
            for layout in layouts:
                name_lower = layout.name.lower()
                if 'title' in name_lower and ('slide' in name_lower or 'only' in name_lower):
                    logger.debug("找到标题布局: %s (索引=%d)", layout.name, layouts.index(layout))
                    return layout
            
            # 回退: 使用第一个布局
            logger.warning("未找到标题布局,使用第一个布局: %s", layouts[0].name)
            return layouts[0]
        
        elif layout_type == 'content':
            # 查找内容布局
            # 常见名称: "Title and Content", "标题和内容", "Content"
            for layout in layouts:
                name_lower = layout.name.lower()
                if ('title' in name_lower and 'content' in name_lower) or \
                   ('标题' in layout.name and '内容' in layout.name):
                    logger.debug("找到内容布局: %s (索引=%d)", layout.name, layouts.index(layout))
                    return layout
            
            # 回退: 使用第二个布局(如果存在)
            if len(layouts) > 1:
                logger.warning("未找到内容布局,使用第二个布局: %s", layouts[1].name)
                return layouts[1]
            else:
                logger.warning("未找到内容布局,使用第一个布局: %s", layouts[0].name)
                return layouts[0]
        
        else:
            # 未知类型,使用第一个布局
            logger.warning("未知布局类型: %s, 使用第一个布局", layout_type)
            return layouts[0]
