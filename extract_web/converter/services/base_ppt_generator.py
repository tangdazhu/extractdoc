# -*- coding: utf-8 -*-
"""
PPT生成器基类

提供通用的PPT生成功能，避免代码重复
"""

import logging
from typing import List, Dict, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from utils.config_manager import config

logger = logging.getLogger(__name__)


class BasePPTGenerator:
    """PPT生成器基类"""
    
    def __init__(self):
        """初始化基类"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
    
    def _clean_markdown_text(self, text: str) -> Tuple[str, bool]:
        """
        清理Markdown标记
        
        Args:
            text: 原始文本
            
        Returns:
            (清理后的文本, 是否加粗)
        """
        is_bold = False
        
        # 处理加粗标记
        if text.startswith("**") and text.endswith("**") and len(text) > 4:
            text = text[2:-2]
            is_bold = True
        
        # 移除其他Markdown标记
        text = text.replace("*", "").replace("_", "")
        
        return text, is_bold
    
    def _add_bullet_content(self, text_frame, content_lines: List[str], font_sizes: Dict[int, int], text_color: RGBColor):
        """
        统一添加带bullet的内容（避免重复代码）
        
        Args:
            text_frame: 文本框对象
            content_lines: 内容行列表
            font_sizes: 字体大小映射 {indent_level: font_size}
            text_color: 文字颜色
        """
        # 从配置读取bullet符号
        bullet_symbols = {
            0: config.get("text_formatting.bullet_level_0", "●"),
            1: config.get("text_formatting.bullet_level_1", "○"),
            2: config.get("text_formatting.bullet_level_2", "▪")
        }
        
        # 清空第一个默认段落，但不使用它
        # PowerPoint的第一个默认段落有无法修改的缩进
        if text_frame.paragraphs:
            first_para = text_frame.paragraphs[0]
            first_para.text = ""  # 清空但保留
        
        logger.debug(f"[缩进修复] 开始添加{len(content_lines)}行内容")
        
        # 处理内容行 - 所有段落都用add_paragraph创建（跳过第一个默认段落）
        for i, line in enumerate(content_lines):
            # 所有段落都通过add_paragraph创建，避免使用默认段落
            para = text_frame.add_paragraph()
            
            # 检测缩进层级
            indent_level = 0
            clean_line = line
            if line.startswith("  - "):
                indent_level = 1
                clean_line = line[4:]
            elif line.startswith("- "):
                indent_level = 0
                clean_line = line[2:]
            elif line.startswith("• "):
                indent_level = 0
                clean_line = line[2:]
            
            # 清理Markdown标记
            clean_line, is_bold = self._clean_markdown_text(clean_line)
            
            # 添加bullet符号（不添加空格，让PowerPoint的默认缩进作为左边距）
            bullet = bullet_symbols.get(indent_level, "●")
            # 1级缩进使用4个空格
            indent_spaces = "    " if indent_level == 1 else ""
            para.text = f"{indent_spaces}{bullet} {clean_line}"
            
            if i == 0:
                logger.debug(f"[缩进修复] 第一行文本: '{para.text[:30]}...'")
            
            # 设置字体（必须在设置text之后）
            font_size = font_sizes.get(indent_level, 18)
            para.font.size = Pt(font_size)
            para.font.color.rgb = text_color
            if is_bold:
                para.font.bold = True
            
            # 不设置para.level！它会自动添加缩进
            # 如果需要缩进，通过在文本前添加空格实现
            # （因为paragraph_format不可用）
    
    def _setup_text_frame_margins(self, text_frame):
        """
        统一设置文本框边距
        
        Args:
            text_frame: 文本框对象
        """
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0)  # 完全去除左边距
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        logger.debug(f"[缩进修复] 文本框边距设置: left=0, right=0.2, top=0.2")
    
    def _get_three_column_config(self) -> Dict:
        """获取三列卡片配置"""
        return {
            "max_cards": config.get("ppt_generation.layout_types.three_column.max_cards"),
            "card_width": config.get("ppt_generation.layout_types.three_column.card_width"),
            "card_gap": config.get("ppt_generation.layout_types.three_column.card_gap"),
            "card_title_font_size": config.get("ppt_generation.layout_types.three_column.card_title_font_size"),
            "card_content_font_size": config.get("ppt_generation.layout_types.three_column.card_content_font_size"),
            "card_content_max_chars": config.get("ppt_generation.layout_types.three_column.card_content_max_chars")
        }
    
    def _get_timeline_config(self) -> Dict:
        """获取时间线配置"""
        return {
            "max_items": config.get("ppt_generation.layout_types.timeline.max_items"),
            "base_item_height": config.get("ppt_generation.layout_types.timeline.base_item_height"),
            "min_item_height": config.get("ppt_generation.layout_types.timeline.min_item_height"),
            "available_height": config.get("ppt_generation.layout_types.timeline.available_height"),
            "start_y": config.get("ppt_generation.layout_types.timeline.start_y"),
            "title_font_size": config.get("ppt_generation.layout_types.timeline.title_font_size"),
            "content_font_size": config.get("ppt_generation.layout_types.timeline.content_font_size"),
            "content_max_chars": config.get("ppt_generation.layout_types.timeline.content_max_chars")
        }
    
    def _truncate_text_smart(self, text: str, max_chars: int) -> str:
        """智能截断文本：优先在标点符号处截断"""
        if len(text) <= max_chars:
            return text
        
        truncate_pos = int(max_chars * 0.9)
        for j in range(truncate_pos, max(truncate_pos - 10, 0), -1):
            if j < len(text) and text[j] in '。，、；':
                return text[:j+1]
        
        return text[:truncate_pos] + "..."
    
    def _calculate_timeline_layout(self, item_count: int, cfg: Dict) -> Tuple[float, int, int, int]:
        """计算时间线动态布局参数
        
        Returns:
            (item_height, title_font_size, content_font_size, content_max_chars)
        """
        if item_count > 0:
            calculated_height = cfg["available_height"] / item_count
            item_height = max(cfg["min_item_height"], min(calculated_height, cfg["base_item_height"]))
        else:
            item_height = cfg["base_item_height"]
        
        title_font_size = cfg["title_font_size"]
        content_font_size = cfg["content_font_size"]
        content_max_chars = cfg["content_max_chars"]
        
        # 根据高度调整字体大小
        if item_height < 1.0:
            title_font_size = int(title_font_size * 0.85)
            content_font_size = int(content_font_size * 0.85)
            content_max_chars = int(content_max_chars * 0.8)
        
        return item_height, title_font_size, content_font_size, content_max_chars
    
    def _get_catalog_config(self) -> Dict:
        """获取目录页配置"""
        return {
            "max_items": config.get("ppt_generation.generation_preferences.catalog_max_items"),
            "min_height": config.get("ppt_generation.generation_preferences.catalog_min_item_height"),
            "max_height": config.get("ppt_generation.generation_preferences.catalog_max_item_height"),
            "available_height": config.get("ppt_generation.generation_preferences.catalog_available_height"),
            "start_y": config.get("ppt_generation.generation_preferences.catalog_start_y")
        }
    
    def _calculate_catalog_layout(self, total_items: int, cfg: Dict) -> Tuple[float, int, int]:
        """计算目录页动态布局参数
        
        Returns:
            (item_height, number_font_size, title_font_size)
        """
        # 动态计算每项高度
        if total_items > 0:
            calculated_height = cfg["available_height"] / total_items
            item_height = max(cfg["min_height"], min(calculated_height, cfg["max_height"]))
        else:
            item_height = cfg["max_height"]
        
        # 动态调整字体大小
        if item_height >= 0.4:
            number_font_size = 20
            title_font_size = 18
        elif item_height >= 0.3:
            number_font_size = 18
            title_font_size = 16
        else:
            number_font_size = 16
            title_font_size = 14
        
        return item_height, number_font_size, title_font_size
    
    def _get_flow_diagram_config(self) -> Dict:
        """获取流程图配置"""
        return {
            "max_steps": config.get("ppt_generation.layout_types.flow_diagram.max_steps"),
            "base_step_width": config.get("ppt_generation.layout_types.flow_diagram.base_step_width"),
            "base_arrow_width": config.get("ppt_generation.layout_types.flow_diagram.base_arrow_width"),
            "min_step_width": config.get("ppt_generation.layout_types.flow_diagram.min_step_width"),
            "min_arrow_width": config.get("ppt_generation.layout_types.flow_diagram.min_arrow_width"),
            "content_area_width": config.get("ppt_generation.layout_types.flow_diagram.content_area_width"),
            "step_title_font_size": config.get("ppt_generation.layout_types.flow_diagram.step_title_font_size"),
            "step_desc_font_size": config.get("ppt_generation.layout_types.flow_diagram.step_desc_font_size"),
            "step_desc_max_chars": config.get("ppt_generation.layout_types.flow_diagram.step_desc_max_chars")
        }
    
    def _calculate_flow_diagram_layout(self, step_count: int, cfg: Dict) -> Tuple[float, float, float, int, int, int]:
        """计算流程图动态布局参数
        
        Returns:
            (step_width, arrow_width, total_width, step_title_font_size, step_desc_font_size, step_desc_max_chars)
        """
        base_step_width = cfg["base_step_width"]
        base_arrow_width = cfg["base_arrow_width"]
        min_step_width = cfg["min_step_width"]
        min_arrow_width = cfg["min_arrow_width"]
        content_area_width = cfg["content_area_width"]
        step_title_font_size = cfg["step_title_font_size"]
        step_desc_font_size = cfg["step_desc_font_size"]
        step_desc_max_chars = cfg["step_desc_max_chars"]
        
        # 计算基础宽度总和
        base_total_width = step_count * base_step_width + (step_count - 1) * base_arrow_width
        
        if base_total_width > content_area_width:
            # 需要缩小，按比例调整
            scale_factor = content_area_width / base_total_width
            step_width = max(min_step_width, base_step_width * scale_factor)
            arrow_width = max(min_arrow_width, base_arrow_width * scale_factor)
            
            # 重新计算实际宽度
            total_width = step_count * step_width + (step_count - 1) * arrow_width
            
            # 根据缩放调整字体大小
            if scale_factor < 0.6:
                step_title_font_size = int(step_title_font_size * 0.7)
                step_desc_font_size = int(step_desc_font_size * 0.7)
                step_desc_max_chars = int(step_desc_max_chars * 0.6)
            elif scale_factor < 0.8:
                step_title_font_size = int(step_title_font_size * 0.85)
                step_desc_font_size = int(step_desc_font_size * 0.85)
                step_desc_max_chars = int(step_desc_max_chars * 0.8)
        else:
            # 不需要缩小，使用基础宽度
            step_width = base_step_width
            arrow_width = base_arrow_width
            total_width = base_total_width
        
        return step_width, arrow_width, total_width, step_title_font_size, step_desc_font_size, step_desc_max_chars
    
    def _create_title_bar(self, slide, y_pos: float, height: float, title: str, font_size: int, title_color: RGBColor):
        """
        创建标题栏（通用方法）
        
        Args:
            slide: 幻灯片对象
            y_pos: Y位置（英寸）
            height: 高度（英寸）
            title: 标题文本
            font_size: 字体大小
            title_color: 标题颜色
        """
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, Inches(y_pos),
            Inches(13.33), Inches(height)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色背景
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(font_size)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = title_color
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        return title_bar
    
    def _create_content_box(self, slide, x: float, y: float, width: float, height: float, 
                           border_color: RGBColor = None, border_width: float = 0):
        """
        创建内容框（通用方法）
        
        Args:
            slide: 幻灯片对象
            x: X位置（英寸）
            y: Y位置（英寸）
            width: 宽度（英寸）
            height: 高度（英寸）
            border_color: 边框颜色（可选）
            border_width: 边框宽度（可选）
        """
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色背景
        
        if border_color and border_width > 0:
            content_box.line.color.rgb = border_color
            content_box.line.width = Pt(border_width)
        else:
            content_box.line.fill.background()
        
        return content_box
    
    def save(self, output_path: str):
        """保存PPT"""
        self.prs.save(output_path)
        logger.info(f"PPT已保存: {output_path}")
