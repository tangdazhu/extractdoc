# -*- coding: utf-8 -*-
"""
创建 PPT 模板文件

生成两个预定义模板：
1. 现代商务风格 (business_template.pptx) - 深蓝渐变背景
2. 学术报告风格 (academic_template.pptx) - 深色专业背景
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path


def set_gradient_background(slide, color1, color2):
    """设置渐变背景"""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


def add_decorative_shape(slide, left, top, width, height, color, opacity=0.3):
    """添加装饰形状"""
    shape = slide.shapes.add_shape(
        1,  # 矩形
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = opacity
    shape.line.fill.background()


def create_business_template():
    """创建现代商务风格模板 - 深蓝渐变背景"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 标题布局
    title_layout = prs.slide_layouts[0]
    # 内容布局
    bullet_layout = prs.slide_layouts[1]
    
    # === 标题页 ===
    slide = prs.slides.add_slide(title_layout)
    
    # 设置深蓝渐变背景
    set_gradient_background(
        slide,
        RGBColor(13, 27, 62),   # 深蓝色
        RGBColor(25, 55, 109)   # 中蓝色
    )
    
    # 添加装饰元素 - 左下角浅蓝色圆形
    add_decorative_shape(
        slide,
        Inches(-1), Inches(3.5),
        Inches(3.5), Inches(3.5),
        RGBColor(41, 128, 185),
        opacity=0.2
    )
    
    # 添加装饰元素 - 右上角金色圆形
    add_decorative_shape(
        slide,
        Inches(8), Inches(-1),
        Inches(3), Inches(3),
        RGBColor(243, 156, 18),
        opacity=0.15
    )
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "现代商务风格"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    subtitle.text = "专业 · 简洁 · 高效"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)  # 浅灰白
    
    # === 内容页 ===
    slide2 = prs.slides.add_slide(bullet_layout)
    
    # 设置相同的渐变背景
    set_gradient_background(
        slide2,
        RGBColor(13, 27, 62),
        RGBColor(25, 55, 109)
    )
    
    # 添加顶部装饰条
    add_decorative_shape(
        slide2,
        Inches(0), Inches(0),
        Inches(10), Inches(0.15),
        RGBColor(52, 152, 219),
        opacity=0.8
    )
    
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "内容页示例"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    tf = body_shape.text_frame
    tf.text = "核心业务指标"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)
    
    p = tf.add_paragraph()
    p.text = "市场增长策略"
    p.level = 0
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(236, 240, 241)
    
    output_path = Path(__file__).parent / "templates" / "business_template.pptx"
    prs.save(str(output_path))
    print(f"[OK] Business template created: {output_path}")


def create_academic_template():
    """创建学术报告风格模板 - 深色专业背景"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 标题布局
    title_layout = prs.slide_layouts[0]
    # 内容布局
    bullet_layout = prs.slide_layouts[1]
    
    # === 标题页 ===
    slide = prs.slides.add_slide(title_layout)
    
    # 设置深灰渐变背景(学术风格)
    set_gradient_background(
        slide,
        RGBColor(44, 62, 80),   # 深灰蓝
        RGBColor(52, 73, 94)    # 中灰蓝
    )
    
    # 添加左侧装饰条 - 学术风格的竖条
    add_decorative_shape(
        slide,
        Inches(0), Inches(0),
        Inches(0.2), Inches(5.625),
        RGBColor(231, 76, 60),  # 学术红
        opacity=0.9
    )
    
    # 添加底部装饰条
    add_decorative_shape(
        slide,
        Inches(0), Inches(5.425),
        Inches(10), Inches(0.2),
        RGBColor(231, 76, 60),
        opacity=0.7
    )
    
    # 添加右上角装饰元素
    add_decorative_shape(
        slide,
        Inches(7.5), Inches(0),
        Inches(2.5), Inches(1.5),
        RGBColor(149, 165, 166),
        opacity=0.15
    )
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "学术报告风格"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色
    
    subtitle.text = "Research · Innovation · Excellence"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)  # 浅灰白
    
    # === 内容页 ===
    slide2 = prs.slides.add_slide(bullet_layout)
    
    # 设置相同的渐变背景
    set_gradient_background(
        slide2,
        RGBColor(44, 62, 80),
        RGBColor(52, 73, 94)
    )
    
    # 添加左侧装饰条
    add_decorative_shape(
        slide2,
        Inches(0), Inches(0),
        Inches(0.2), Inches(5.625),
        RGBColor(231, 76, 60),
        opacity=0.9
    )
    
    # 添加标题下方装饰线
    add_decorative_shape(
        slide2,
        Inches(0.5), Inches(1.1),
        Inches(3), Inches(0.05),
        RGBColor(231, 76, 60),
        opacity=0.8
    )
    
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "研究内容"
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    tf = body_shape.text_frame
    tf.text = "研究背景与动机"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)
    
    p = tf.add_paragraph()
    p.text = "研究方法与实验设计"
    p.level = 0
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(236, 240, 241)
    
    p2 = tf.add_paragraph()
    p2.text = "结果分析与讨论"
    p2.level = 0
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(236, 240, 241)
    
    output_path = Path(__file__).parent / "templates" / "academic_template.pptx"
    prs.save(str(output_path))
    print(f"[OK] Academic template created: {output_path}")


if __name__ == "__main__":
    create_business_template()
    create_academic_template()
    print("\n[SUCCESS] All templates created!")
