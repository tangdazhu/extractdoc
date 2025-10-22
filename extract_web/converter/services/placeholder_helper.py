# -*- coding: utf-8 -*-
"""
PPT占位符辅助工具

正确使用PPT模板的占位符机制，保留模板样式
"""

import logging
from typing import Optional, List
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE, MSO_SHAPE_TYPE
from pptx.util import Pt as PptPt

logger = logging.getLogger("converter")


class PlaceholderHelper:
    """
    PPT占位符辅助工具
    
    核心理念：
    1. 使用模板的占位符，而不是手动创建形状
    2. 保留模板的所有样式（背景、字体、颜色等）
    3. 只填充内容，不修改样式
    """
    
    @staticmethod
    def get_title_placeholder(slide) -> Optional[SlidePlaceholder]:
        """
        获取标题占位符
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            标题占位符，如果没有返回None
        """
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER_TYPE.TITLE or phf.type == PP_PLACEHOLDER_TYPE.CENTER_TITLE:
                    return shape
        return None
    
    @staticmethod
    def get_content_placeholder(slide) -> Optional[SlidePlaceholder]:
        """
        获取内容占位符（BODY或OBJECT类型）
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            内容占位符，如果没有返回None
        """
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type in (PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.OBJECT):
                    return shape
        return None
    
    @staticmethod
    def get_subtitle_placeholder(slide) -> Optional[SlidePlaceholder]:
        """
        获取副标题占位符
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            副标题占位符，如果没有返回None
        """
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                    return shape
        return None
    
    @staticmethod
    def get_picture_placeholder(slide) -> Optional[SlidePlaceholder]:
        """
        获取图片占位符
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            图片占位符，如果没有返回None
        """
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER_TYPE.PICTURE:
                    return shape
        return None
    
    @staticmethod
    def get_all_content_placeholders(slide) -> List[SlidePlaceholder]:
        """
        获取所有内容占位符（用于多列布局）
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            内容占位符列表
        """
        placeholders = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type in (PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.OBJECT):
                    placeholders.append(shape)
        return placeholders
    
    @staticmethod
    def fill_title(slide, title_text: str) -> bool:
        """
        填充标题占位符
        
        Args:
            slide: 幻灯片对象
            title_text: 标题文本
            
        Returns:
            是否成功填充
        """
        title_ph = PlaceholderHelper.get_title_placeholder(slide)
        if title_ph and hasattr(title_ph, 'text_frame'):
            title_ph.text = title_text
            logger.debug(f"填充标题: {title_text}")
            return True
        else:
            logger.warning("未找到标题占位符")
            return False
    
    @staticmethod
    def fill_text_content(slide, content_lines: List[str], use_bullets: bool = True) -> bool:
        """
        填充文本内容到占位符
        
        Args:
            slide: 幻灯片对象
            content_lines: 内容行列表
            use_bullets: 是否使用项目符号
            
        Returns:
            是否成功填充
        """
        content_ph = PlaceholderHelper.get_content_placeholder(slide)
        if not content_ph or not hasattr(content_ph, 'text_frame'):
            logger.warning("未找到内容占位符")
            return False
        
        text_frame = content_ph.text_frame
        text_frame.clear()  # 清空默认文本
        
        for idx, line in enumerate(content_lines):
            if not line.strip():
                continue
                
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            
            # 设置项目符号
            if use_bullets:
                p.level = 0  # 一级项目符号
        
        logger.debug(f"填充内容: {len(content_lines)}行")
        return True
    
    @staticmethod
    def fill_subtitle(slide, subtitle_text: str) -> bool:
        """
        填充副标题占位符
        
        Args:
            slide: 幻灯片对象
            subtitle_text: 副标题文本
            
        Returns:
            是否成功填充
        """
        subtitle_ph = PlaceholderHelper.get_subtitle_placeholder(slide)
        if subtitle_ph and hasattr(subtitle_ph, 'text_frame'):
            subtitle_ph.text = subtitle_text
            logger.debug(f"填充副标题: {subtitle_text}")
            return True
        else:
            logger.warning("未找到副标题占位符")
            return False
    
    @staticmethod
    def insert_picture_to_placeholder(slide, image_path: str) -> bool:
        """
        插入图片到图片占位符
        
        Args:
            slide: 幻灯片对象
            image_path: 图片文件路径
            
        Returns:
            是否成功插入
        """
        picture_ph = PlaceholderHelper.get_picture_placeholder(slide)
        if not picture_ph:
            logger.warning("未找到图片占位符")
            return False
        
        try:
            # 使用占位符的insert_picture方法
            picture_ph.insert_picture(image_path)
            logger.debug(f"插入图片到占位符: {image_path}")
            return True
        except Exception as e:
            logger.error(f"插入图片失败: {e}")
            return False
    
    @staticmethod
    def get_layout_info(slide_layout) -> dict:
        """
        获取布局信息（用于调试）
        
        Args:
            slide_layout: 布局对象
            
        Returns:
            布局信息字典
        """
        info = {
            "name": slide_layout.name,
            "placeholders": []
        }
        
        for shape in slide_layout.placeholders:
            ph_info = {
                "idx": shape.placeholder_format.idx,
                "type": shape.placeholder_format.type,
                "name": shape.name
            }
            info["placeholders"].append(ph_info)
        
        return info
