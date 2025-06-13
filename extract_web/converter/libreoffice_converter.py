import subprocess
import os
import shutil  # For moving file if needed
import logging
import re  # Add re import at module level

logger = logging.getLogger("converter")


def convert_to_pdf(input_path, output_dir):
    """
    Converts a document to PDF using LibreOffice (soffice).

    Args:
        input_path (str): Absolute path to the input document.
        output_dir (str): Absolute path to the directory where the PDF should be saved.

    Returns:
        tuple: (bool_success, pdf_output_path_or_error_msg, original_pdf_filename_or_None)
               - bool_success: True if conversion was successful, False otherwise.
               - pdf_output_path_or_error_msg: Absolute path to the converted PDF if successful,
                                               or an error message string if failed.
               - original_pdf_filename_or_None: The original filename of the PDF as created by LibreOffice.
    """
    if not os.path.exists(input_path):
        return False, f"Input file not found: {input_path}", None
    if not os.path.isdir(output_dir):
        return False, f"Output directory not found: {output_dir}", None

    input_filename_stem = os.path.splitext(os.path.basename(input_path))[0]
    expected_pdf_filename = f"{input_filename_stem}.pdf"
    potentially_converted_pdf_path = os.path.join(output_dir, expected_pdf_filename)

    if os.path.exists(potentially_converted_pdf_path):
        try:
            os.remove(potentially_converted_pdf_path)
            logger.debug(
                f"Removed existing file before conversion: {potentially_converted_pdf_path}"
            )
        except OSError as e:
            logger.error(
                f"Error removing existing file {potentially_converted_pdf_path}: {e}"
            )
            return (
                False,
                f"Error removing existing file {potentially_converted_pdf_path}: {e}",
                None,
            )

    command = [
        "soffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        output_dir,
        input_path,
    ]

    logger.info(f"Executing LibreOffice command: {' '.join(command)}")

    try:
        process = subprocess.run(
            command, capture_output=True, text=True, timeout=120, check=False
        )

        if process.returncode == 0:
            if os.path.exists(potentially_converted_pdf_path):
                logger.info(
                    f"LibreOffice successfully converted '{input_path}' to '{potentially_converted_pdf_path}'"
                )
                return True, potentially_converted_pdf_path, expected_pdf_filename
            else:
                error_message = f"LibreOffice exited successfully (code 0) but the expected output PDF was not found: {potentially_converted_pdf_path}. stdout: {process.stdout}, stderr: {process.stderr}"
                logger.error(error_message)
                return False, error_message, None
        else:
            error_message = f"LibreOffice conversion failed for '{input_path}'. Return code: {process.returncode}. stdout: {process.stdout}, stderr: {process.stderr}"
            logger.error(error_message)
            return False, error_message, None

    except FileNotFoundError:
        error_msg = "'soffice' command not found. Please ensure LibreOffice is installed and in your system's PATH."
        logger.error(error_msg)
        return False, error_msg, None
    except subprocess.TimeoutExpired:
        error_msg = f"LibreOffice conversion timed out for '{input_path}'."
        logger.error(error_msg)
        return False, error_msg, None
    except Exception as e:
        error_msg = f"An unexpected error occurred during LibreOffice conversion of '{input_path}': {e}"
        logger.error(error_msg, exc_info=True)
        return False, error_msg, None


def convert_to_pptx(input_path, output_dir, skip_default_content=False):
    """
    Converts a document to PPTX using python-pptx library (primary) or LibreOffice (fallback).

    Args:
        input_path (str): Absolute path to the input document (e.g., a DOCX file).
        output_dir (str): Absolute path to the directory where the PPTX should be saved.
        skip_default_content (bool): If True, skip adding default "文档内容" slide when no content is found.

    Returns:
        tuple: (bool_success, pptx_output_path_or_error_msg, original_pptx_filename_or_None)
               - bool_success: True if conversion was successful, False otherwise.
               - pptx_output_path_or_error_msg: Absolute path to the converted PPTX if successful,
                                               or an error message string if failed.
               - original_pptx_filename_or_None: The original filename of the PPTX as created by LibreOffice.
    """
    print(f"!!! CONVERT_TO_PPTX CALLED !!! {input_path}")  # Force visible log
    logger.error(
        f"FORCE LOG: convert_to_pptx called with {input_path}"
    )  # Force error level log

    if not os.path.exists(input_path):
        return False, f"Input file not found: {input_path}", None
    if not os.path.isdir(output_dir):
        return (
            False,
            f"Output directory not found: {output_dir}",
            None,
        )  # Try Python-based PPTX creation first (more reliable for DOCX to PPTX)
    logger.info(
        f"Attempting PPTX creation using python-pptx for: {input_path}"
    )  # Check if this is a PPT-derived file (should skip default content)
    input_filename = os.path.basename(input_path)
    filename_based_skip = (
        "ppt" in input_filename.lower()
        or "powerpoint" in input_filename.lower()
        or input_filename.lower().endswith((".ppt", ".pptx"))
    )

    # Use parameter if explicitly passed, otherwise use filename-based detection
    final_skip_default_content = skip_default_content or filename_based_skip

    logger.error(
        f"FORCE LOG: skip_default_content param = {skip_default_content}, filename_based = {filename_based_skip}, final = {final_skip_default_content} for filename {input_filename}"
    )
    print(
        f"!!! SKIP_DEFAULT_CONTENT = {final_skip_default_content} for {input_filename} !!!"
    )

    success, result, filename = create_pptx_from_docx(
        input_path, output_dir, final_skip_default_content
    )

    if success:
        return success, result, filename

    # Log the python-pptx failure but continue to LibreOffice fallback
    logger.warning(f"Python-pptx method failed: {result}")

    # LibreOffice fallback - but DOCX to PPTX conversion may not be supported
    logger.info("Attempting LibreOffice fallback for DOCX to PPTX conversion...")

    input_filename_stem = os.path.splitext(os.path.basename(input_path))[0]
    expected_pptx_filename = f"{input_filename_stem}.pptx"
    potentially_converted_pptx_path = os.path.join(output_dir, expected_pptx_filename)

    if os.path.exists(potentially_converted_pptx_path):
        try:
            os.remove(potentially_converted_pptx_path)
            logger.debug(
                f"Removed existing file before conversion: {potentially_converted_pptx_path}"
            )
        except OSError as e:
            logger.error(
                f"Error removing existing file {potentially_converted_pptx_path}: {e}"
            )
            return (
                False,
                f"Error removing existing file {potentially_converted_pptx_path}: {e}",
                None,
            )

    # Try LibreOffice conversion
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        "pptx",
        "--outdir",
        output_dir,
        input_path,
    ]

    logger.info(f"LibreOffice fallback command: {' '.join(command)}")

    try:
        process = subprocess.run(
            command, capture_output=True, text=True, timeout=120, check=False
        )

        if process.returncode == 0:
            if os.path.exists(potentially_converted_pptx_path):
                logger.info(
                    f"LibreOffice successfully converted '{input_path}' to '{potentially_converted_pptx_path}'"
                )
                return True, potentially_converted_pptx_path, expected_pptx_filename
            else:
                error_message = f"LibreOffice exited successfully (code 0) but the expected output PPTX was not found: {potentially_converted_pptx_path}. stdout: {process.stdout}, stderr: {process.stderr}"
                logger.error(error_message)
                # Since both methods failed, return the more informative python-pptx error
                return (
                    False,
                    f"PPTX conversion failed. Python-pptx error: {result}. LibreOffice error: {process.stderr}",
                    None,
                )
        else:
            error_message = f"LibreOffice conversion to PPTX failed for '{input_path}'. Return code: {process.returncode}. stdout: {process.stdout}, stderr: {process.stderr}"
            logger.error(error_message)
            # Return the more informative python-pptx error since LibreOffice also failed
            return (
                False,
                f"PPTX conversion failed. Python-pptx error: {result}. LibreOffice error: {process.stderr}",
                None,
            )

    except FileNotFoundError:
        error_msg = "'soffice' command not found. Please ensure LibreOffice is installed and in your system's PATH."
        logger.error(error_msg)
        return (
            False,
            f"PPTX conversion failed. Python-pptx error: {result}. LibreOffice not found: {error_msg}",
            None,
        )
    except subprocess.TimeoutExpired:
        error_msg = f"LibreOffice conversion to PPTX timed out for '{input_path}'."
        logger.error(error_msg)
        return (
            False,
            f"PPTX conversion failed. Python-pptx error: {result}. LibreOffice timeout: {error_msg}",
            None,
        )
    except Exception as e:
        error_msg = f"An unexpected error occurred during LibreOffice conversion of '{input_path}' to PPTX: {e}"
        logger.error(error_msg, exc_info=True)
        return (
            False,
            f"PPTX conversion failed. Python-pptx error: {result}. LibreOffice error: {error_msg}",
            None,
        )


def create_pptx_from_docx(docx_path, output_dir, skip_default_content=False):
    """
    Creates a PPTX file from a DOCX file using python-pptx library.
    Each page (separated by page break) in the DOCX becomes a separate slide in the PPTX.
    All content (paragraphs, tables, headings) between page breaks is grouped into one slide.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from docx import Document
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.table import CT_Tbl
        from docx.text.paragraph import Paragraph
        from docx.table import Table

        doc = Document(docx_path)
        prs = Presentation()
        input_filename_stem = os.path.splitext(os.path.basename(docx_path))[0]
        expected_pptx_filename = f"{input_filename_stem}.pptx"
        output_path = os.path.join(output_dir, expected_pptx_filename)

        # 1. 按分页分组内容，确保分页符之间所有内容都进同一slide
        def iter_block_items(parent):
            for child in parent.element.body.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, parent)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, parent)

        W_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        slides_content = []
        current_slide_content = []
        for block in iter_block_items(doc):
            is_page_break = False
            if isinstance(block, Paragraph):
                # 检查是否是仅包含分页符的段落（即段落内容为空但有分页符）
                is_page_break = any(
                    run
                    for run in block.runs
                    if any(
                        br.get(
                            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type"
                        )
                        == "page"
                        for br in run._element.findall(".//w:br", namespaces=W_NS)
                    )
                )
            if is_page_break and not (block.text and block.text.strip()):
                # 分割slide，分页符本身不进slide内容
                if current_slide_content:
                    slides_content.append(current_slide_content)
                current_slide_content = []
                continue
            # 段落、表格都要收集
            current_slide_content.append(block)
        if current_slide_content:
            slides_content.append(current_slide_content)

        slide_count = 0
        for slide_items in slides_content:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            slide_count += 1
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame if content_shape else None
            if text_frame:
                text_frame.clear()
            first_title = None
            for idx, item in enumerate(slide_items):
                if isinstance(item, Paragraph):
                    text = item.text.strip()
                    if idx == 0 and title_shape and text:
                        title_shape.text = text
                        first_title = text
                    elif text_frame and text:
                        p = (
                            text_frame.add_paragraph()
                            if text_frame.paragraphs
                            else text_frame.paragraphs[0]
                        )
                        p.text = text
                        p.font.size = Pt(14)
                elif isinstance(item, Table):
                    # 在内容区插入表格（简单方式：转文本）
                    table_text = []
                    for row in item.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        table_text.append("\t".join(row_text))
                    if text_frame:
                        for t in table_text:
                            p = (
                                text_frame.add_paragraph()
                                if text_frame.paragraphs
                                else text_frame.paragraphs[0]
                            )
                            p.text = t
                            p.font.size = Pt(12)
            # 如果没有标题，给个默认标题
            if title_shape and not (first_title and first_title.strip()):
                title_shape.text = f"Slide {slide_count}"
        if slide_count == 0 and not skip_default_content:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = "文档内容"
            if slide.placeholders[1]:
                slide.placeholders[1].text = "从图片中提取的文本内容"
        prs.save(output_path)
        logger.info(f"Successfully created PPTX using python-pptx: {output_path}")
        return True, output_path, expected_pptx_filename
    except ImportError as e:
        error_msg = f"Required libraries not available for PPTX creation: {e}. Please install python-pptx and python-docx."
        logger.error(error_msg)
        return False, error_msg, None
    except Exception as e:
        error_msg = f"Error creating PPTX from DOCX '{docx_path}': {e}"
        logger.error(error_msg, exc_info=True)
        return False, error_msg, None


if __name__ == "__main__":
    # Example usage (for testing this script directly)
    # Create dummy files and dirs for testing
    test_output_dir = "./test_output_lo"
    test_input_file = "./test_input_lo.docx"

    if not os.path.exists(test_output_dir):
        os.makedirs(test_output_dir)

    # Create a simple DOCX file for testing if it doesn't exist
    if not os.path.exists(test_input_file):
        try:
            from docx import (
                Document as DocxDocument,
            )  # Use a different alias to avoid confusion if Document is used elsewhere

            doc = DocxDocument()
            doc.add_paragraph(
                "This is a test docx for LibreOffice conversion created by script."
            )
            doc.save(test_input_file)
            logger.info(f"Created dummy test file: {test_input_file}")
        except ImportError:
            logger.warning(
                "python-docx library is not installed. Cannot create a dummy .docx file for testing. Please create it manually."
            )
            # As a very basic fallback, create a text file that soffice might still process or error out on
            with open(test_input_file, "w") as f:
                f.write("This is a test docx for LibreOffice conversion (plain text).")
        except Exception as e_create:
            logger.error(
                f"Failed to create dummy test file {test_input_file}: {e_create}"
            )

    print("Testing LibreOffice converter...")
    if os.path.exists(test_input_file):  # Only test if input file exists
        success, result_path_or_msg, _ = convert_to_pdf(
            os.path.abspath(test_input_file), os.path.abspath(test_output_dir)
        )

        if success:
            print(f"Conversion successful. PDF at: {result_path_or_msg}")
            # Optional: Clean up test files by uncommenting below
            # print(f"To cleanup, manually remove: {test_input_file}, {result_path_or_msg}, and directory {test_output_dir}")
            # os.remove(test_input_file)
            # os.remove(result_path_or_msg)
            # shutil.rmtree(test_output_dir)
        else:
            print(f"Conversion failed. Error: {result_path_or_msg}")
    else:
        print(
            f"Skipping test, input file {test_input_file} does not exist or could not be created."
        )
