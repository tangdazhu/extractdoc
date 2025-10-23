# -*- coding: utf-8 -*-
"""
商务风格PPT生成器

专业的商务设计风格，动态生成精美的PPT页面
蓝色配色方案，适合商务汇报、产品介绍
"""
import logging
from typing import List, Dict, Any, Optional
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from utils.config_manager import config
from .base_ppt_generator import BasePPTGenerator

logger = logging.getLogger(__name__)


class BusinessStylePPTGenerator(BasePPTGenerator):
    """
    商务风格PPT生成器
    
    设计特点:
    - 蓝色渐变背景
    - 白色内容区域
    - 圆角装饰元素
    - 清晰的视觉层次
    """
    
    # 配色方案
    COLOR_PRIMARY_DARK = RGBColor(1, 93, 187)  # 深蓝
    COLOR_PRIMARY_LIGHT = RGBColor(100, 180, 255)  # 浅蓝
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(50, 50, 50)
    COLOR_TEXT_LIGHT = RGBColor(200, 200, 200)
    
    def __init__(self):
        """初始化生成器"""
        super().__init__()  # 调用基类初始化
        logger.info("初始化商务风格PPT生成器")
    
    
    def _add_gradient_background(self, slide, angle=135.0):
        """添加渐变背景"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = angle
        fill.gradient_stops[0].color.rgb = self.COLOR_PRIMARY_DARK
        fill.gradient_stops[1].color.rgb = self.COLOR_PRIMARY_LIGHT
    
    def _add_decorative_dots(self, slide, x=11, y=0.5, rows=6, cols=6):
        """添加装饰圆点"""
        for i in range(cols):
            for j in range(rows):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x + i*0.15), Inches(y + j*0.15),
                    Inches(0.08), Inches(0.08)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = self.COLOR_WHITE
                dot.fill.fore_color.brightness = -0.5
                dot.line.fill.background()
    
    def create_cover_slide(self, title: str, subtitle: str = "", 
                          reporter: str = "", date: str = ""):
        """
        创建封面页
        
        Args:
            title: 主标题
            subtitle: 副标题
            reporter: 汇报人
            date: 日期
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        self._add_gradient_background(slide)
        self._add_decorative_dots(slide)
        
        # 从配置读取封面标题参数
        cover_config = config.get("ppt_generation.cover_title", {})
        max_font_size = cover_config.get("max_font_size", 66)
        min_font_size = cover_config.get("min_font_size", 36)
        max_chars_per_line = cover_config.get("max_chars_per_line", 20)
        text_box_width = cover_config.get("text_box_width", 10.0)
        text_box_height = cover_config.get("text_box_height", 1.5)
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5),
            Inches(text_box_width), Inches(text_box_height)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True  # 启用自动换行
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 根据标题长度动态调整字体大小
        title_length = len(title)
        if title_length > max_chars_per_line * 2:
            # 超长标题：使用最小字体
            font_size = min_font_size
        elif title_length > max_chars_per_line:
            # 中等长度：使用中等字体
            font_size = (max_font_size + min_font_size) // 2
        else:
            # 短标题：使用最大字体
            font_size = max_font_size
        
        title_frame.paragraphs[0].font.size = Pt(font_size)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        logger.debug(f"封面标题长度={title_length}，使用字体大小={font_size}pt")
        
        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(2), Inches(4.2),
                Inches(9.33), Inches(0.6)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # 汇报人信息框
        if reporter:
            reporter_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(2.5), Inches(6),
                Inches(3), Inches(0.4)
            )
            reporter_shape.fill.solid()
            reporter_shape.fill.fore_color.rgb = self.COLOR_WHITE
            reporter_shape.fill.fore_color.brightness = -0.3
            reporter_shape.line.fill.background()
            
            reporter_text = reporter_shape.text_frame
            reporter_text.text = f"汇报人: {reporter}"
            reporter_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            reporter_text.paragraphs[0].font.size = Pt(18)
            reporter_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # 日期信息框
        if date:
            date_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7.5), Inches(6),
                Inches(3), Inches(0.4)
            )
            date_shape.fill.solid()
            date_shape.fill.fore_color.rgb = self.COLOR_WHITE
            date_shape.fill.fore_color.brightness = -0.3
            date_shape.line.fill.background()
            
            date_text = date_shape.text_frame
            date_text.text = f"日期: {date}"
            date_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            date_text.paragraphs[0].font.size = Pt(18)
            date_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        logger.info(f"创建封面页: {title}")
    
    def create_content_slide(self, title: str, content_lines: List[str]):
        """
        创建内容页
        
        Args:
            title: 页面标题
            content_lines: 内容行列表
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # 使用基类方法创建标题栏
        self._create_title_bar(slide, 0, 1.2, title, 36, self.COLOR_PRIMARY_DARK)
        
        # 使用基类方法创建内容框
        content_box = self._create_content_box(slide, 0.8, 1.8, 11.73, 5)
        
        # 内容文字（使用基类统一方法）
        content_text = content_box.text_frame
        self._setup_text_frame_margins(content_text)
        
        # 从配置读取字体大小
        font_sizes = {
            0: config.get("text_formatting.font_size_level_0", 20),
            1: config.get("text_formatting.font_size_level_1", 18),
            2: config.get("text_formatting.font_size_level_2", 16)
        }
        
        # 使用基类统一方法添加内容
        self._add_bullet_content(content_text, content_lines, font_sizes, self.COLOR_TEXT_DARK)
        
        logger.debug(f"创建内容页: {title} ({len(content_lines)}行)")
    
    def create_two_column_slide(self, title: str, left_content: List[str], right_content: List[str], 
                                left_title: str = "\u4f20\u7edf\u65b9\u5f0f", right_title: str = "AI\u65b9\u5f0f"):
        """
        \u521b\u5efa\u5de6\u53f3\u5bf9\u6bd4\u9875
        
        Args:
            title: \u9875\u9762\u6807\u9898
            left_content: \u5de6\u4fa7\u5185\u5bb9\u5217\u8868
            right_content: \u53f3\u4fa7\u5185\u5bb9\u5217\u8868
            left_title: \u5de6\u4fa7\u6807\u9898
            right_title: \u53f3\u4fa7\u6807\u9898
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # \u4f7f\u7528\u57fa\u7c7b\u65b9\u6cd5\u521b\u5efa\u6807\u9898\u680f
        self._create_title_bar(slide, 0, 1.2, title, 36, self.COLOR_PRIMARY_DARK)
        
        # 从配置读取分栏比例
        split_ratio = config.get("ppt_generation.layout_types.two_column.split_ratio", 0.5)
        
        # 计算左右宽度
        total_width = 11.73
        gap = 0.5
        left_width = (total_width - gap) * split_ratio
        right_width = total_width - gap - left_width
        
        # 使用基类方法创建左侧内容区
        left_box = self._create_content_box(slide, 0.8, 1.8, left_width, 5, self.COLOR_PRIMARY_DARK, 2)
        left_box.fill.fore_color.brightness = -0.05
        
        # 左侧标题
        left_title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(left_width), Inches(0.4)
        )
        left_title_text = left_title_box.text_frame
        left_title_text.text = left_title
        left_title_text.paragraphs[0].font.size = Pt(20)
        left_title_text.paragraphs[0].font.bold = True
        left_title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        left_title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 左侧内容
        left_text = left_box.text_frame
        left_text.word_wrap = True
        left_text.margin_left = Inches(0.3)
        left_text.margin_right = Inches(0.3)
        left_text.margin_top = Inches(0.5)
        
        for i, line in enumerate(left_content):
            if i > 0:
                left_text.add_paragraph()
            para = left_text.paragraphs[i]
            para.text = f"● {line}"
            para.font.size = Pt(16)
            para.font.color.rgb = self.COLOR_TEXT_DARK
            para.space_after = Pt(8)
        
        # 使用基类方法创建右侧内容区
        right_box = self._create_content_box(slide, 0.8 + left_width + gap, 1.8, right_width, 5, self.COLOR_PRIMARY_DARK, 2)
        right_box.fill.fore_color.brightness = -0.05
        
        # 右侧标题
        right_title_box = slide.shapes.add_textbox(
            Inches(0.8 + left_width + gap), Inches(1.5),
            Inches(right_width), Inches(0.4)
        )
        right_title_text = right_title_box.text_frame
        right_title_text.text = right_title
        right_title_text.paragraphs[0].font.size = Pt(20)
        right_title_text.paragraphs[0].font.bold = True
        right_title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        right_title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 右侧内容
        right_text = right_box.text_frame
        right_text.word_wrap = True
        right_text.margin_left = Inches(0.3)
        right_text.margin_right = Inches(0.3)
        right_text.margin_top = Inches(0.5)
        
        for i, line in enumerate(right_content):
            if i > 0:
                right_text.add_paragraph()
            para = right_text.paragraphs[i]
            para.text = f"● {line}"
            para.font.size = Pt(16)
            para.font.color.rgb = self.COLOR_TEXT_DARK
            para.space_after = Pt(8)
        
        logger.debug(f"创建左右对比页: {title}")
    
    def create_three_column_slide(self, title: str, cards: List[Dict[str, str]]):
        """
        创建三列卡片页
        
        Args:
            title: 页面标题
            cards: 卡片列表，每个卡片包含 {"icon": "图标", "title": "标题", "content": "内容"}
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # 使用基类方法创建标题栏
        self._create_title_bar(slide, 0, 1.2, title, 36, self.COLOR_PRIMARY_DARK)
        
        # 从基类获取配置
        cfg = self._get_three_column_config()
        max_cards = cfg["max_cards"]
        card_width = cfg["card_width"]
        card_gap = cfg["card_gap"]
        card_title_font_size = cfg["card_title_font_size"]
        card_content_font_size = cfg["card_content_font_size"]
        card_content_max_chars = cfg["card_content_max_chars"]
        
        start_x = 1.0
        
        # 限制显示数量
        cards_to_show = cards[:max_cards]
        
        for i, card in enumerate(cards_to_show):
            x_pos = start_x + i * (card_width + card_gap)
            
            # 使用基类方法创建卡片容器
            card_box = self._create_content_box(slide, x_pos, 2.0, card_width, 4.5, self.COLOR_PRIMARY_DARK, 2)
            card_box.fill.fore_color.brightness = -0.05
            
            # 图标/编号（大号圆形）
            icon_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos + 1.0), Inches(2.5),
                Inches(1.5), Inches(1.5)
            )
            icon_circle.fill.solid()
            icon_circle.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            icon_circle.line.fill.background()
            
            icon_text = icon_circle.text_frame
            icon_text.text = card.get("icon", f"{i+1}")
            icon_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            icon_text.paragraphs[0].font.size = Pt(48)
            icon_text.paragraphs[0].font.bold = True
            icon_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            icon_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 卡片标题
            card_title_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.3), Inches(4.2),
                Inches(card_width - 0.6), Inches(0.5)
            )
            card_title_text = card_title_box.text_frame
            card_title_text.text = card.get("title", "")
            card_title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            card_title_text.paragraphs[0].font.size = Pt(card_title_font_size)
            card_title_text.paragraphs[0].font.bold = True
            card_title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
            card_title_text.word_wrap = True
            
            # 卡片内容
            card_content_box = slide.shapes.add_textbox(
                Inches(x_pos + 0.2), Inches(4.8),
                Inches(card_width - 0.4), Inches(1.8)
            )
            card_content_text = card_content_box.text_frame
            # 使用基类的智能截断方法
            content = self._truncate_text_smart(card.get("content", ""), card_content_max_chars)
            card_content_text.text = content
            card_content_text.word_wrap = True
            card_content_text.paragraphs[0].alignment = PP_ALIGN.LEFT
            card_content_text.paragraphs[0].font.size = Pt(card_content_font_size)
            card_content_text.paragraphs[0].font.color.rgb = self.COLOR_TEXT_DARK
            card_content_text.paragraphs[0].line_spacing = 1.2
        
        logger.debug(f"创建三列卡片页: {title}, {len(cards_to_show)}张卡片")
    
    def create_flow_diagram_slide(self, title: str, steps: List[Dict[str, str]]):
        """
        创建流程图页
        
        Args:
            title: 页面标题
            steps: 步骤列表，每个步骤包含 {"title": "步骤名", "description": "说明"}
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # 使用基类方法创建标题栏
        self._create_title_bar(slide, 0, 1.2, title, 36, self.COLOR_PRIMARY_DARK)
        
        # 从基类获取配置并计算布局
        cfg = self._get_flow_diagram_config()
        steps_to_show = steps[:cfg["max_steps"]]
        step_count = len(steps_to_show)
        
        step_width, arrow_width, total_width, step_title_font_size, step_desc_font_size, step_desc_max_chars = self._calculate_flow_diagram_layout(step_count, cfg)
        
        start_x = (13.33 - total_width) / 2
        
        for i, step in enumerate(steps_to_show):
            x_pos = start_x + i * (step_width + arrow_width)
            
            # 步骤框
            step_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(2.5),
                Inches(step_width), Inches(1.5)
            )
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            step_box.line.fill.background()
            
            # 步骤标题
            step_title_text = step_box.text_frame
            step_title_text.text = step.get("title", f"步骤{i+1}")
            step_title_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            step_title_text.paragraphs[0].font.size = Pt(step_title_font_size)
            step_title_text.paragraphs[0].font.bold = True
            step_title_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            step_title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            step_title_text.word_wrap = True
            
            # 步骤说明
            desc_box = slide.shapes.add_textbox(
                Inches(x_pos), Inches(4.2),
                Inches(step_width), Inches(2.0)
            )
            desc_text = desc_box.text_frame
            # 使用基类的智能截断方法
            description = self._truncate_text_smart(step.get("description", ""), step_desc_max_chars)
            desc_text.text = description
            desc_text.word_wrap = True
            desc_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_text.paragraphs[0].font.size = Pt(step_desc_font_size)
            desc_text.paragraphs[0].font.color.rgb = self.COLOR_TEXT_DARK
            
            # 箭头（除了最后一个步骤）
            if i < step_count - 1:
                arrow_x = x_pos + step_width
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(arrow_x), Inches(2.9),
                    Inches(arrow_width), Inches(0.7)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.COLOR_PRIMARY_LIGHT
                arrow.line.fill.background()
        
        logger.debug(f"创建流程图页: {title}, {step_count}个步骤")
    
    def create_timeline_slide(self, title: str, timeline_items: List[Dict[str, str]]):
        """
        创建时间线页
        
        Args:
            title: 页面标题
            timeline_items: 时间线项目列表，每项包含 {"title": "标题", "content": "内容"}
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # 使用基类方法创建标题栏
        self._create_title_bar(slide, 0, 1.2, title, 36, self.COLOR_PRIMARY_DARK)
        
        # 从基类获取配置并计算布局
        cfg = self._get_timeline_config()
        items_to_show = timeline_items[:cfg["max_items"]]
        item_count = len(items_to_show)
        
        item_height, title_font_size, content_font_size, content_max_chars = self._calculate_timeline_layout(item_count, cfg)
        start_y = cfg["start_y"]
        
        # 垂直时间线
        line_x = 3.0
        
        # 绘制时间线主线
        for i in range(len(items_to_show)):
            y_pos = start_y + i * item_height
            
            # 节点圆圈
            node_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(line_x - 0.15), Inches(y_pos),
                Inches(0.3), Inches(0.3)
            )
            node_circle.fill.solid()
            node_circle.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            node_circle.line.fill.background()
            
            # 连接线（除了最后一个）
            if i < len(items_to_show) - 1:
                connector = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(line_x - 0.02), Inches(y_pos + 0.3),
                    Inches(0.04), Inches(item_height - 0.3)
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = self.COLOR_PRIMARY_LIGHT
                connector.line.fill.background()
            
            # 使用基类方法创建内容框（动态高度）
            item = items_to_show[i]
            box_height = item_height * 0.8
            content_box = self._create_content_box(slide, line_x + 0.5, y_pos - 0.1, 9, box_height, self.COLOR_PRIMARY_DARK, 1)
            content_box.fill.fore_color.brightness = -0.05
            
            # 内容文字
            content_text = content_box.text_frame
            content_text.margin_left = Inches(0.3)
            content_text.margin_top = Inches(0.1)
            
            # 标题
            content_text.text = item.get("title", "")
            content_text.paragraphs[0].font.size = Pt(title_font_size)
            content_text.paragraphs[0].font.bold = True
            content_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
            
            # 内容
            if item.get("content"):
                content_text.add_paragraph()
                # 使用基类的智能截断方法
                content = self._truncate_text_smart(item.get("content", ""), content_max_chars)
                content_text.paragraphs[1].text = content
                content_text.paragraphs[1].font.size = Pt(content_font_size)
                content_text.paragraphs[1].font.color.rgb = self.COLOR_TEXT_DARK
            
            # 启用自动换行
            content_text.word_wrap = True
        
        logger.debug(f"创建时间线页: {title}, {len(items_to_show)}项")
    
    def create_section_slide(self, number: str, title: str):
        """
        创建章节分隔页
        
        Args:
            number: 章节编号（如 "01"）
            title: 章节标题
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_decorative_dots(slide, x=1, y=1)
        
        # 大号数字
        number_box = slide.shapes.add_textbox(
            Inches(3), Inches(2),
            Inches(7.33), Inches(2)
        )
        number_frame = number_box.text_frame
        number_frame.text = number
        number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        number_frame.paragraphs[0].font.size = Pt(120)
        number_frame.paragraphs[0].font.bold = True
        number_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        # 章节标题
        section_title = slide.shapes.add_textbox(
            Inches(3), Inches(4.2),
            Inches(7.33), Inches(1)
        )
        section_frame = section_title.text_frame
        section_frame.text = title
        section_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        section_frame.paragraphs[0].font.size = Pt(48)
        section_frame.paragraphs[0].font.bold = True
        section_frame.paragraphs[0].font.color.rgb = self.COLOR_WHITE
        
        logger.debug(f"创建章节页: {number} - {title}")
    
    def create_catalog_slide(self, catalog_items: List[Dict[str, str]]):
        """
        创建目录页（动态高度，单列布局）
        
        Args:
            catalog_items: 目录项列表，每项包含 {"number": "01", "title": "标题"}
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        self._add_decorative_dots(slide, x=10.5, y=0.3)
        self._add_decorative_dots(slide, x=10.5, y=5.5)
        
        # 左侧标题区域
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.7),
            Inches(2.5), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = "目录"
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        
        # "CONTENTS" 副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(2.8), Inches(1.1),
            Inches(3), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "CONTENTS"
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = self.COLOR_TEXT_LIGHT
        
        # 右侧配图（圆形）
        pic_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.5), Inches(1),
            Inches(4), Inches(4)
        )
        pic_circle.fill.solid()
        pic_circle.fill.fore_color.rgb = self.COLOR_WHITE
        pic_circle.fill.fore_color.brightness = -0.2
        pic_circle.line.color.rgb = self.COLOR_PRIMARY_DARK
        pic_circle.line.width = Pt(3)
        
        # 从基类获取配置并计算布局
        cfg = self._get_catalog_config()
        items_to_show = catalog_items[:cfg["max_items"]]
        total_items = len(items_to_show)
        
        item_height, number_font_size, title_font_size = self._calculate_catalog_layout(total_items, cfg)
        start_y = cfg["start_y"]
        
        for i, item in enumerate(items_to_show):
            number = item.get("number", f"{i+1:02d}")
            title = item.get("title", "")
            
            y_pos = start_y + i * item_height
            
            # 编号框（蓝色圆角矩形）
            number_box_height = item_height * 0.7  # 编号框占70%高度
            number_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.95), Inches(y_pos),
                Inches(0.7), Inches(number_box_height)
            )
            number_shape.fill.solid()
            number_shape.fill.fore_color.rgb = self.COLOR_PRIMARY_DARK
            number_shape.line.fill.background()
            
            number_text = number_shape.text_frame
            number_text.text = number
            number_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            number_text.paragraphs[0].font.size = Pt(number_font_size)
            number_text.paragraphs[0].font.bold = True
            number_text.paragraphs[0].font.color.rgb = self.COLOR_WHITE
            number_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 标题框（白色圆角矩形）
            title_box_height = item_height * 0.7
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.75), Inches(y_pos),
                Inches(5.5), Inches(title_box_height)
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = self.COLOR_WHITE
            title_shape.line.color.rgb = self.COLOR_PRIMARY_DARK
            title_shape.line.width = Pt(1)
            
            title_text = title_shape.text_frame
            title_text.text = title
            title_text.paragraphs[0].font.size = Pt(title_font_size)
            title_text.paragraphs[0].font.color.rgb = self.COLOR_TEXT_DARK
            title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            title_text.margin_left = Inches(0.2)
        
        logger.info(f"创建目录页: {total_items}项，每项高度{item_height:.2f}英寸")
    
    def create_picture_slide(self, title: str, image_path: str, caption: str = ""):
        """
        创建图片页（左右布局：左侧图片，右侧文字说明）
        
        Args:
            title: 页面标题
            image_path: 图片路径
            caption: 图片说明（文字内容）
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_gradient_background(slide)
        
        # 标题栏
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(13.33), Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLOR_WHITE
        title_bar.line.fill.background()
        
        title_text = title_bar.text_frame
        title_text.text = title
        title_text.paragraphs[0].font.size = Pt(36)
        title_text.paragraphs[0].font.bold = True
        title_text.paragraphs[0].font.color.rgb = self.COLOR_PRIMARY_DARK
        title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_text.margin_left = Inches(0.5)
        
        # 左侧：图片容器
        pic_container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.8),
            Inches(6.5), Inches(5.2)
        )
        pic_container.fill.solid()
        pic_container.fill.fore_color.rgb = self.COLOR_WHITE
        pic_container.line.fill.background()
        
        # 插入图片
        try:
            pic = slide.shapes.add_picture(
                image_path,
                Inches(1), Inches(2),
                width=Inches(6)
            )
            # 调整图片大小以适应容器
            if pic.height > Inches(4.8):
                ratio = Inches(4.8) / pic.height
                pic.height = Inches(4.8)
                pic.width = int(pic.width * ratio)
            
            # 居中图片
            pic.left = Inches(4.05) - int(pic.width / 2)
            pic.top = Inches(4.4) - int(pic.height / 2)
            
        except Exception as e:
            logger.error(f"插入图片失败: {e}")
        
        # 右侧：文字说明区域
        if caption:
            text_container = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(7.5), Inches(1.8),
                Inches(5.5), Inches(5.2)
            )
            text_container.fill.solid()
            text_container.fill.fore_color.rgb = self.COLOR_WHITE
            text_container.line.fill.background()
            
            text_frame = text_container.text_frame
            self._setup_text_frame_margins(text_frame)
            text_frame.margin_bottom = Inches(0.3)
            
            # 使用基类统一方法添加内容
            lines = caption.split('\n') if caption else []
            font_sizes = {0: 14, 1: 12}
            self._add_bullet_content(text_frame, lines, font_sizes, self.COLOR_TEXT_DARK)
        
        logger.debug(f"创建图片页: {title}")
    
    def save(self, output_path: str):
        """保存PPT"""
        self.prs.save(output_path)
        logger.info(f"PPT已保存: {output_path}")
