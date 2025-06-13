from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os


def copy_images_to_pptx(image_paths, output_pptx_path):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # 空白页

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        img = Image.open(img_path)
        width, height = img.size

        # 获取幻灯片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 计算缩放比例，保证图片完整显示且居中
        scale = min(slide_width / width, slide_height / height)
        img_width = int(width * scale)
        img_height = int(height * scale)
        left = int((slide_width - img_width) / 2)
        top = int((slide_height - img_height) / 2)

        slide.shapes.add_picture(
            img_path, left, top, width=img_width, height=img_height
        )

    prs.save(output_pptx_path)
    return output_pptx_path
