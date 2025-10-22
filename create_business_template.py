# -*- coding: utf-8 -*-
"""
基于Kimi设计创建专业的占位符模板
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_professional_template():
    """创建专业模板"""
    
    # 创建新演示文稿（16:9比例）
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9比例
    prs.slide_height = Inches(7.5)
    
    print("创建专业PPT模板...")
    print(f"幻灯片尺寸: {prs.slide_width} x {prs.slide_height}")
    
    # 获取母版
    master = prs.slide_master
    
    # 设置母版背景（蓝色渐变）
    background = master.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = RGBColor(1, 93, 187)  # 深蓝
    fill.gradient_stops[1].color.rgb = RGBColor(100, 180, 255)  # 浅蓝
    
    print("\n创建布局...")
    
    # ========== 布局1: 封面页 ==========
    print("1. 封面布局 (Title Slide)")
    title_slide_layout = master.slide_layouts[0]
    
    # 清空现有占位符
    for shape in list(title_slide_layout.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # 添加标题占位符（大标题，居中）
    title_placeholder = title_slide_layout.shapes.add_placeholder(
        PP_PLACEHOLDER.CENTER_TITLE,
        Inches(1.5), Inches(2.5),  # 位置
        Inches(10.33), Inches(1.5)  # 大小
    )
    
    # 添加副标题占位符
    subtitle_placeholder = title_slide_layout.shapes.add_placeholder(
        PP_PLACEHOLDER.SUBTITLE,
        Inches(2), Inches(4.2),
        Inches(9.33), Inches(0.8)
    )
    
    # 添加页脚占位符（汇报人）
    footer_left = title_slide_layout.shapes.add_placeholder(
        PP_PLACEHOLDER.FOOTER,
        Inches(2.5), Inches(6),
        Inches(3), Inches(0.4)
    )
    
    # 添加日期占位符
    date_placeholder = title_slide_layout.shapes.add_placeholder(
        PP_PLACEHOLDER.DATE,
        Inches(7.5), Inches(6),
        Inches(3), Inches(0.4)
    )
    
    # ========== 布局2: 内容页 ==========
    print("2. 内容布局 (Title and Content)")
    content_layout = master.slide_layouts[1]
    
    # ========== 布局3: 章节页 ==========
    print("3. 章节布局 (Section Header)")
    section_layout = master.slide_layouts[2]
    
    # ========== 布局4: 图片页 ==========
    print("4. 图片布局 (Picture with Caption)")
    # 使用现有的图片布局
    
    # 保存模板
    output_path = 'config/templates/professional_template.pptx'
    prs.save(output_path)
    print(f"\n✅ 模板已保存: {output_path}")
    
    return output_path

def create_simple_professional_template():
    """
    创建简化版专业模板
    使用形状模拟Kimi的设计风格
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    print("创建简化版专业模板...")
    
    # ========== 创建封面示例页 ==========
    print("\n1. 创建封面示例页")
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景渐变
    background = slide1.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = RGBColor(1, 93, 187)  # 深蓝
    fill.gradient_stops[1].color.rgb = RGBColor(100, 180, 255)  # 浅蓝
    
    # 添加装饰圆点（右上角）
    for i in range(6):
        for j in range(6):
            dot = slide1.shapes.add_shape(
                1,  # 圆形
                Inches(11 + i*0.15), Inches(0.5 + j*0.15),
                Inches(0.08), Inches(0.08)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(255, 255, 255)
            dot.fill.fore_color.brightness = -0.5
            dot.line.fill.background()
    
    # 主标题文本框
    title_box = slide1.shapes.add_textbox(
        Inches(1.5), Inches(2.5),
        Inches(10), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "演示文稿标题"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(66)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 副标题
    subtitle_box = slide1.shapes.add_textbox(
        Inches(2), Inches(4.2),
        Inches(9.33), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "副标题或描述信息"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 汇报人信息框（圆角矩形）
    reporter_shape = slide1.shapes.add_shape(
        5,  # 圆角矩形
        Inches(2.5), Inches(6),
        Inches(3), Inches(0.4)
    )
    reporter_shape.fill.solid()
    reporter_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    reporter_shape.fill.fore_color.brightness = -0.3
    reporter_shape.line.fill.background()
    
    reporter_text = reporter_shape.text_frame
    reporter_text.text = "汇报人: XXX"
    reporter_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    reporter_text.paragraphs[0].font.size = Pt(18)
    reporter_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 日期信息框
    date_shape = slide1.shapes.add_shape(
        5,  # 圆角矩形
        Inches(7.5), Inches(6),
        Inches(3), Inches(0.4)
    )
    date_shape.fill.solid()
    date_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    date_shape.fill.fore_color.brightness = -0.3
    date_shape.line.fill.background()
    
    date_text = date_shape.text_frame
    date_text.text = "日期: 2025/01/01"
    date_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    date_text.paragraphs[0].font.size = Pt(18)
    date_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # ========== 创建内容页示例 ==========
    print("2. 创建内容页示例")
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide2.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = RGBColor(1, 93, 187)
    fill.gradient_stops[1].color.rgb = RGBColor(100, 180, 255)
    
    # 标题栏（白色背景）
    title_bar = slide2.shapes.add_shape(
        1,  # 矩形
        0, 0,
        Inches(13.33), Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(255, 255, 255)
    title_bar.line.fill.background()
    
    # 标题文字
    title_text = title_bar.text_frame
    title_text.text = "内容页标题"
    title_text.paragraphs[0].font.size = Pt(36)
    title_text.paragraphs[0].font.bold = True
    title_text.paragraphs[0].font.color.rgb = RGBColor(1, 93, 187)
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 内容区域（白色背景）
    content_box = slide2.shapes.add_shape(
        1,  # 矩形
        Inches(0.8), Inches(1.8),
        Inches(11.73), Inches(5)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box.line.fill.background()
    
    # 内容文字
    content_text = content_box.text_frame
    content_text.text = "• 要点1\n• 要点2\n• 要点3"
    content_text.paragraphs[0].font.size = Pt(24)
    content_text.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
    
    # ========== 创建章节页示例 ==========
    print("3. 创建章节页示例")
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide3.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].color.rgb = RGBColor(1, 93, 187)
    fill.gradient_stops[1].color.rgb = RGBColor(100, 180, 255)
    
    # 大号数字
    number_box = slide3.shapes.add_textbox(
        Inches(3), Inches(2),
        Inches(7.33), Inches(2)
    )
    number_frame = number_box.text_frame
    number_frame.text = "01"
    number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    number_frame.paragraphs[0].font.size = Pt(120)
    number_frame.paragraphs[0].font.bold = True
    number_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 章节标题
    section_title = slide3.shapes.add_textbox(
        Inches(3), Inches(4.2),
        Inches(7.33), Inches(1)
    )
    section_frame = section_title.text_frame
    section_frame.text = "章节标题"
    section_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    section_frame.paragraphs[0].font.size = Pt(48)
    section_frame.paragraphs[0].font.bold = True
    section_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 保存
    output_path = 'config/templates/professional_blue_template.pptx'
    prs.save(output_path)
    print(f"\n[OK] Professional template saved: {output_path}")
    
    return output_path

if __name__ == '__main__':
    # 创建简化版专业模板
    template_path = create_simple_professional_template()
    print(f"\nTemplate created successfully!")
    print(f"Path: {template_path}")
    print(f"\nPlease set this template as default in config file")
