"""文档生成服务模块。

负责加载模板配置、整理输入内容，并生成 PPT / Word 文件。
"""

from __future__ import annotations

import json
import logging
import os
import re
import io
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import requests
from django.conf import settings
from docx import Document as WordDocument
from docx.shared import Pt
from PIL import Image
import dashscope
from dashscope import Generation

from .ai_document_analyzer import AIDocumentAnalyzer
from .smart_ppt_generator import SmartPPTGenerator
from .ocr_table_extractor import get_ocr_extractor

logger = logging.getLogger("converter")


def _normalize_table_structure(table: List[List[str]]) -> List[List[str]]:
    """
    规范化表格结构,确保所有行的列数一致
    
    Args:
        table: 原始表格数据
        
    Returns:
        规范化后的表格
    """
    if not table:
        return []
    
    # 找出最大列数
    max_cols = max(len(row) for row in table)
    
    # 补齐所有行到最大列数
    normalized_table = []
    for row_idx, row in enumerate(table):
        if len(row) < max_cols:
            # 补空字符串
            normalized_row = row + [''] * (max_cols - len(row))
            logger.debug(f"规范化第{row_idx}行: {len(row)}列 → {max_cols}列")
        else:
            normalized_row = row
        normalized_table.append(normalized_row)
    
    return normalized_table


def _align_ocr_row_to_target(
    ocr_row: List[str],
    target_row: List[str],
    header_row: Optional[List[str]] = None
) -> List[str]:
    """
    智能对齐OCR提取的行到目标列数
    
    策略:
    1. 识别日期列(格式如2025-01-25),保留到最后一列
    2. 识别第一列(版本号/序号),保留不变
    3. 合并中间的空列或内容较短的列
    
    Args:
        ocr_row: OCR提取的行数据
        target_row: 目标行(用于确定列数)
        header_row: 表头行(可选,用于辅助判断)
        
    Returns:
        对齐后的行数据
    """
    target_cols = len(target_row)
    ocr_cols = len(ocr_row)
    
    if ocr_cols <= target_cols:
        # 列数相同或更少,直接补齐
        return ocr_row + [''] * (target_cols - ocr_cols)
    
    # OCR列数更多,需要智能合并
    # 策略1: 查找日期列(格式: YYYY-MM-DD 或 YYYY-MM)
    import re
    date_pattern = re.compile(r'^\d{4}-\d{2}(-\d{2})?$')
    date_col_idx = -1
    for idx, cell in enumerate(ocr_row):
        if date_pattern.match(cell.strip()):
            date_col_idx = idx
            break
    
    # 策略2: 构建对齐后的行
    aligned_row = []
    
    # 保留第一列(版本号/序号)
    aligned_row.append(ocr_row[0])
    
    # 如果找到日期列,特殊处理
    if date_col_idx > 0:
        # 中间列: 合并第1列到日期列之前的所有列(除了第一列)
        middle_cols = [ocr_row[i] for i in range(1, date_col_idx) if ocr_row[i].strip()]
        
        # 根据目标列数分配
        if target_cols == 5:  # 版本、内容、团队、校核、时间
            if len(middle_cols) >= 2:
                aligned_row.append(middle_cols[0])  # 内容
                aligned_row.append(middle_cols[1] if len(middle_cols) > 1 else '')  # 团队
                aligned_row.append(middle_cols[2] if len(middle_cols) > 2 else '')  # 校核
            elif len(middle_cols) == 1:
                aligned_row.append(middle_cols[0])  # 内容
                aligned_row.append('')  # 团队
                aligned_row.append('')  # 校核
            else:
                aligned_row.extend(['', '', ''])
        else:
            # 其他列数: 简单合并
            aligned_row.extend(middle_cols)
            aligned_row.extend([''] * (target_cols - len(aligned_row) - 1))
        
        # 最后一列是日期
        aligned_row.append(ocr_row[date_col_idx])
    else:
        # 没有日期列: 保留前target_cols-1列,合并剩余列到最后一列
        aligned_row.extend(ocr_row[1:target_cols-1])
        aligned_row.append(' '.join(ocr_row[target_cols-1:]))
    
    # 确保列数正确
    if len(aligned_row) < target_cols:
        aligned_row.extend([''] * (target_cols - len(aligned_row)))
    elif len(aligned_row) > target_cols:
        aligned_row = aligned_row[:target_cols]
    
    return aligned_row


try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx 未安装，PPT 生成功能不可用。")

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("PyMuPDF 未安装，PDF 转图片功能不可用。")

TEMPLATES_CONFIG_PATH = Path(settings.BASE_DIR).parent / "config" / "document_generation_templates.json"


def load_generation_templates() -> Dict[str, Dict[str, dict]]:
    """加载文档生成模板配置。

    返回结构:
        {
            "ppt": {"style_key": {...}},
            "word": {"style_key": {...}}
        }
    若文件缺失，返回空配置并记录日志。
    """

    if not TEMPLATES_CONFIG_PATH.exists():
        logger.warning("未找到模板配置文件 %s。", TEMPLATES_CONFIG_PATH)
        return {"ppt": {}, "word": {}}

    try:
        with TEMPLATES_CONFIG_PATH.open("r", encoding="utf-8") as fp:
            data = json.load(fp)
        if not isinstance(data, dict):
            logger.warning("模板配置文件结构非法，已忽略。")
            return {"ppt": {}, "word": {}}
        return {
            "ppt": data.get("ppt", {}),
            "word": data.get("word", {}),
        }
    except Exception as exc:
        logger.error("加载模板配置失败: %s", exc, exc_info=True)
        return {"ppt": {}, "word": {}}


def _extract_pdf_multimodal(pdf_path: Path, temp_dir: Path) -> Dict[str, any]:
    """
    从 PDF 提取多模态内容：文本、表格、图片。
    
    Returns:
        {
            "text": "全文文本",
            "tables": [{"page": 1, "data": [[row1], [row2], ...]}, ...],
            "images": [{"page": 1, "path": Path(...), "bbox": (x0,y0,x1,y1)}, ...]
        }
    """
    
    def _merge_single_char_rows(table: List[List[str]]) -> List[List[str]]:
        """
        合并OCR表格中的单字行(可能是被错误分离的姓名)
        
        策略:
        1. 查找只有1-2个字符的行
        2. 向前查找最近的版本号行(可能不是紧邻的前一行)
        3. 将单字行合并到版本号行的第3列(团队列)
        
        Args:
            table: OCR提取的表格
            
        Returns:
            合并后的表格
        """
        if not table or len(table) < 2:
            return table
        
        import re
        version_pattern = re.compile(r'^\d+\.\d+$')  # 匹配版本号: 0.1, 1.0等
        
        # 第一步: 找出所有单字行的索引
        single_char_rows = []
        for i, row in enumerate(table):
            is_single_char = (
                len(row) > 0 and
                len(row[0].strip()) <= 2 and
                len(row[0].strip()) > 0 and  # 不是空行
                all(not cell.strip() for cell in row[1:])
            )
            if is_single_char:
                single_char_rows.append(i)
                logger.debug("  发现单字行[%d]: '%s'", i, row[0])
        
        if not single_char_rows:
            logger.debug("  没有发现单字行,无需合并")
            return table
        
        # 第二步: 对每个单字行,向前查找最近的版本号行
        merge_map = {}  # {单字行索引: 目标版本号行索引}
        for single_idx in single_char_rows:
            # 向前查找版本号行(最多向前查找5行)
            for j in range(single_idx - 1, max(-1, single_idx - 6), -1):
                if j >= 0 and len(table[j]) > 0:
                    if version_pattern.match(table[j][0].strip()):
                        merge_map[single_idx] = j
                        logger.debug("  单字行[%d]'%s' → 合并到版本行[%d]'%s'", 
                                   single_idx, table[single_idx][0], j, table[j][0])
                        break
        
        # 第三步: 执行合并
        merged_table = []
        skip_indices = set()
        
        for i, row in enumerate(table):
            if i in skip_indices:
                continue
            
            # 如果当前行是版本号行,检查是否有单字行需要合并到它
            current_row = list(row)  # 复制一份
            for single_idx, target_idx in merge_map.items():
                if target_idx == i:
                    # 合并单字到第3列
                    single_char = table[single_idx][0].strip()
                    if len(current_row) > 2:
                        if current_row[2].strip():
                            current_row[2] = current_row[2] + single_char
                        else:
                            current_row[2] = single_char
                    else:
                        # 补齐列数
                        while len(current_row) < 3:
                            current_row.append('')
                        current_row[2] = single_char
                    
                    skip_indices.add(single_idx)
                    logger.info("  已合并单字'%s'到版本%s第3列: '%s'", 
                               single_char, current_row[0], current_row[2])
            
            merged_table.append(current_row)
        
        logger.info("  单字行合并完成: 原%d行 → %d行", len(table), len(merged_table))
        return merged_table
    
    result = {"text": "", "tables": [], "images": [], "pages": []}
    
    if not PYMUPDF_AVAILABLE:
        logger.warning("PyMuPDF 未安装，无法提取图片。")
        return result
    
    try:
        import pdfplumber
        
        # 使用 pdfplumber 提取文本和表格
        with pdfplumber.open(str(pdf_path)) as pdf:
            text_parts = []
            for page_num, page in enumerate(pdf.pages, start=1):
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
                
                # 保存每页的文本（用于生成单独的幻灯片）
                result["pages"].append({
                    "page": page_num,
                    "text": page_text if page_text else ""
                })
                
                # 提取表格(使用更精确的设置)
                # 优化表格提取参数,提高复杂表格的识别准确率
                tables = page.extract_tables(table_settings={
                    "vertical_strategy": "lines",
                    "horizontal_strategy": "lines",
                    "snap_tolerance": 5,  # 增加容差,更宽松地识别表格线
                    "join_tolerance": 5,
                    "edge_min_length": 5,
                    "min_words_vertical": 2,  # 降低阈值,避免漏掉短文本
                    "min_words_horizontal": 1,
                    "text_tolerance": 5,  # 增加文本容差,允许更大的Y坐标偏差
                    "text_x_tolerance": 5,
                    "text_y_tolerance": 5,
                    "intersection_tolerance": 5,
                })
                
                # 如果lines策略失败,尝试text策略
                if not tables or len(tables) == 0:
                    tables = page.extract_tables(table_settings={
                        "vertical_strategy": "text",
                        "horizontal_strategy": "text",
                    })
                
                for table in tables:
                    if table and len(table) > 1:  # 至少有标题行和一行数据
                        # 检查第一行是否为空(可能是纯背景色的表头)
                        first_row = table[0]
                        first_row_empty = all(not cell or not str(cell).strip() for cell in first_row)
                        
                        if first_row_empty and len(table) > 2:
                            # 第一行为空,使用第二行作为表头
                            logger.debug("检测到空表头行,使用第二行作为表头,页面=%d", page_num)
                            table = table[1:]  # 跳过空的第一行
                        
                        # 处理单元格中的换行符
                        # 第一步:检查第一行是否包含表头和数据(用换行符分隔)
                        first_row = table[0]
                        has_newline_in_first_row = any('\n' in str(cell) for cell in first_row if cell)
                        
                        # 添加详细日志查看原始数据
                        logger.debug("页面%d原始表格所有行:", page_num)
                        for row_idx, row in enumerate(table):
                            logger.debug("  第%d行: %s", row_idx, [str(cell)[:30] for cell in row])
                            
                            # 智能修复: 如果某行只有"LLM",尝试从下一行获取完整内容
                            # 这是因为PDF提取时可能将内容分散到多行
                            for row_idx in range(len(table) - 1):
                                row = table[row_idx]
                                next_row = table[row_idx + 1]
                                # 检查当前行的第二列是否只有"LLM"(或其他短关键词)
                                if len(row) > 1 and str(row[1]).strip() == 'LLM':
                                    # 检查下一行的第二列是否有内容
                                    if len(next_row) > 1 and next_row[1] and str(next_row[1]).strip():
                                        next_content = str(next_row[1]).strip()
                                        # 如果下一行内容以"增加"或其他动词开头,可能是被分离的内容
                                        if next_content.startswith('增加') or next_content.startswith('修改') or next_content.startswith('优化'):
                                            # 合并内容: "增加 模型..." → "增加 LLM 模型..."
                                            if '\n' in next_content:
                                                # 移除末尾的"\nLLM"
                                                next_content = next_content.replace('\nLLM', '').strip()
                                            # 将"LLM"插入到内容中
                                            parts = next_content.split(' ', 1)
                                            if len(parts) == 2:
                                                row[1] = f"{parts[0]} LLM {parts[1]}"
                                            else:
                                                row[1] = f"{next_content} LLM"
                                            logger.debug("  修复第%d行: 从下一行合并内容 → '%s'", row_idx, str(row[1])[:50])
                                            
                                            # TODO: 下一行的内容已经被合并到当前行,但下一行本应有自己的内容
                                            # 这是pdfplumber的提取bug,无法从错误数据中恢复
                                            # 可能的解决方案: 使用OCR重新提取表格内容
                                            # 当前workaround: 保持下一行的原始内容(虽然可能不正确)
                        
                        if has_newline_in_first_row:
                            # 第一行包含换行符,可能是表头和数据合并
                            # 尝试分离表头和数据
                            header_row = []
                            data_row = []
                            for col_idx, cell in enumerate(first_row):
                                if cell and '\n' in str(cell):
                                    parts = str(cell).split('\n', 1)
                                    header_part = parts[0].strip()
                                    data_part = parts[1].strip() if len(parts) > 1 else ''
                                    
                                    # 智能重组: 如果数据部分以短关键词开头,将其移到后面
                                    # 例如: "LLM\n首次发布..." → "首次发布: LLM ..."
                                    if data_part and '\n' in data_part:
                                        data_lines = data_part.split('\n')
                                        # 如果第一行很短(<=10字符)且第二行以中文开头,可能是关键词错位
                                        if len(data_lines) >= 2 and len(data_lines[0]) <= 10:
                                            first_line = data_lines[0].strip()
                                            rest_lines = '\n'.join(data_lines[1:]).strip()
                                            # 检查第二行是否以中文或"首次"等开头
                                            if rest_lines and (rest_lines[0] >= '\u4e00' or rest_lines.startswith('首次') or rest_lines.startswith('增加')):
                                                # 重组: 将短关键词移到后面
                                                # "首次发布：..." → "首次发布： LLM ..."
                                                if '：' in rest_lines or ':' in rest_lines:
                                                    # 在冒号后插入关键词
                                                    for sep in ['：', ':']:
                                                        if sep in rest_lines:
                                                            prefix, suffix = rest_lines.split(sep, 1)
                                                            data_part = f"{prefix}{sep} {first_line} {suffix}".strip()
                                                            logger.debug("  重组列%d数据: '%s' + '%s' → '%s'", 
                                                                       col_idx, first_line[:20], rest_lines[:30], data_part[:50])
                                                            break
                                    
                                    header_row.append(header_part)
                                    data_row.append(data_part)
                                    logger.debug("  页面%d分离列%d: 表头='%s', 数据='%s'", 
                                               page_num, col_idx, header_part[:50], data_part[:50])
                                else:
                                    header_row.append(str(cell).strip() if cell else '')
                                    data_row.append('')
                            
                            # 检查是否成功分离(数据行不全为空)
                            if any(data_row):
                                logger.debug("分离表头和数据行,页面=%d", page_num)
                                # 替换第一行为表头,插入数据行
                                table[0] = header_row
                                table.insert(1, data_row)
                        
                        # 第二步:清理所有行的换行符(包括后续数据行,也需要智能重组)
                        processed_table = []
                        prev_content = {}  # 记录前一行的内容,用于检测重复
                        for row_idx, row in enumerate(table):
                            cleaned_row = []
                            for col_idx, cell in enumerate(row):
                                if cell and '\n' in str(cell):
                                    cell_str = str(cell)
                                    
                                    # 对所有单元格应用智能重组逻辑(不只是第一行)
                                    if '\n' in cell_str:
                                        lines = cell_str.split('\n')
                                        # 如果最后一行是"LLM"(或其他短关键词),将其移到前面
                                        if len(lines) >= 2 and lines[-1].strip() == 'LLM':
                                            keyword = lines[-1].strip()
                                            content_lines = lines[:-1]  # 移除最后一行
                                            content = ' '.join(line.strip() for line in content_lines if line.strip())
                                            # 将关键词插入到内容中
                                            # 例如: "增加 模型..." + "LLM" → "增加 LLM 模型..."
                                            parts = content.split(' ', 1)
                                            if len(parts) == 2:
                                                cell_str = f"{parts[0]} {keyword} {parts[1]}".strip()
                                            else:
                                                cell_str = f"{content} {keyword}".strip()
                                            if row_idx > 0:
                                                logger.debug("  页面%d重组第%d行列%d: '%s' + '%s' → '%s'", 
                                                           page_num, row_idx, col_idx, content[:20], keyword, cell_str[:50])
                                        # 如果第一行很短(<=10字符)且第二行以中文开头,可能是关键词错位
                                        elif len(lines) >= 2 and len(lines[0]) <= 10:
                                            first_line = lines[0].strip()
                                            rest_lines = '\n'.join(lines[1:]).strip()
                                            # 检查第二行是否以中文或"首次"等开头
                                            if rest_lines and (rest_lines[0] >= '\u4e00' or rest_lines.startswith('首次') or rest_lines.startswith('增加')):
                                                # 重组: 将短关键词移到后面
                                                if '：' in rest_lines or ':' in rest_lines:
                                                    # 在冒号后插入关键词
                                                    for sep in ['：', ':']:
                                                        if sep in rest_lines:
                                                            prefix, suffix = rest_lines.split(sep, 1)
                                                            cell_str = f"{prefix}{sep} {first_line} {suffix}".strip()
                                                            if row_idx > 0:
                                                                logger.debug("  页面%d重组第%d行列%d: '%s' + '%s' → '%s'", 
                                                                           page_num, row_idx, col_idx, first_line[:20], rest_lines[:30], cell_str[:50])
                                                            break
                                    
                                    # 清理换行符
                                    cleaned_cell = ' '.join(cell_str.split('\n'))
                                    cleaned_row.append(cleaned_cell.strip())
                                    
                                    # 检测重复内容: 如果当前单元格内容与前一行相同,可能是PDF提取错误
                                    # 这种情况下,我们无法从错误数据中恢复正确数据,只能标记为可疑
                                    if col_idx == 1 and row_idx > 0:
                                        if col_idx in prev_content and cleaned_cell.strip() == prev_content[col_idx]:
                                            # 内容重复!这是PDF提取的bug
                                            # 无法修复,因为真实内容已经丢失
                                            logger.warning("  页面%d检测到第%d行列%d内容与前一行重复: '%s' (PDF提取错误,无法修复)", 
                                                         page_num, row_idx, col_idx, cleaned_cell[:50])
                                    
                                    # 记录当前内容
                                    prev_content[col_idx] = cleaned_cell.strip()
                                else:
                                    cleaned_row.append(str(cell).strip() if cell else '')
                                    # 记录当前内容
                                    prev_content[col_idx] = str(cell).strip() if cell else ''
                            
                            processed_table.append(cleaned_row)
                        
                        # OCR验证: 如果检测到重复内容,只对有问题的行使用OCR
                        # 不再整个表格重新提取,而是保留pdfplumber的表结构
                        ocr_extractor = get_ocr_extractor()
                        if ocr_extractor:
                            # 检测有问题的行(忽略空行)
                            duplicate_rows = []
                            for i in range(1, len(processed_table) - 1):
                                if len(processed_table[i]) > 0 and len(processed_table[i+1]) > 0:
                                    # 检查所有列是否有重复内容(而不是只检查第2列)
                                    has_duplicate = False
                                    for col_idx in range(min(len(processed_table[i]), len(processed_table[i+1]))):
                                        content1 = processed_table[i][col_idx].strip()
                                        content2 = processed_table[i+1][col_idx].strip()
                                        # 只有当两列都有实际内容,且内容相同时,才认为是重复
                                        if content1 and content2 and content1 == content2:
                                            has_duplicate = True
                                            logger.debug("  第%d行第%d列与前一行重复: '%s'", i + 1, col_idx, content1[:30])
                                            break
                                    
                                    if has_duplicate:
                                        # 标记第i+1行为有问题的行(因为它是重复的)
                                        duplicate_rows.append(i + 1)
                                        logger.warning("页面%d第%d行内容与前一行重复,标记为需要OCR修复", page_num, i + 1)
                            
                            # 如果有重复行,使用OCR提取整个表格,但只替换有问题的行
                            if duplicate_rows:
                                try:
                                    ocr_table = ocr_extractor.extract_table_from_page(pdf_path, page_num)
                                    if ocr_table and len(ocr_table) > 0:
                                        logger.info("OCR提取成功,将修复%d个有问题的行", len(duplicate_rows))
                                        # 调试: 输出pdfplumber和OCR提取的完整表格对比
                                        logger.debug("pdfplumber表格(%d行):", len(processed_table))
                                        for idx, row in enumerate(processed_table):
                                            logger.debug("  PDF第%d行: %s", idx, row[:3] if len(row) > 3 else row)
                                        logger.debug("OCR提取的完整表格(%d行):", len(ocr_table))
                                        for idx, row in enumerate(ocr_table):
                                            logger.debug("  OCR第%d行: %s", idx, row[:3] if len(row) > 3 else row)
                                        
                                        # 预处理: 合并OCR表格中的单字行(可能是被错误分离的姓名)
                                        ocr_table = _merge_single_char_rows(ocr_table)
                                        
                                        # 只替换有问题的行,保留pdfplumber的表结构
                                        # 通过第一列的值(版本号)来匹配OCR行,而不是用行索引
                                        for row_idx in duplicate_rows:
                                            if row_idx >= len(processed_table):
                                                continue
                                            
                                            target_row = processed_table[row_idx]
                                            target_cols = len(target_row)
                                            
                                            # 通过第一列的值查找OCR表格中的对应行
                                            first_col_value = target_row[0].strip() if len(target_row) > 0 else ''
                                            ocr_row = None
                                            ocr_row_idx = -1
                                            
                                            for idx, row in enumerate(ocr_table):
                                                if len(row) > 0 and row[0].strip() == first_col_value:
                                                    ocr_row = row
                                                    ocr_row_idx = idx
                                                    break
                                            
                                            if not ocr_row:
                                                logger.warning("  第%d行未在OCR表格中找到匹配(第一列='%s'),跳过", row_idx, first_col_value)
                                                continue
                                            
                                            # 调试: 输出OCR提取的原始行数据
                                            logger.debug("  通过第一列'%s'匹配到OCR第%d行", first_col_value, ocr_row_idx)
                                            logger.debug("  OCR原始第%d行(%d列): %s", ocr_row_idx, len(ocr_row), ocr_row)
                                            logger.debug("  目标第%d行(%d列): %s", row_idx, target_cols, target_row)
                                            
                                            # 智能修复: 如果OCR团队列只有单字,尝试从pdfplumber其他行中查找完整姓名
                                            if len(ocr_row) > 2 and len(ocr_row[2].strip()) <= 2:
                                                single_char = ocr_row[2].strip()
                                                # 在pdfplumber表格中查找包含该单字的完整姓名
                                                for pdf_row in processed_table:
                                                    if len(pdf_row) > 2:
                                                        team_col = pdf_row[2].strip()
                                                        # 如果找到包含该单字的姓名(如"李赟")
                                                        if single_char and single_char in team_col:
                                                            # 提取完整姓名(假设姓名是2个字)
                                                            import re
                                                            # 查找单字后面紧跟的一个字
                                                            pattern = single_char + r'[\u4e00-\u9fff]'
                                                            match = re.search(pattern, team_col)
                                                            if match:
                                                                full_name = match.group()
                                                                ocr_row[2] = full_name
                                                                logger.info("  智能修复团队列: '%s' → '%s'", single_char, full_name)
                                                                break
                                            
                                            # 如果列数一致,直接替换
                                            if len(ocr_row) == target_cols:
                                                processed_table[row_idx] = ocr_row
                                                logger.debug("  替换第%d行: %s", row_idx, ocr_row[:3])
                                            # 如果OCR列数更多,智能对齐列
                                            elif len(ocr_row) > target_cols:
                                                # 尝试找到最佳的列对齐方式
                                                aligned_row = _align_ocr_row_to_target(
                                                    ocr_row, 
                                                    target_row,
                                                    processed_table[0] if len(processed_table) > 0 else None
                                                )
                                                processed_table[row_idx] = aligned_row
                                                logger.info("  第%d行列数对齐: OCR %d列 → %d列, 对齐后: %s", 
                                                           row_idx, len(ocr_row), target_cols, aligned_row[:3])
                                            # 如果OCR列数更少,补齐空字符串
                                            else:
                                                aligned_row = ocr_row + [''] * (target_cols - len(ocr_row))
                                                processed_table[row_idx] = aligned_row
                                                logger.info("  第%d行列数对齐: OCR %d列 → %d列, 补齐后: %s", 
                                                           row_idx, len(ocr_row), target_cols, aligned_row[:3])
                                    else:
                                        logger.warning("OCR提取失败,保留pdfplumber结果")
                                except Exception as e:
                                    logger.error("OCR提取异常: %s", e)
                        
                        # Fallback: 如果OCR失败,尝试从校核列或前一行推断团队列的缺失数据
                        import re
                        for row_idx in range(1, len(processed_table)):  # 跳过表头
                            row = processed_table[row_idx]
                            if len(row) > 3:  # 至少有4列(版本、内容、团队、校核)
                                team_col = row[2].strip() if len(row) > 2 else ''
                                review_col = row[3].strip() if len(row) > 3 else ''
                                
                                # 策略1: 如果团队列为空,但校核列有内容
                                if not team_col and review_col:
                                    # 提取校核列的第一个姓名(假设姓名是2-3个中文字符)
                                    name_pattern = r'[\u4e00-\u9fff]{2,3}'
                                    match = re.search(name_pattern, review_col)
                                    if match:
                                        first_name = match.group()
                                        row[2] = first_name
                                        logger.info("  Fallback修复(策略1): 从校核列'%s'推断团队列 → '%s' (第%d行)", 
                                                   review_col[:20], first_name, row_idx)
                                
                                # 策略2: 如果团队列和校核列都为空,从前一行的团队列推断
                                if not team_col and not review_col and row_idx > 1:
                                    prev_row = processed_table[row_idx - 1]
                                    if len(prev_row) > 2:
                                        prev_team = prev_row[2].strip()
                                        if prev_team:
                                            # 提取前一行团队列的第一个姓名
                                            name_pattern = r'[\u4e00-\u9fff]{2,3}'
                                            match = re.search(name_pattern, prev_team)
                                            if match:
                                                first_name = match.group()
                                                row[2] = first_name
                                                logger.info("  Fallback修复(策略2): 从前一行团队列'%s'推断 → '%s' (第%d行)", 
                                                           prev_team[:20], first_name, row_idx)
                        
                        # 最后再次规范化表格结构(确保所有行列数一致)
                        processed_table = _normalize_table_structure(processed_table)
                        
                        result["tables"].append({
                            "page": page_num,
                            "data": processed_table
                        })
                        logger.debug("提取表格: 页面=%d, 行数=%d, 列数=%d, 首行=%s", 
                                   page_num, len(processed_table), len(processed_table[0]) if processed_table else 0, 
                                   processed_table[0] if processed_table else [])
            
            result["text"] = "\n".join(text_parts)
        
        # 使用 PyMuPDF 提取图片（过滤小图）
        # 按页面去重:同一页面内的相同xref只提取一次,但不同页面可以有相同图片
        pdf_document = fitz.open(str(pdf_path))
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            image_list = page.get_images(full=True)
            
            page_extracted_xrefs = set()  # 当前页面已提取的xref
            img_count = 0  # 当前页面有效图片计数
            for img_index, img in enumerate(image_list):
                xref = img[0]
                
                # 跳过当前页面已提取的图片
                if xref in page_extracted_xrefs:
                    logger.debug("跳过页面内重复图片: xref=%d (页面 %d)", xref, page_num + 1)
                    continue
                
                try:
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # 检查图片尺寸
                    img_pil = Image.open(io.BytesIO(image_bytes))
                    width, height = img_pil.size
                    
                    # 过滤条件：宽度或高度小于 200px 的图片（通常是装饰性元素）
                    if width < 200 or height < 200:
                        logger.debug("跳过小图片: %dx%d (页面 %d)", width, height, page_num + 1)
                        continue
                    
                    # 过滤条件：宽高比极端的图片（如横幅、分隔线）
                    aspect_ratio = max(width, height) / min(width, height)
                    if aspect_ratio > 10:
                        logger.debug("跳过极端宽高比图片: %.1f (页面 %d)", aspect_ratio, page_num + 1)
                        continue
                    
                    # 保存图片到临时目录
                    img_count += 1
                    img_filename = f"page{page_num + 1}_img{img_count}.{image_ext}"
                    img_path = temp_dir / img_filename
                    img_path.write_bytes(image_bytes)
                    
                    result["images"].append({
                        "page": page_num + 1,
                        "path": img_path,
                        "index": img_count,
                        "width": width,
                        "height": height,
                        "bytes": image_bytes,  # 用于去重
                        "xref": xref  # 记录xref用于调试
                    })
                    
                    # 标记为当前页面已提取
                    page_extracted_xrefs.add(xref)
                    
                    logger.debug("提取有效图片: %s (%dx%d, 页面 %d, xref=%d)", 
                                img_filename, width, height, page_num + 1, xref)
                    
                except Exception as e:
                    logger.warning("图片提取失败 xref=%d (页面 %d): %s", xref, page_num + 1, e)
                    continue
        
        pdf_document.close()
        logger.info("PDF 多模态提取完成: 文本 %d 字符, 表格 %d 个, 图片 %d 张",
                    len(result["text"]), len(result["tables"]), len(result["images"]))
        return result
    
    except Exception as exc:
        logger.error("PDF 多模态提取失败: %s", exc, exc_info=True)
        return result


def _read_text_from_file(file_path: Path) -> str:
    if not file_path or not file_path.exists():
        return ""
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(str(file_path)) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                    return "\n".join(text_parts)
            except ImportError:
                logger.warning("pdfplumber 未安装，无法提取 PDF 文本。")
                return ""
        elif suffix == ".docx":
            doc = WordDocument(str(file_path))
            return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        elif suffix in {".txt", ".md", ".csv"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")
        else:
            return file_path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        logger.error("读取本地文件失败: %s", exc, exc_info=True)
        return ""


def _read_text_from_url(url: str) -> str:
    if not url:
        return ""
    try:
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        text = resp.text
        # 简单去除 HTML 标签
        text = re.sub(r"<script[\s\S]*?</script>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<style[\s\S]*?</style>", "", text, flags=re.IGNORECASE)
        text = re.sub(r"<[^>]+>", "\n", text)
        return re.sub(r"\n{2,}", "\n", text)
    except Exception as exc:
        logger.error("下载 URL 文本失败，URL=%s，错误=%s", url, exc, exc_info=True)
        return ""


def _collect_source_text(source_file_path: Optional[Path], source_url: str) -> str:
    collected = []
    if source_file_path:
        collected.append(_read_text_from_file(source_file_path))
    if source_url:
        collected.append(_read_text_from_url(source_url))
    text = "\n".join(part.strip() for part in collected if part and part.strip())
    return text.strip()


def _ensure_text_chunks(text: str) -> List[str]:
    if not text:
        return []
    paragraphs = [seg.strip() for seg in text.split("\n") if seg.strip()]
    chunks: List[str] = []
    current = []
    length = 0
    for para in paragraphs:
        para_len = len(para)
        if length + para_len > 500 and current:
            chunks.append("\n".join(current))
            current = [para]
            length = para_len
        else:
            current.append(para)
            length += para_len
    if current:
        chunks.append("\n".join(current))
    return chunks if chunks else [text]


def _analyze_content_with_ai(text: str, request_id: str) -> Dict[str, any]:
    """
    使用 DashScope AI 分析文档内容，提取结构化信息。
    
    Args:
        text: 原始文本内容
        request_id: 请求 ID
    
    Returns:
        结构化数据：
        {
            "title": "文档标题",
            "subtitle": "副标题",
            "sections": [
                {
                    "title": "章节标题",
                    "points": ["要点1", "要点2", ...]
                },
                ...
            ]
        }
    """
    logger.info("开始 AI 内容分析，RequestID=%s", request_id)
    
    prompt = f"""请分析以下文档内容，提取结构化信息并以 JSON 格式返回。

要求：
1. 识别文档的主标题和副标题
2. 将内容分为多个章节（每个章节包含标题和3-5个要点）
3. 每个要点应简洁明了，适合在 PPT 中展示
4. 返回格式必须是有效的 JSON

文档内容：
{text[:4000]}

请返回 JSON 格式：
{{
  "title": "主标题",
  "subtitle": "副标题",
  "sections": [
    {{
      "title": "章节1标题",
      "points": ["要点1", "要点2", "要点3"]
    }}
  ]
}}"""
    
    try:
        response = Generation.call(
            model="qwen-plus",
            prompt=prompt,
            result_format="message"
        )
        
        if response.status_code == 200:
            ai_output = response.output.choices[0].message.content
            logger.debug("AI 返回内容: %s", ai_output)
            
            # 提取 JSON 部分（去除可能的 Markdown 代码块标记）
            json_match = re.search(r'```json\s*({.*?})\s*```', ai_output, re.DOTALL)
            if json_match:
                json_str = json_match.group(1)
            else:
                # 尝试直接解析
                json_str = ai_output.strip()
            
            result = json.loads(json_str)
            logger.info("AI 内容分析成功，提取 %d 个章节。", len(result.get("sections", [])))
            return result
        else:
            logger.error("AI 调用失败: %s", response.message)
            return _fallback_structure(text)
    
    except Exception as exc:
        logger.error("AI 内容分析失败: %s", exc, exc_info=True)
        return _fallback_structure(text)


def _fallback_structure(text: str) -> Dict[str, any]:
    """
    当 AI 分析失败时的降级方案：简单分段。
    """
    logger.warning("使用降级方案进行内容结构化。")
    
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    
    # 简单启发式：第一行作为标题
    title = lines[0] if lines else "文档内容"
    subtitle = lines[1] if len(lines) > 1 else "自动生成"
    
    # 将剩余内容分为若干段
    sections = []
    chunk_size = 5
    for i in range(2, len(lines), chunk_size):
        chunk = lines[i:i+chunk_size]
        sections.append({
            "title": f"内容概览 {len(sections) + 1}",
            "points": chunk
        })
    
    return {
        "title": title,
        "subtitle": subtitle,
        "sections": sections if sections else [{"title": "内容", "points": lines[2:]}]
    }


def generate_word_document(
    *,
    request_id: str,
    username: str,
    upload_dir: Path,
    converted_dir: Path,
    temp_dir: Path,
    source_file_path: Optional[Path],
    source_url: str,
    template_config: Optional[dict] = None,
) -> Tuple[str, str]:
    """生成 Word 文档，返回 (文件名, 提示信息)。"""

    # 1. 检测是否为 PDF 并提取多模态内容
    is_pdf = source_file_path and source_file_path.suffix.lower() == ".pdf"
    multimodal_data = None
    
    if is_pdf:
        logger.info("检测到 PDF 文件，使用AI智能分析模式...")
        multimodal_data = _extract_pdf_multimodal(source_file_path, temp_dir)
        text = multimodal_data["text"]
    else:
        logger.info("提取文本内容...")
        text = _collect_source_text(source_file_path, source_url)
        
        if not text:
            raise ValueError("未能从输入源提取到有效文本。")
    
    # 2. 根据文件类型选择生成模式
    if is_pdf and multimodal_data:
        # PDF 文件：使用AI智能分析模式
        logger.info("使用AI分析文档结构...")
        
        # 初始化AI分析器
        ai_analyzer = AIDocumentAnalyzer(model="qwen-max")
        
        # 分析文档整体结构
        document_structure = ai_analyzer.analyze_document_structure(multimodal_data, request_id)
        
        # 分析每一页的内容
        page_analyses = []
        for page_data in multimodal_data.get("pages", []):
            page_num = page_data["page"]
            page_text = page_data.get("text", "")
            page_tables = [t for t in multimodal_data.get("tables", []) if t["page"] == page_num]
            page_images = [i for i in multimodal_data.get("images", []) if i["page"] == page_num]
            
            page_analysis = ai_analyzer.analyze_page_content(
                page_num, page_text, page_tables, page_images, request_id
            )
            page_analyses.append(page_analysis)
        
        # 使用智能Word生成器
        logger.info("使用AI智能生成器生成Word...")
        from .smart_word_generator import SmartWordGenerator
        
        smart_generator = SmartWordGenerator()
        doc = smart_generator.generate_word(
            document_structure,
            page_analyses,
            multimodal_data,
            request_id
        )
    else:
        # 非 PDF 文件：传统文本模式
        logger.info("使用传统文本模式生成Word...")
        
        if template_config and template_config.get("template_path"):
            template_path = Path(template_config["template_path"])
            if template_path.exists():
                doc = WordDocument(str(template_path))
            else:
                logger.warning("模板 %s 不存在，使用空白文档。", template_path)
                doc = WordDocument()
        else:
            doc = WordDocument()

        style_name = template_config.get("paragraph_style") if template_config else None
        for chunk in _ensure_text_chunks(text):
            paragraph = doc.add_paragraph(chunk)
            if style_name and style_name in doc.styles:
                paragraph.style = style_name
            else:
                run = paragraph.runs[0]
                run.font.size = Pt(12)

    # 3. 保存文档
    output_filename = f"{request_id}_document.docx"
    output_path = converted_dir / output_filename
    doc.save(output_path)

    logger.info(
        "Word 文档生成完成: %s (user=%s, request=%s)",
        output_path,
        username,
        request_id,
    )
    return output_filename, "Word 文档生成成功。"




def generate_ppt_document(
    *,
    request_id: str,
    username: str,
    upload_dir: Path,
    converted_dir: Path,
    temp_dir: Path,
    source_file_path: Optional[Path],
    source_url: str,
    template_config: dict,
) -> Tuple[str, str]:
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx 未安装，无法生成 PPT。")

    # 1. 检测是否为 PDF 并提取多模态内容
    is_pdf = source_file_path and source_file_path.suffix.lower() == ".pdf"
    multimodal_data = None
    
    if is_pdf:
        logger.info("检测到 PDF 文件，使用AI智能分析模式...")
        multimodal_data = _extract_pdf_multimodal(source_file_path, temp_dir)
        text = multimodal_data["text"]
    else:
        logger.info("提取文本内容...")
        text = _collect_source_text(source_file_path, source_url)
        
        if not text:
            raise ValueError("未能从输入源提取到有效文本。")
    
    # 2. 使用AI分析文档结构(所有文件类型统一使用AI)
    logger.info("使用AI分析文档结构...")
    
    # 初始化AI分析器
    ai_analyzer = AIDocumentAnalyzer(model="qwen-max")
    
    # 如果是PDF,使用AI分析整体结构
    if is_pdf and multimodal_data:
        document_structure = ai_analyzer.analyze_document_structure(multimodal_data, request_id)
        
        # 分析每一页的内容
        page_analyses = []
        for page_data in multimodal_data.get("pages", []):
            page_num = page_data["page"]
            page_text = page_data.get("text", "")
            page_tables = [t for t in multimodal_data.get("tables", []) if t["page"] == page_num]
            page_images = [i for i in multimodal_data.get("images", []) if i["page"] == page_num]
            
            page_analysis = ai_analyzer.analyze_page_content(
                page_num, page_text, page_tables, page_images, request_id
            )
            page_analyses.append(page_analysis)
    else:
        # 非PDF文件,使用传统AI分析
        structure = _analyze_content_with_ai(text, request_id)
        document_structure = None
        page_analyses = None
    
    # 3. 加载预定义模板
    template_path = template_config.get("template_path")
    if template_path:
        tpl_path = Path(settings.BASE_DIR).parent / template_path
        if not tpl_path.exists():
            logger.warning("PPT 模板 %s 不存在，使用空白模板。", tpl_path)
            tpl_path = None
    else:
        tpl_path = None
    
    # 4. 根据文件类型选择生成模式
    if is_pdf and multimodal_data and document_structure and page_analyses:
        # PDF 文件：使用AI智能生成器
        logger.info("使用AI智能生成器生成PPT...")
        
        smart_generator = SmartPPTGenerator()
        presentation = smart_generator.generate_ppt(
            tpl_path if tpl_path else Path(settings.BASE_DIR).parent / "config" / "templates" / "academic_template.pptx",
            document_structure,
            page_analyses,
            multimodal_data,
            request_id
        )
    else:
        # 非 PDF 文件：传统AI分析模式
        logger.info("使用传统AI分析模式生成PPT...")
        
        # 加载模板
        if tpl_path and tpl_path.exists():
            presentation = Presentation(str(tpl_path))
            # 删除模板示例页
            if len(presentation.slides) > 1:
                xml_slides = presentation.slides._sldIdLst
                for idx in reversed(range(1, len(xml_slides))):
                    rId = xml_slides[idx].rId
                    presentation.part.drop_rel(rId)
                    del xml_slides[idx]
        else:
            presentation = Presentation()
        
        # 创建标题页
        title_text = structure.get("title", template_config.get("title", "文档演示"))
        subtitle_text = structure.get("subtitle", template_config.get("subtitle", "AI 智能生成"))
        
        if len(presentation.slides) > 0:
            title_slide = presentation.slides[0]
            if title_slide.shapes.title:
                title_slide.shapes.title.text = title_text
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = subtitle_text
        else:
            title_layout = presentation.slide_layouts[0]
            title_slide = presentation.slides.add_slide(title_layout)
            title_slide.shapes.title.text = title_text
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = subtitle_text
        
        logger.info("已创建标题页: %s", title_text)
        
        # 创建内容页
        sections = structure.get("sections", [])
        bullet_layout_index = template_config.get("bullet_layout_index", 1)
        try:
            body_layout = presentation.slide_layouts[bullet_layout_index]
        except IndexError:
            body_layout = presentation.slide_layouts[1]

        for idx, section in enumerate(sections, start=1):
            section_title = section.get("title", f"章节 {idx}")
            points = section.get("points", [])
            
            if not points:
                continue
            
            slide = presentation.slides.add_slide(body_layout)
            
            if slide.shapes.title:
                slide.shapes.title.text = section_title
            
            body_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    body_shape = shape
                    break
            
            if not body_shape:
                logger.warning("幻灯片 %d 未找到内容占位符，已跳过。", idx)
                continue
            
            text_frame = body_shape.text_frame
            text_frame.clear()
            
            for point_idx, point in enumerate(points):
                if not point.strip():
                    continue
                
                if point_idx == 0:
                    text_frame.text = point
                    if text_frame.paragraphs:
                        text_frame.paragraphs[0].font.size = PptPt(18)
                else:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = PptPt(16)
            
            logger.debug("已创建章节页: %s (%d 个要点)", section_title, len(points))

    # 6. 保存 PPT
    output_filename = f"{request_id}_slides.pptx"
    output_path = converted_dir / output_filename
    presentation.save(str(output_path))

    total_pages = len(presentation.slides)
    logger.info(
        "PPT 文档生成完成: %s，共 %d 页 (user=%s, request=%s)",
        output_path,
        total_pages,
        username,
        request_id,
    )
    return output_filename, f"PPT 文档生成成功，共 {total_pages} 页。"


def _generate_ppt_from_pdf_pages(
    presentation: Presentation,
    multimodal_data: Dict[str, any],
    template_config: dict
) -> int:
    """
    按 PDF 页面结构生成 PPT（页面保留模式）。
    每个 PDF 页面对应一个或多个 PPT 页面（表格+图片）。
    
    Returns:
        生成的页面数量
    """
    # 使用模板的内容布局(通常是索引1或5),而非空白布局(索引6)
    # 索引1通常是"标题和内容"布局,有模板样式
    try:
        content_layout = presentation.slide_layouts[1]  # 标题和内容布局
    except IndexError:
        content_layout = presentation.slide_layouts[0]  # 回退到标题布局
    pages_added = 0
    
    # 按页面分组
    pages_with_content = {}
    
    # 初始化所有页面（包含文本）
    for page_data in multimodal_data.get("pages", []):
        page_num = page_data["page"]
        pages_with_content[page_num] = {
            "text": page_data["text"],
            "tables": [],
            "images": []
        }
    
    # 收集表格(智能合并连续小表格页)
    # 如果连续两页都只有表格且第一页是小表格,则合并到第一页
    tables_by_page = {}
    for table_data in multimodal_data["tables"]:
        page_num = table_data["page"]
        if page_num not in tables_by_page:
            tables_by_page[page_num] = []
        tables_by_page[page_num].append(table_data)
        logger.debug("收集表格: 第 %d 页, %d 行 x %d 列", 
                    page_num, len(table_data["data"]), len(table_data["data"][0]) if table_data["data"] else 0)
    
    # 不进行表格合并,保持PDF原始页面结构
    # (之前的合并逻辑导致错误地将两个独立页面的表格合并)
    
    # 将修正后的表格分配到页面
    for page_num, tables in tables_by_page.items():
        if page_num not in pages_with_content:
            pages_with_content[page_num] = {"text": "", "tables": [], "images": []}
        pages_with_content[page_num]["tables"].extend(tables)
    
    # 收集图片（过滤背景装饰图 + 智能去重）
    import hashlib
    
    # 第一步：先收集所有图片（仅过滤背景图）
    temp_images_by_page = {}  # {page_num: [img_data, ...]}
    
    for img_data in multimodal_data["images"]:
        page_num = img_data["page"]
        width = img_data["width"]
        height = img_data["height"]
        
        # 通用背景图过滤规则：
        # 1. 常见演示文稿尺寸（16:9 或 4:3 比例，且尺寸较大）
        # 2. 宽高比接近 16:9 (1.77) 或 4:3 (1.33)，且面积超过 1.5M 像素
        aspect_ratio = width / height if height > 0 else 0
        area = width * height
        
        # 判断是否为全屏背景图（常见演示文稿尺寸）
        is_16_9 = 1.7 <= aspect_ratio <= 1.8  # 16:9 比例
        is_4_3 = 1.3 <= aspect_ratio <= 1.4   # 4:3 比例
        is_large = area > 1500000  # 面积超过 1.5M 像素
        
        if (is_16_9 or is_4_3) and is_large:
            logger.debug("跳过背景装饰图: %dx%d (页面 %d, 宽高比=%.2f, 面积=%d)", 
                        width, height, page_num, aspect_ratio, area)
            continue
        
        if page_num not in temp_images_by_page:
            temp_images_by_page[page_num] = []
        temp_images_by_page[page_num].append(img_data)
    
    # 第二步：全局去重 - 每张图片只保留在最后一个无表格的页面
    global_image_registry = {}  # {hash: [page_num1, page_num2, ...]}
    
    # 先收集每张图片出现在哪些页面
    for page_num in sorted(temp_images_by_page.keys()):
        page_images = temp_images_by_page[page_num]
        
        for img_data in page_images:
            img_bytes = img_data.get("bytes")
            if img_bytes:
                img_hash = hashlib.md5(img_bytes).hexdigest()
                if img_hash not in global_image_registry:
                    global_image_registry[img_hash] = []
                global_image_registry[img_hash].append({
                    "page": page_num,
                    "data": img_data
                })
    
    # 决定每张图片应该保留在哪个页面
    for img_hash, occurrences in global_image_registry.items():
        # 找到最后一个没有表格的页面
        target_page = None
        for occurrence in reversed(occurrences):
            page_num = occurrence["page"]
            page_has_table = page_num in [t["page"] for t in multimodal_data["tables"]]
            if not page_has_table:
                target_page = page_num
                break
        
        # 如果所有页面都有表格,则不保留此图片
        if target_page is None:
            logger.debug("图片 hash=%s 在所有页面都有表格,不保留", img_hash[:8])
            continue
        
        # 只在目标页面保留此图片
        for occurrence in occurrences:
            page_num = occurrence["page"]
            img_data = occurrence["data"]
            
            if page_num not in pages_with_content:
                pages_with_content[page_num] = {"text": "", "tables": [], "images": []}
            
            if page_num == target_page:
                pages_with_content[page_num]["images"].append(img_data)
                logger.debug("保留图片: %dx%d (hash=%s, 保留在第 %d 页)", 
                            img_data["width"], img_data["height"], img_hash[:8], page_num)
            else:
                logger.debug("跳过重复图片: %dx%d (hash=%s, 首次出现于第 %d 页，当前第 %d 页)", 
                            img_data["width"], img_data["height"], img_hash[:8], target_page, page_num)
    
    # 按页面顺序处理,支持智能合并小内容页面
    sorted_pages = sorted(pages_with_content.keys())
    i = 0
    while i < len(sorted_pages):
        page_num = sorted_pages[i]
        page_content = pages_with_content[page_num]
        page_text = page_content.get("text", "")
        tables = page_content["tables"]
        images = page_content["images"]
        
        # 跳过第1页(如果只有小表格,已添加到标题页)
        if page_num == 1 and len(tables) == 1 and len(tables[0]["data"]) <= 3 and not images:
            logger.debug("跳过第 1 页(小表格已添加到标题页)")
            i += 1
            continue
        
        # 跳过完全空白的页面
        if not page_text.strip() and not tables and not images:
            logger.debug("跳过空白页面: 第 %d 页", page_num)
            i += 1
            continue
        
        # 判断当前页面是否为"小内容页"(可以与下一页合并)
        # 条件: 只有一个小表格(<=3行),无图片
        # 注意:不检查文本,因为pdfplumber会把表格内容也提取为文本
        is_small_content = (
            len(tables) == 1 and 
            len(tables[0]["data"]) <= 3 and 
            not images
        )
        
        logger.debug("第 %d 页判断: 表格数=%d, 表格行数=%d, 图片数=%d, 是否小内容页=%s",
                    page_num, len(tables), 
                    len(tables[0]["data"]) if tables else 0,
                    len(images), is_small_content)
        
        # 不进行页面合并,保持PDF原始页面结构
        pages_to_merge = [page_num]
        
        # 创建一个PPT页面,包含合并后的所有内容
        slide = presentation.slides.add_slide(content_layout)
        
        # 从PDF文本中提取页面标题
        page_title = None
        if page_text:
            lines = [line.strip() for line in page_text.split('\n') if line.strip()]
            if lines:
                # 查找可能的标题(长度<50字符,不包含过多标点)
                for line in lines[:5]:  # 检查前5行
                    # 过滤掉页眉页脚特征的文本
                    line_lower = line.lower()
                    is_header_footer = (
                        'proprietary' in line_lower or
                        'confidential' in line_lower or
                        'content' in line_lower or
                        line.startswith('第') and '页' in line or  # "第X页"
                        line.isdigit() or  # 纯数字(页码)
                        len(line) < 3  # 太短
                    )
                    
                    if not is_header_footer and len(line) < 50 and line.count(',') < 3 and line.count('，') < 3:
                        page_title = line
                        break
        
        # 使用模板的标题占位符(而不是创建新textbox)
        if slide.shapes.title:
            if page_title:
                slide.shapes.title.text = page_title
            elif len(pages_to_merge) > 1:
                slide.shapes.title.text = f"第 {pages_to_merge[0]}-{pages_to_merge[-1]} 页"
            else:
                slide.shapes.title.text = f"第 {page_num} 页"
        
        current_top = Inches(1.0)  # 当前垂直位置
        
        # 处理所有合并的页面内容
        for merge_page_num in pages_to_merge:
            merge_content = pages_with_content[merge_page_num]
            merge_tables = merge_content["tables"]
            merge_images = merge_content["images"]
            merge_text = merge_content.get("text", "")
            
            # 1. 添加表格（如果有）
            for table_data in merge_tables:
                rows_data = table_data["data"]
                rows = len(rows_data)
                cols = len(rows_data[0]) if rows_data else 0
                
                if rows > 0 and cols > 0:
                    left = Inches(0.5)
                    width = Inches(9)
                    # 根据行数动态计算高度
                    row_height = min(0.4, 2.5 / rows)  # 每行最多 0.4 英寸，总高度不超过 2.5 英寸
                    height = Inches(row_height * rows)
                    
                    table = slide.shapes.add_table(rows, cols, left, current_top, width, height).table
                    
                    # 填充表格数据
                    for row_idx, row_data in enumerate(rows_data):
                        for col_idx, cell_value in enumerate(row_data):
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(cell_value) if cell_value else ""
                            cell.text_frame.paragraphs[0].font.size = PptPt(11)
                            
                            # 标题行加粗
                            if row_idx == 0:
                                cell.text_frame.paragraphs[0].font.bold = True
                    
                    current_top += height + Inches(0.3)  # 表格后留间距
                    logger.debug("已添加表格: 第 %d 页, %d 行 x %d 列, 位置=%.2f英寸, 高度=%.2f英寸", 
                                merge_page_num, rows, cols, (current_top - height - Inches(0.3)) / 914400, height / 914400)
            
            # 2. 添加图片（如果有，且空间足够）
            # PPT标准页面高度7.5英寸,预留底部1.5英寸给文本,可用高度约5.5英寸
            max_content_height = Inches(5.5)
            
            if merge_images and current_top < max_content_height:
                # 如果有2张图片,左右并排显示
                if len(merge_images) == 2:
                    available_height = max_content_height - current_top
                    available_width_per_img = Inches(4.5)  # 每张图片最大宽度4.5英寸
                    
                    img_positions = []  # 存储每张图片的位置和尺寸
                    max_img_height = 0
                    
                    for img_idx, img_data in enumerate(merge_images):
                        img_path = img_data["path"]
                        if not img_path.exists():
                            continue
                        
                        try:
                            img_width_px = img_data["width"]
                            img_height_px = img_data["height"]
                            
                            # 将像素转换为英寸
                            img_width_inch = img_width_px / 96.0
                            img_height_inch = img_height_px / 96.0
                            
                            # 计算缩放比例
                            width_ratio = available_width_per_img / Inches(img_width_inch)
                            height_ratio = available_height / Inches(img_height_inch)
                            scale_ratio = min(width_ratio, height_ratio, 1.0)
                            
                            final_width = Inches(img_width_inch * scale_ratio)
                            final_height = Inches(img_height_inch * scale_ratio)
                            
                            # 计算水平位置(左右并排)
                            if img_idx == 0:
                                left = Inches(0.5)  # 左侧图片
                            else:
                                left = Inches(5.5)  # 右侧图片
                            
                            img_positions.append({
                                "path": img_path,
                                "left": left,
                                "top": current_top,
                                "width": final_width,
                                "height": final_height
                            })
                            
                            max_img_height = max(max_img_height, final_height)
                            
                        except Exception as e:
                            logger.warning("图片尺寸计算失败: %s", e)
                    
                    # 添加所有图片
                    for img_pos in img_positions:
                        slide.shapes.add_picture(
                            str(img_pos["path"]),
                            left=img_pos["left"],
                            top=img_pos["top"],
                            width=img_pos["width"],
                            height=img_pos["height"]
                        )
                        logger.debug("已添加图片: 第 %d 页 (左右并排)")
                    
                    current_top += max_img_height + Inches(0.3)
                else:
                    # 单张图片或多张图片,垂直排列
                    for img_idx, img_data in enumerate(merge_images, start=1):
                        if current_top >= max_content_height - Inches(0.5):
                            logger.warning("空间不足，跳过剩余 %d 张图片", len(merge_images) - img_idx + 1)
                            break
                        
                        img_path = img_data["path"]
                        if not img_path.exists():
                            continue
                        
                        try:
                            available_height = max_content_height - current_top - Inches(0.2)
                            max_width = Inches(9)
                            
                            img_width_px = img_data["width"]
                            img_height_px = img_data["height"]
                            img_width_inch = img_width_px / 96.0
                            img_height_inch = img_height_px / 96.0
                            
                            width_ratio = max_width / Inches(img_width_inch)
                            height_ratio = available_height / Inches(img_height_inch)
                            scale_ratio = min(width_ratio, height_ratio, 1.0)
                            
                            final_width = Inches(img_width_inch * scale_ratio)
                            final_height = Inches(img_height_inch * scale_ratio)
                            left = (Inches(10) - final_width) / 2
                            
                            slide.shapes.add_picture(
                                str(img_path),
                                left=left,
                                top=current_top,
                                width=final_width,
                                height=final_height
                            )
                            
                            current_top += final_height + Inches(0.2)
                            logger.debug("已添加图片: 第 %d 页 (%dx%d)", merge_page_num, img_width_px, img_height_px)
                        except Exception as img_error:
                            logger.error("插入图片失败: %s", img_error, exc_info=True)
            
            # 3. 添加文本内容（如果有）
            if merge_text.strip() and not merge_tables:
                # 提取第一行作为标题（如果是纯大写或包含关键词）
                lines = [line.strip() for line in merge_text.split('\n') if line.strip()]
                
                # 检查是否有明显的标题行
                first_line = lines[0] if lines else ""
                if first_line and (first_line.isupper() or len(first_line) < 50):
                    # 使用第一行作为标题(如果之前没有设置标题)
                    if slide.shapes.title and not slide.shapes.title.text:
                        slide.shapes.title.text = first_line
                    content_lines = lines[1:]
                else:
                    content_lines = lines
                
                # 添加文本内容
                if content_lines:
                    text_box = slide.shapes.add_textbox(
                        Inches(0.5), current_top, Inches(9), Inches(4.0)
                    )
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    
                    # 添加内容（限制行数，避免溢出）
                    max_lines = 15
                    for idx, line in enumerate(content_lines[:max_lines]):
                        if idx == 0:
                            text_frame.text = line
                            text_frame.paragraphs[0].font.size = PptPt(14)
                        else:
                            p = text_frame.add_paragraph()
                            p.text = line
                            p.font.size = PptPt(12)
                            p.level = 0
                    
                    current_top += Inches(4.2)
                    logger.debug("已添加文本内容: 第 %d 页, %d 行", merge_page_num, min(len(content_lines), max_lines))
        
        pages_added += 1
        # 统计合并页面的内容
        total_tables = sum(len(pages_with_content[p]["tables"]) for p in pages_to_merge)
        total_images = sum(len(pages_with_content[p]["images"]) for p in pages_to_merge)
        total_text = sum(1 if pages_with_content[p].get("text", "").strip() else 0 for p in pages_to_merge)
        logger.debug("已创建综合页面: 第 %s 页 (文本:%d, 表格:%d, 图片:%d)", 
                    "-".join(map(str, pages_to_merge)) if len(pages_to_merge) > 1 else str(pages_to_merge[0]),
                    total_text, total_tables, total_images)
        
        i += 1  # 移动到下一个未处理的页面
    
    return pages_added


# 旧代码已删除，新逻辑在 _generate_ppt_from_pdf_pages() 中实现
