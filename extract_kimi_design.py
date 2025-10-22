# -*- coding: utf-8 -*-
"""
提取Kimi PPT的设计元素
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os


def analyze_shape_style(shape):
    """分析形状样式"""
    info = {}

    # 填充颜色
    if hasattr(shape, "fill"):
        fill = shape.fill
        if fill.type == 1:  # SOLID
            if hasattr(fill, "fore_color") and hasattr(fill.fore_color, "rgb"):
                info["fill_color"] = fill.fore_color.rgb

    # 线条颜色
    if hasattr(shape, "line"):
        line = shape.line
        if hasattr(line, "color") and hasattr(line.color, "rgb"):
            info["line_color"] = line.color.rgb
        if hasattr(line, "width"):
            info["line_width"] = line.width

    # 文字样式
    if hasattr(shape, "text_frame"):
        tf = shape.text_frame
        if len(tf.paragraphs) > 0:
            para = tf.paragraphs[0]
            if len(para.runs) > 0:
                run = para.runs[0]
                font = run.font
                info["font_name"] = font.name
                info["font_size"] = font.size
                if hasattr(font, "color") and hasattr(font.color, "rgb"):
                    info["font_color"] = font.color.rgb
                info["font_bold"] = font.bold

    return info


def extract_images(prs, output_dir="extracted_images"):
    """提取PPT中的图片"""
    os.makedirs(output_dir, exist_ok=True)

    image_count = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob

                # 保存图片
                ext = image.ext
                filename = f"slide{slide_idx+1}_shape{shape_idx}_{shape.width}x{shape.height}.{ext}"
                filepath = os.path.join(output_dir, filename)

                with open(filepath, "wb") as f:
                    f.write(image_bytes)

                print(f"提取图片: {filename}")
                print(f"  位置: ({shape.left}, {shape.top})")
                print(f"  大小: {shape.width} x {shape.height}")
                image_count += 1

    return image_count


def analyze_kimi_template():
    """分析Kimi模板"""
    prs = Presentation("config/templates/kimi_AI_style_template.pptx")

    print("=" * 80)
    print("Kimi PPT 设计元素分析")
    print("=" * 80)

    # 幻灯片尺寸
    print(f"\n【幻灯片尺寸】")
    print(f"宽度: {prs.slide_width} ({prs.slide_width/914400:.2f}英寸)")
    print(f"高度: {prs.slide_height} ({prs.slide_height/914400:.2f}英寸)")
    print(
        f"比例: 16:9"
        if abs(prs.slide_width / prs.slide_height - 16 / 9) < 0.01
        else "其他"
    )

    # 分析封面页（第1页）
    print(f"\n{'='*80}")
    print("【第1页 - 封面设计】")
    print(f"{'='*80}")

    slide1 = prs.slides[0]
    for i, shape in enumerate(slide1.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"\n形状{i}: 图片")
            print(f"  位置: ({shape.left}, {shape.top})")
            print(f"  大小: {shape.width} x {shape.height}")
        elif hasattr(shape, "text") and shape.text:
            print(f"\n形状{i}: 文本框")
            print(f"  内容: {shape.text[:50]}")
            print(f"  位置: ({shape.left}, {shape.top})")
            print(f"  大小: {shape.width} x {shape.height}")

            # 分析样式
            style = analyze_shape_style(shape)
            if "fill_color" in style:
                rgb = style["fill_color"]
                print(f"  填充色: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")
            if "font_size" in style and style["font_size"]:
                print(f"  字体大小: {style['font_size'].pt}pt")
            if "font_color" in style:
                rgb = style["font_color"]
                print(f"  字体颜色: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")

    # 分析目录页（第2页）
    print(f"\n{'='*80}")
    print("【第2页 - 目录设计】")
    print(f"{'='*80}")

    slide2 = prs.slides[1]
    print(f"总形状数: {len(slide2.shapes)}")

    # 找出编号框
    print("\n编号框样式:")
    for i, shape in enumerate(slide2.shapes):
        if hasattr(shape, "text") and shape.text in ["01", "02", "03", "04", "05"]:
            print(f"\n编号 {shape.text}:")
            print(f"  位置: ({shape.left}, {shape.top})")
            print(f"  大小: {shape.width} x {shape.height}")
            style = analyze_shape_style(shape)
            if "fill_color" in style:
                rgb = style["fill_color"]
                print(f"  填充色: RGB({rgb[0]}, {rgb[1]}, {rgb[2]})")

    # 分析章节页（第3页）
    print(f"\n{'='*80}")
    print("【第3页 - 章节页设计】")
    print(f"{'='*80}")

    slide3 = prs.slides[2]
    for i, shape in enumerate(slide3.shapes):
        if hasattr(shape, "text") and shape.text:
            print(f"\n形状{i}: {shape.text[:30]}")
            print(f"  位置: ({shape.left}, {shape.top})")
            print(f"  大小: {shape.width} x {shape.height}")
            style = analyze_shape_style(shape)
            if "font_size" in style and style["font_size"]:
                print(f"  字体大小: {style['font_size'].pt}pt")

    # 提取图片
    print(f"\n{'='*80}")
    print("【提取图片资源】")
    print(f"{'='*80}")
    image_count = extract_images(prs)
    print(f"\n共提取 {image_count} 张图片")


if __name__ == "__main__":
    analyze_kimi_template()
