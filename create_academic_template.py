# -*- coding: utf-8 -*-
"""
创建学术风格PPT模板

基于商务模板，调整为学术风格：
- 配色：深绿色系（学术、专业）
- 更简洁的设计
- 适合论文、研究报告
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_academic_template():
    """创建学术风格模板"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    
    print("Creating Academic Style Template...")
    
    # 学术风格配色
    COLOR_PRIMARY_DARK = RGBColor(0, 102, 68)  # 深绿
    COLOR_PRIMARY_LIGHT = RGBColor(76, 175, 80)  # 浅绿
    COLOR_ACCENT = RGBColor(255, 193, 7)  # 金色强调
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_DARK = RGBColor(50, 50, 50)
    COLOR_TEXT_LIGHT = RGBColor(150, 150, 150)
    
    # ========== 创建封面示例页 ==========
    print("\n1. Creating cover slide...")
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景渐变（绿色系）
    background = slide1.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = COLOR_PRIMARY_DARK
    fill.gradient_stops[1].color.rgb = COLOR_PRIMARY_LIGHT
    
    # 顶部装饰条（金色）
    top_bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(13.33), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ACCENT
    top_bar.line.fill.background()
    
    # 主标题文本框
    title_box = slide1.shapes.add_textbox(
        Inches(1.5), Inches(2.8),
        Inches(10), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "学术研究报告标题"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # 副标题
    subtitle_box = slide1.shapes.add_textbox(
        Inches(2), Inches(4.2),
        Inches(9.33), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "副标题或研究方向"
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # 作者信息（左下）
    author_box = slide1.shapes.add_textbox(
        Inches(1.5), Inches(6.2),
        Inches(4), Inches(0.8)
    )
    author_frame = author_box.text_frame
    author_frame.text = "作者: XXX\n单位: XXX大学"
    author_frame.paragraphs[0].font.size = Pt(16)
    author_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # 日期信息（右下）
    date_box = slide1.shapes.add_textbox(
        Inches(8), Inches(6.5),
        Inches(3.5), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_frame.text = "日期: 2025/01/01"
    date_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    date_frame.paragraphs[0].font.size = Pt(16)
    date_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # ========== 创建目录页示例 ==========
    print("2. Creating catalog slide...")
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide2.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = COLOR_PRIMARY_DARK
    fill.gradient_stops[1].color.rgb = COLOR_PRIMARY_LIGHT
    
    # 顶部装饰条
    top_bar = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(13.33), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ACCENT
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide2.shapes.add_textbox(
        Inches(1), Inches(0.8),
        Inches(3), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "目录"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # CONTENTS副标题
    subtitle_box = slide2.shapes.add_textbox(
        Inches(4), Inches(1),
        Inches(3), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "CONTENTS"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_LIGHT
    
    # 目录项示例
    catalog_items = [
        "01  研究背景与意义",
        "02  文献综述",
        "03  研究方法",
        "04  实验结果",
        "05  结论与展望"
    ]
    
    start_y = 2.2
    item_height = 0.8
    
    for i, item in enumerate(catalog_items):
        y_pos = start_y + i * item_height
        
        # 目录项框
        item_shape = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_pos),
            Inches(10), Inches(0.6)
        )
        item_shape.fill.solid()
        item_shape.fill.fore_color.rgb = COLOR_WHITE
        item_shape.fill.fore_color.brightness = -0.1
        item_shape.line.fill.background()
        
        item_text = item_shape.text_frame
        item_text.text = item
        item_text.paragraphs[0].font.size = Pt(20)
        item_text.paragraphs[0].font.color.rgb = COLOR_TEXT_DARK
        item_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        item_text.margin_left = Inches(0.3)
    
    # ========== 创建内容页示例 ==========
    print("3. Creating content slide...")
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide3.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = COLOR_PRIMARY_DARK
    fill.gradient_stops[1].color.rgb = COLOR_PRIMARY_LIGHT
    
    # 顶部装饰条
    top_bar = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(13.33), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ACCENT
    top_bar.line.fill.background()
    
    # 标题栏
    title_bar = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, Inches(0.3),
        Inches(13.33), Inches(0.9)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_WHITE
    title_bar.line.fill.background()
    
    # 标题文字
    title_text = title_bar.text_frame
    title_text.text = "内容页标题"
    title_text.paragraphs[0].font.size = Pt(32)
    title_text.paragraphs[0].font.bold = True
    title_text.paragraphs[0].font.color.rgb = COLOR_PRIMARY_DARK
    title_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_text.margin_left = Inches(0.5)
    
    # 内容区域
    content_box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.8),
        Inches(11.73), Inches(5)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COLOR_WHITE
    content_box.line.fill.background()
    
    # 内容文字
    content_text = content_box.text_frame
    content_text.text = "• 研究要点1\n• 研究要点2\n• 研究要点3"
    content_text.paragraphs[0].font.size = Pt(22)
    content_text.paragraphs[0].font.color.rgb = COLOR_TEXT_DARK
    content_text.margin_left = Inches(0.5)
    content_text.margin_top = Inches(0.3)
    
    # ========== 创建章节页示例 ==========
    print("4. Creating section slide...")
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide4.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = COLOR_PRIMARY_DARK
    fill.gradient_stops[1].color.rgb = COLOR_PRIMARY_LIGHT
    
    # 顶部装饰条
    top_bar = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(13.33), Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_ACCENT
    top_bar.line.fill.background()
    
    # 大号数字
    number_box = slide4.shapes.add_textbox(
        Inches(3), Inches(2.5),
        Inches(7.33), Inches(1.5)
    )
    number_frame = number_box.text_frame
    number_frame.text = "01"
    number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    number_frame.paragraphs[0].font.size = Pt(100)
    number_frame.paragraphs[0].font.bold = True
    number_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
    
    # 章节标题
    section_title = slide4.shapes.add_textbox(
        Inches(3), Inches(4.2),
        Inches(7.33), Inches(0.8)
    )
    section_frame = section_title.text_frame
    section_frame.text = "章节标题"
    section_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    section_frame.paragraphs[0].font.size = Pt(40)
    section_frame.paragraphs[0].font.bold = True
    section_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    
    # 保存
    output_path = 'config/templates/academic_template.pptx'
    prs.save(output_path)
    print(f"\n[OK] Academic template saved: {output_path}")
    
    return output_path

if __name__ == '__main__':
    template_path = create_academic_template()
    print(f"\nTemplate created successfully!")
    print(f"Path: {template_path}")
    print(f"\nColor scheme: Green (Academic)")
    print(f"- Primary Dark: RGB(0, 102, 68)")
    print(f"- Primary Light: RGB(76, 175, 80)")
    print(f"- Accent: RGB(255, 193, 7) - Gold")
