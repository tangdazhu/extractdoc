# -*- coding: utf-8 -*-
"""
基于模板的PPT生成器

正确使用PPT模板的占位符机制，保留所有样式
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Pt as PptPt

from .placeholder_helper import PlaceholderHelper
from .text_formatter import TextFormatter

logger = logging.getLogger("converter")


class TemplateBasedPPTGenerator:
    """
    基于模板的PPT生成器
    
    核心理念：
    1. 使用模板的布局（slide_layouts）
    2. 通过占位符（placeholders）填充内容
    3. 保留模板的所有样式（背景、字体、颜色等）
    4. 不手动创建形状，不修改样式
    """
    
    # 布局索引映射
    LAYOUT_TITLE_SLIDE = 0          # 封面
    LAYOUT_TITLE_AND_CONTENT = 1    # 标题+内容
    LAYOUT_SECTION_HEADER = 2       # 章节标题
    LAYOUT_TWO_CONTENT = 3          # 两列内容
    LAYOUT_TITLE_ONLY = 5           # 仅标题
    LAYOUT_PICTURE_WITH_CAPTION = 8 # 图片+说明
    
    def __init__(self, template_path: Path):
        """
        初始化生成器
        
        Args:
            template_path: PPT模板文件路径
        """
        self.template_path = template_path
        self.prs = None
        
        if not template_path or not template_path.exists():
            raise FileNotFoundError(f"模板文件不存在: {template_path}")
        
        logger.info(f"加载PPT模板: {template_path}")
        self.prs = Presentation(str(template_path))
        
        # 删除模板中的示例页（保留第一页作为封面）
        self._remove_example_slides()
        
        # 标记是否已使用第一页作为封面
        self.cover_created = False
    
    def _remove_example_slides(self):
        """删除模板中第2页及之后的示例页，保留第1页作为封面"""
        if len(self.prs.slides) > 1:
            xml_slides = self.prs.slides._sldIdLst
            # 只删除第2页及之后的页面
            for idx in reversed(range(1, len(xml_slides))):
                rId = xml_slides[idx].rId
                self.prs.part.drop_rel(rId)
                del xml_slides[idx]
            logger.debug("已删除模板示例页（保留第1页）")
    
    def create_cover_slide(self, title: str, subtitle: str = "") -> bool:
        """
        创建封面页（使用模板第一页）
        
        Args:
            title: 标题
            subtitle: 副标题
            
        Returns:
            是否成功创建
        """
        # 使用模板第一页（已经有样式）
        if len(self.prs.slides) > 0 and not self.cover_created:
            slide = self.prs.slides[0]
            self.cover_created = True
            logger.debug("使用模板第一页作为封面")
        else:
            # 如果没有示例页或已经使用过，创建新的
            layout = self.prs.slide_layouts[self.LAYOUT_TITLE_SLIDE]
            slide = self.prs.slides.add_slide(layout)
            logger.debug("创建新封面页")
        
        # 填充标题
        success = PlaceholderHelper.fill_title(slide, title)
        
        # 填充副标题
        if subtitle:
            PlaceholderHelper.fill_subtitle(slide, subtitle)
        
        logger.info(f"创建封面页: {title}")
        return success
    
    def create_content_slide(self, title: str, content_lines: List[str]) -> bool:
        """
        创建内容页（标题+列表）
        
        Args:
            title: 标题
            content_lines: 内容行列表
            
        Returns:
            是否成功创建
        """
        layout = self.prs.slide_layouts[self.LAYOUT_TITLE_AND_CONTENT]
        slide = self.prs.slides.add_slide(layout)
        
        # 填充标题
        PlaceholderHelper.fill_title(slide, title)
        
        # 解析Markdown格式
        parsed_lines = TextFormatter.parse_markdown_text("\n".join(content_lines))
        
        # 填充内容到占位符
        content_ph = PlaceholderHelper.get_content_placeholder(slide)
        if content_ph and hasattr(content_ph, 'text_frame'):
            text_frame = content_ph.text_frame
            text_frame.clear()
            
            for idx, (line_text, indent_level, is_bold) in enumerate(parsed_lines):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line_text
                p.level = min(indent_level, 8)
                
                if is_bold:
                    p.font.bold = True
            
            logger.debug(f"创建内容页: {title} ({len(parsed_lines)}行)")
            return True
        else:
            logger.warning(f"未找到内容占位符: {title}")
            return False
    
    def create_section_slide(self, section_title: str, description: str = "") -> bool:
        """
        创建章节页
        
        Args:
            section_title: 章节标题
            description: 描述文本
            
        Returns:
            是否成功创建
        """
        layout = self.prs.slide_layouts[self.LAYOUT_SECTION_HEADER]
        slide = self.prs.slides.add_slide(layout)
        
        # 填充标题
        PlaceholderHelper.fill_title(slide, section_title)
        
        # 填充描述（使用内容占位符）
        if description:
            content_ph = PlaceholderHelper.get_content_placeholder(slide)
            if content_ph and hasattr(content_ph, 'text_frame'):
                content_ph.text = description
        
        logger.info(f"创建章节页: {section_title}")
        return True
    
    def create_picture_slide(self, title: str, image_path: Path, caption: str = "") -> bool:
        """
        创建图片页
        
        Args:
            title: 标题
            image_path: 图片文件路径
            caption: 图片说明
            
        Returns:
            是否成功创建
        """
        layout = self.prs.slide_layouts[self.LAYOUT_PICTURE_WITH_CAPTION]
        slide = self.prs.slides.add_slide(layout)
        
        # 填充标题
        PlaceholderHelper.fill_title(slide, title)
        
        # 插入图片到占位符
        if image_path and image_path.exists():
            PlaceholderHelper.insert_picture_to_placeholder(slide, str(image_path))
        
        # 填充说明
        if caption:
            content_ph = PlaceholderHelper.get_content_placeholder(slide)
            if content_ph and hasattr(content_ph, 'text_frame'):
                content_ph.text = caption
        
        logger.debug(f"创建图片页: {title}")
        return True
    
    def create_two_column_slide(self, title: str, left_content: List[str], right_content: List[str]) -> bool:
        """
        创建两列内容页
        
        Args:
            title: 标题
            left_content: 左列内容
            right_content: 右列内容
            
        Returns:
            是否成功创建
        """
        layout = self.prs.slide_layouts[self.LAYOUT_TWO_CONTENT]
        slide = self.prs.slides.add_slide(layout)
        
        # 填充标题
        PlaceholderHelper.fill_title(slide, title)
        
        # 获取两个内容占位符
        content_phs = PlaceholderHelper.get_all_content_placeholders(slide)
        
        if len(content_phs) >= 2:
            # 填充左列
            self._fill_placeholder_with_lines(content_phs[0], left_content)
            # 填充右列
            self._fill_placeholder_with_lines(content_phs[1], right_content)
            
            logger.debug(f"创建两列页: {title}")
            return True
        else:
            logger.warning(f"两列布局占位符不足: {len(content_phs)}")
            return False
    
    def _fill_placeholder_with_lines(self, placeholder, lines: List[str]):
        """填充占位符（支持Markdown）"""
        if not hasattr(placeholder, 'text_frame'):
            return
        
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        parsed_lines = TextFormatter.parse_markdown_text("\n".join(lines))
        
        for idx, (line_text, indent_level, is_bold) in enumerate(parsed_lines):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line_text
            p.level = min(indent_level, 8)
            
            if is_bold:
                p.font.bold = True
    
    def save(self, output_path: Path):
        """
        保存PPT文件
        
        Args:
            output_path: 输出文件路径
        """
        self.prs.save(str(output_path))
        logger.info(f"PPT已保存: {output_path}")
    
    def get_slide_count(self) -> int:
        """获取幻灯片数量"""
        return len(self.prs.slides)
