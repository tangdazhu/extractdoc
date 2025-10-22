# -*- coding: utf-8 -*-
"""
Auto-Fit渲染器
使用PowerPoint原生的auto-fit功能,让内容自动适应区域
"""

import logging
from pathlib import Path

from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import MSO_AUTO_SIZE

from .text_formatter import TextFormatter

logger = logging.getLogger(__name__)


class AutoFitRenderer:
    """
    内容渲染器
    
    核心理念:
    1. 优先使用placeholder(auto-size只对placeholder有效)
    2. 设置合理的初始字体大小
    3. 让PowerPoint的auto-size自动调整
    """
    
    @staticmethod
    def render_text(slide, text: str, zone: dict, enable_autofit: bool = True):
        """
        渲染文本(优先使用placeholder,因为auto-fit只对placeholder有效)
        
        Args:
            slide: 幻灯片对象
            text: 文本内容
            zone: 区域定义 (包含left, top, width, height的Inches对象)
            enable_autofit: 是否启用auto-fit
        """
        if not text or not text.strip():
            logger.debug("文本为空,跳过渲染")
            return
        
        # 尝试找到内容placeholder
        content_placeholder = None
        for shape in slide.shapes:
            if (hasattr(shape, 'placeholder_format') and 
                shape.has_text_frame and 
                shape != slide.shapes.title):
                content_placeholder = shape
                break
        
        # 如果找到placeholder,使用它;否则创建textbox
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            logger.debug("使用placeholder渲染文本")
        else:
            # 创建文本框
            textbox = slide.shapes.add_textbox(
                zone['left'],
                zone['top'],
                zone['width'],
                zone['height']
            )
            text_frame = textbox.text_frame
            logger.debug("使用textbox渲染文本")
        
        text_frame.word_wrap = True
        
        # 清除placeholder的内边距(释放更多空间)
        if content_placeholder:
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
        
        # 不使用auto-size(太复杂且不可靠),使用固定小字体
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        
        # 解析Markdown格式文本
        parsed_lines = TextFormatter.parse_markdown_text(text)
        
        # 添加文本内容
        for idx, (line_text, indent_level, is_bold) in enumerate(parsed_lines):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line_text
            
            # 设置缩进级别（PPT支持0-8级）
            p.level = min(indent_level, 8)
            
            # 使用固定小字体(保证不溢出)
            p.font.size = PptPt(9)
            
            # 应用加粗
            if is_bold:
                p.font.bold = True
        
        logger.debug("已渲染文本: %d行, 字体=9pt, 使用placeholder=%s", len(parsed_lines), content_placeholder is not None)
    
    @staticmethod
    def render_table(slide, table_data: list, zone: dict, enable_autofit: bool = True):
        """
        渲染表格(使用auto-fit)
        
        Args:
            slide: 幻灯片对象
            table_data: 表格数据 (二维列表)
            zone: 区域定义
            enable_autofit: 是否启用auto-fit
        """
        if not table_data:
            logger.debug("表格数据为空,跳过渲染")
            return
        
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows == 0 or cols == 0:
            logger.debug("表格行列为0,跳过渲染")
            return
        
        # 创建表格
        table = slide.shapes.add_table(
            rows, cols,
            zone['left'],
            zone['top'],
            zone['width'],
            zone['height']
        ).table
        
        # 填充数据
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value) if cell_value else ""
                
                # 清除单元格内边距
                cell.text_frame.margin_left = PptPt(2)
                cell.text_frame.margin_right = PptPt(2)
                cell.text_frame.margin_top = PptPt(2)
                cell.text_frame.margin_bottom = PptPt(2)
                
                # 设置所有段落的字体大小(不只是第一个段落)
                for para in cell.text_frame.paragraphs:
                    para.font.size = PptPt(7)  # 使用更小的字体
                    para.font.name = 'Arial'  # 统一字体
                    
                    # 表头不加粗(保持字体大小一致)
                    # if row_idx == 0:
                    #     para.font.bold = True
                
                # 不使用auto-size,使用固定小字体
                cell.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                cell.text_frame.word_wrap = True
        
        # 调整行高以适应内容(减小行高,避免浪费空间)
        for row_idx in range(rows):
            row = table.rows[row_idx]
            
            # 计算该行需要的最小高度
            # 7pt字体 + 4pt边距(上下各2pt) + 2pt额外空间 = 13pt ≈ 0.18英寸
            if row_idx == 0:
                # 表头稍微高一点
                row.height = Inches(0.25)
            else:
                # 内容行根据文本量调整
                max_lines = 1
                for col_idx in range(cols):
                    cell = table.cell(row_idx, col_idx)
                    text = cell.text
                    # 估算行数(每30个字符换一行)
                    estimated_lines = max(1, len(text) // 30 + 1)
                    max_lines = max(max_lines, estimated_lines)
                
                # 每行7pt + 2pt行间距 = 9pt ≈ 0.125英寸
                row.height = Inches(0.15 + max_lines * 0.125)
        
        logger.debug("已渲染表格: %d行x%d列, 字体=7pt, 边距=2pt, 行高已自适应", rows, cols)
    
    @staticmethod
    def render_image(slide, image_path: Path, zone: dict, maintain_aspect: bool = True):
        """
        渲染图片
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            zone: 区域定义
            maintain_aspect: 是否保持宽高比
        """
        if not image_path or not image_path.exists():
            logger.warning("图片不存在: %s", image_path)
            return
        
        try:
            # 添加图片
            if maintain_aspect:
                # 保持宽高比,图片可能不会填满整个区域
                pic = slide.shapes.add_picture(
                    str(image_path),
                    zone['left'],
                    zone['top'],
                    width=zone['width']
                )
                
                # 如果高度超过区域,调整大小
                if pic.height > zone['height']:
                    pic.height = zone['height']
            else:
                # 填满整个区域,可能拉伸
                pic = slide.shapes.add_picture(
                    str(image_path),
                    zone['left'],
                    zone['top'],
                    width=zone['width'],
                    height=zone['height']
                )
            
            logger.debug("已渲染图片: %s", image_path.name)
        except Exception as e:
            logger.error("渲染图片失败: %s, 错误: %s", image_path, e)
    
    @staticmethod
    def render_images_side_by_side(slide, image_paths: list, zone: dict):
        """
        渲染多张图片(左右并排)
        
        Args:
            slide: 幻灯片对象
            image_paths: 图片路径列表
            zone: 区域定义
        """
        if not image_paths:
            logger.debug("图片列表为空,跳过渲染")
            return
        
        num_images = len(image_paths)
        if num_images == 0:
            return
        
        # 计算每张图片的宽度
        spacing = Inches(0.2)
        total_spacing = spacing * (num_images - 1)
        available_width = zone['width'] - total_spacing
        img_width = available_width / num_images
        
        # 逐个添加图片
        for idx, img_path in enumerate(image_paths):
            if not img_path or not img_path.exists():
                logger.warning("图片不存在: %s", img_path)
                continue
            
            try:
                img_left = zone['left'] + idx * (img_width + spacing)
                
                pic = slide.shapes.add_picture(
                    str(img_path),
                    img_left,
                    zone['top'],
                    width=img_width
                )
                
                # 如果高度超过区域,调整大小
                if pic.height > zone['height']:
                    pic.height = zone['height']
                
                logger.debug("已渲染图片%d: %s", idx + 1, img_path.name)
            except Exception as e:
                logger.error("渲染图片%d失败: %s, 错误: %s", idx + 1, img_path, e)
        
        logger.debug("已渲染%d张并排图片", num_images)
